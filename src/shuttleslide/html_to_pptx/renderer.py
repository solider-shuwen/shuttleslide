"""
PPTX Renderer — converts PPT-DSL JSON to a PowerPoint file.

Consumes the PresentationDSL data model (from schema.py) and produces
a .pptx file using python-pptx.  Each element type has its own render method.
"""

from __future__ import annotations

import io
import logging
from pathlib import Path
from typing import Optional

from lxml import etree
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn

from shuttleslide.html_to_pptx.schema import (
    PresentationDSL, SlideDSL, ThemeDef,
    TitleBarElement, TextBoxElement, ImageElement, ShapeElement,
    GradientOverlayElement, BlurGlowElement, IconTextElement,
    CardElement, NumberedStepElement, DividerLineElement,
    BadgeElement, BulletItem, BulletListElement, TableElement, SVGElement,
    BackgroundDef, GradientDef, BorderDef, ShadowDef,
    TextBlock, TextRun, PositionPercent,
)
from shuttleslide.html_to_pptx.layouts import (
    SLIDE_WIDTH_EMU, SLIDE_HEIGHT_EMU,
    position_percent_to_region, Region, px_to_emu, pct_to_emu,
)
from shuttleslide.html_to_pptx.style_mapper import (
    hex_to_rgbcolor, parse_hex_color, color_opacity, px_to_pt, pt_to_emu,
    map_alignment, gradient_angle_deg, opacity_to_transparency, icon_to_unicode,
)
from shuttleslide.html_to_pptx.image_utils import ImageCache
from shuttleslide.html_to_pptx.fonts import parse_css_font_family

logger = logging.getLogger(__name__)

# DrawingML namespace
_A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'

# Schema vertical_align string → python-pptx MSO_ANCHOR.
_VERTICAL_ANCHOR_MAP = {
    "top": MSO_ANCHOR.TOP,
    "middle": MSO_ANCHOR.MIDDLE,
    "bottom": MSO_ANCHOR.BOTTOM,
}


def _set_typeface_child(rPr, tag: str, value: str, after: str) -> None:
    """Insert or update ``<a:{tag} typeface="value"/>`` in a run's ``rPr``.

    Positions the element right after ``<a:{after}>`` to keep the ECMA-376
    schema order (latin → ea → cs). Updates in place if the element exists.
    """
    if not value:
        return
    q = qn(f"a:{tag}")
    existing = rPr.find(q)
    if existing is not None:
        existing.set("typeface", value)
        return
    new = etree.SubElement(rPr, q)
    new.set("typeface", value)
    anchor = rPr.find(qn(f"a:{after}"))
    if anchor is not None:
        rPr.remove(new)
        anchor.addnext(new)


def _apply_run_fonts(run, font_family_raw: Optional[str]) -> None:
    """Set ``<a:latin>/<a:ea>/<a:cs>`` on a run from a CSS font-family string.

    Replaces the buggy pattern ``run.font.name = elem.font_name`` that wrote
    raw CSS strings (e.g. ``'Nunito, sans-serif'``) verbatim into the
    ``typeface`` attribute, producing unparseable DrawingML.
    """
    fonts = parse_css_font_family(font_family_raw)
    if not (fonts["latin"] or fonts["ea"] or fonts["cs"]):
        return
    # latin via python-pptx API (it manages <a:latin> element lifecycle).
    if fonts["latin"]:
        run.font.name = fonts["latin"]
    # ea/cs: python-pptx doesn't expose these. Insert via lxml, keeping
    # schema order latin → ea → cs.
    rPr = run.font._rPr
    _set_typeface_child(rPr, "ea", fonts["ea"], after="latin")
    _set_typeface_child(rPr, "cs", fonts["cs"], after="ea")


def _blank_layout(prs) -> object:
    """Find the Blank slide layout by name; fall back to index 6.

    `prs.slide_layouts[6]` is Blank in python-pptx's default template, but
    custom templates may number layouts differently. Looking up by name
    is robust to that.
    """
    for layout in prs.slide_layouts:
        if layout.name == "Blank":
            return layout
    return prs.slide_layouts[6]


class PPTXRenderer:
    """Renders a PresentationDSL into a .pptx file."""

    def __init__(
        self,
        image_cache: Optional[ImageCache] = None,
        base_dir: Optional[Path] = None,
    ):
        # Explicit ImageCache wins; otherwise construct one bound to
        # base_dir so relative <img src="images/foo.png"> resolves
        # against the HTML file's parent directory.
        if image_cache is not None:
            self._image_cache = image_cache
        else:
            self._image_cache = ImageCache(base_dir=base_dir)

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def render(self, dsl: PresentationDSL, output_path: str) -> str:
        """Render a PresentationDSL to a PPTX file.

        Args:
            dsl: The presentation model.
            output_path: Where to save the .pptx file.

        Returns:
            The output path as a string.
        """
        prs = Presentation()
        prs.slide_width = Emu(dsl.slide_width_emu)
        prs.slide_height = Emu(dsl.slide_height_emu)

        blank_layout = _blank_layout(prs)

        for slide_dsl in dsl.slides:
            slide = prs.slides.add_slide(blank_layout)
            self._render_slide(slide, slide_dsl, dsl.theme, dsl.slide_width_emu, dsl.slide_height_emu)

        prs.save(output_path)

        # Phase 3b: embed fonts referenced by the DSL. Failures here must
        # not invalidate the PPTX — the file is already saved correctly,
        # only font embedding (a UX nicety) failed.
        try:
            self._embed_fonts(output_path, dsl)
        except Exception as exc:
            logger.warning("Font embedding failed (continuing without embedded fonts): %s", exc)

        return output_path

    # ------------------------------------------------------------------
    # Font embedding
    # ------------------------------------------------------------------

    # Fonts that ship with Windows/Office — no point embedding, PowerPoint
    # will resolve them locally on any stock install.
    _SYSTEM_INSTALLED_FONTS = {
        "Segoe UI", "Segoe UI Historic", "Microsoft YaHei", "Microsoft YaHei UI",
        "Microsoft JhengHei", "SimSun", "SimHei", "KaiTi", "FangSong",
        "Arial", "Times New Roman", "Georgia", "Calibri", "Cambria",
        "Consolas", "Courier New", "Tahoma", "Verdana", "Trebuchet MS",
        "DengXian", "Yu Gothic", "Meiryo", "Malgun Gothic",
    }

    def _embed_fonts(self, output_path: str, dsl: PresentationDSL) -> None:
        """Collect fonts used in the DSL, download them, embed into PPTX.

        Called after ``prs.save()`` so the file already exists on disk; we
        reopen the ZIP and inject font binaries + OpenXML references.
        """
        from dataclasses import asdict
        from shuttleslide.html_to_pptx.fonts import fetch_text_font_bytes
        from shuttleslide.html_to_pptx.font_embedder import embed_fonts

        font_names: set[str] = set()
        for slide in dsl.slides:
            self._collect_font_names(asdict(slide), font_names)

        if not font_names:
            return

        fonts_to_embed: dict[str, bytes] = {}
        for raw_name in font_names:
            # raw_name might still be a CSS stack — parse it to single names.
            parsed = parse_css_font_family(raw_name)
            for slot in ("latin", "ea", "cs"):
                candidate = parsed[slot]
                if not candidate or candidate in self._SYSTEM_INSTALLED_FONTS:
                    continue
                if candidate in fonts_to_embed:
                    continue
                font_bytes = fetch_text_font_bytes(candidate)
                if font_bytes:
                    fonts_to_embed[candidate] = font_bytes

        if fonts_to_embed:
            embed_fonts(output_path, fonts_to_embed)
            logger.info("Embedded %d font(s): %s",
                        len(fonts_to_embed), ", ".join(fonts_to_embed.keys()))

    # Whitelist of DSL field names whose string value is a real font name
    # (something fetch_text_font_bytes can download). Avoids false positives
    # like font_color='#ffffff', text_color='#5a5a5a', icon_font='material-icons',
    # font_size_pt=18, etc. — all of which contain 'font' in the key but are
    # not font names.
    _FONT_NAME_FIELDS = frozenset({
        "font_name",      # TextRun.font_name, TitleBarElement.font_name
        "text_font_name", # IconTextElement.text_font_name
        "font_title",     # ThemeDef.font_title
        "font_body",      # ThemeDef.font_body
    })

    @classmethod
    def _collect_font_names(cls, obj, out: set[str]) -> None:
        """Walk a dataclass-as-dict tree, collect whitelisted font-name fields."""
        if isinstance(obj, dict):
            for k, v in obj.items():
                if k in cls._FONT_NAME_FIELDS and isinstance(v, str) and v:
                    out.add(v)
                else:
                    cls._collect_font_names(v, out)
        elif isinstance(obj, (list, tuple)):
            for item in obj:
                cls._collect_font_names(item, out)
        # Primitives and other types: ignored.

    def render_to_bytes(self, dsl: PresentationDSL) -> bytes:
        """Render to an in-memory bytes buffer."""
        prs = Presentation()
        prs.slide_width = Emu(dsl.slide_width_emu)
        prs.slide_height = Emu(dsl.slide_height_emu)
        blank_layout = _blank_layout(prs)

        for slide_dsl in dsl.slides:
            slide = prs.slides.add_slide(blank_layout)
            self._render_slide(slide, slide_dsl, dsl.theme, dsl.slide_width_emu, dsl.slide_height_emu)

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    # ------------------------------------------------------------------
    # Slide-level rendering
    # ------------------------------------------------------------------

    def _render_slide(
        self, slide, slide_dsl: SlideDSL, theme: ThemeDef,
        slide_w: int, slide_h: int,
    ):
        """Render all elements of a single slide."""
        # Background
        self._render_background(slide, slide_dsl.background, theme)

        # Elements (sorted by z_order)
        sorted_elements = sorted(slide_dsl.elements, key=lambda e: getattr(e, 'z_order', 0))
        for elem in sorted_elements:
            self._render_element(slide, elem, theme, slide_w, slide_h)

    def _render_background(self, slide, bg: Optional[BackgroundDef], theme: ThemeDef):
        """Render slide background."""
        if bg is None:
            # Use theme bg color
            fill = slide.background.fill
            fill.solid()
            rgb = hex_to_rgbcolor(theme.bg_color)
            if rgb:
                fill.fore_color.rgb = rgb
            return

        fill = slide.background.fill
        if bg.type == "solid":
            fill.solid()
            rgb = hex_to_rgbcolor(bg.color or theme.bg_color)
            if rgb:
                fill.fore_color.rgb = rgb
        elif bg.type == "gradient" and bg.gradient:
            self._apply_gradient_fill(fill, bg.gradient)
        elif bg.type == "image" and bg.image_url:
            image_bytes = self._image_cache.get(bg.image_url)
            if image_bytes:
                stream = io.BytesIO(image_bytes)
                slide.shapes.add_picture(
                    stream, 0, 0,
                    Emu(SLIDE_WIDTH_EMU), Emu(SLIDE_HEIGHT_EMU),
                )

    # ------------------------------------------------------------------
    # Element dispatch
    # ------------------------------------------------------------------

    def _render_element(self, slide, elem, theme: ThemeDef,
                        slide_w: int, slide_h: int):
        """Dispatch an element to its specific render method."""
        t = elem.type
        if t == "title_bar":
            self._render_title_bar(slide, elem, theme, slide_w, slide_h)
        elif t == "text_box":
            self._render_text_box(slide, elem, theme, slide_w, slide_h)
        elif t == "image":
            self._render_image(slide, elem, slide_w, slide_h)
        elif t == "shape":
            self._render_shape(slide, elem, slide_w, slide_h)
        elif t == "gradient_overlay":
            self._render_gradient_overlay(slide, elem, slide_w, slide_h)
        elif t == "blur_glow":
            self._render_blur_glow(slide, elem, slide_w, slide_h)
        elif t == "icon_text":
            self._render_icon_text(slide, elem, slide_w, slide_h)
        elif t == "card":
            self._render_card(slide, elem, theme, slide_w, slide_h)
        elif t == "numbered_step":
            self._render_numbered_step(slide, elem, slide_w, slide_h)
        elif t == "divider_line":
            self._render_divider_line(slide, elem, slide_w, slide_h)
        elif t == "badge":
            self._render_badge(slide, elem, slide_w, slide_h)
        elif t == "bullet_list":
            self._render_bullet_list(slide, elem, slide_w, slide_h)
        elif t == "table":
            self._render_table(slide, elem, theme, slide_w, slide_h)
        elif t == "svg":
            self._render_svg(slide, elem, slide_w, slide_h)
        else:
            logger.warning("Unknown element type: %s", t)

    # ------------------------------------------------------------------
    # Element renderers
    # ------------------------------------------------------------------

    def _render_title_bar(self, slide, elem: TitleBarElement, theme: ThemeDef,
                          slide_w: int, slide_h: int):
        """Render a gradient title bar at the top of the slide."""
        bar_h = pct_to_emu(elem.height_pct, slide_h)

        # Background shape
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, Emu(slide_w), Emu(bar_h),
        )
        # Fill
        if elem.background:
            self._apply_background(shape.fill, elem.background)
        else:
            shape.fill.gradient()
            shape.fill.gradient_stops[0].color.rgb = RGBColor(0x13, 0x3E, 0xFF)
            shape.fill.gradient_stops[0].position = 0.0
            shape.fill.gradient_stops[1].color.rgb = RGBColor(0x0d, 0x28, 0xb5)
            shape.fill.gradient_stops[1].position = 1.0

        # Remove border
        shape.line.fill.background()

        # Title text
        tf = shape.text_frame
        tf.word_wrap = True
        # Vertically centre the title within the bar (CSS `items-center`).
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = elem.text
        run.font.size = Pt(float(elem.font_size_pt))
        run.font.bold = elem.font_bold
        rgb = hex_to_rgbcolor(elem.font_color)
        if rgb:
            run.font.color.rgb = rgb
        if elem.font_name:
            _apply_run_fonts(run, elem.font_name)

        # No extra paragraph spacing — the anchor handles vertical placement.
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)

    def _render_text_box(self, slide, elem: TextBoxElement,
                         theme: ThemeDef, slide_w: int, slide_h: int):
        """Render a text box with styled paragraphs."""
        region = self._resolve_region(elem, slide_w, slide_h)

        txBox = slide.shapes.add_textbox(
            Emu(region.left), Emu(region.top),
            Emu(region.width), Emu(region.height),
        )
        tf = txBox.text_frame
        # add_textbox defaults to SHAPE_TO_FIT_TEXT, which makes the shape
        # grow to fit text on one line — defeating word_wrap. Force NONE so
        # the box keeps its specified width and text wraps inside it.
        tf.auto_size = MSO_AUTO_SIZE.NONE
        # Atomic-token overflow (URLs, decorative numbers, product codes):
        # text glyph is wider than the textbox. Browser would show visual
        # overflow (CSS overflow: visible); PPT matches via word_wrap=False
        # so text stays on one line and overflows the textbox right edge,
        # instead of wrapping inside it.
        tf.word_wrap = not elem.no_wrap

        # Minimize internal margins so text fills the box width accurately
        tf.margin_left = Emu(9144)    # 0.01 inch
        tf.margin_right = Emu(9144)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)

        # Vertical alignment — maps schema's vertical_align string to the
        # text-frame's vertical anchor. Per-paragraph horizontal alignment
        # is still applied below via block.alignment.
        tf.vertical_anchor = _VERTICAL_ANCHOR_MAP.get(
            elem.vertical_align, MSO_ANCHOR.TOP
        )

        # Render each text block as a paragraph
        first = True
        for block in elem.content:
            if first:
                para = tf.paragraphs[0]
                first = False
            else:
                para = tf.add_paragraph()

            para.alignment = map_alignment(block.alignment)

            if block.line_spacing_pt:
                para.line_spacing = Pt(float(block.line_spacing_pt))
            if block.spacing_before_pt:
                para.space_before = Pt(float(block.spacing_before_pt))
            if block.spacing_after_pt:
                para.space_after = Pt(float(block.spacing_after_pt))

            # If runs exist, use them; otherwise use plain text
            if block.runs:
                for r in block.runs:
                    run = para.add_run()
                    self._set_run_text(run, r.text)
                    if r.font_size_pt:
                        run.font.size = Pt(float(r.font_size_pt))
                    if r.color:
                        rgb = hex_to_rgbcolor(r.color)
                        if rgb:
                            run.font.color.rgb = rgb
                        if r.opacity < 1.0:
                            self._apply_text_run_alpha(run, r.opacity)
                    run.font.bold = r.bold
                    run.font.italic = r.italic
                    if r.font_name:
                        _apply_run_fonts(run, r.font_name)
            else:
                run = para.add_run()
                self._set_run_text(run, block.text)
                # Apply block-level defaults
                if block.level in ("title", "h3", "h4"):
                    run.font.bold = True

    def _render_image(self, slide, elem: ImageElement,
                      slide_w: int, slide_h: int):
        """Render an image element."""
        if not elem.url:
            return

        region = self._resolve_region(elem, slide_w, slide_h)
        image_bytes = self._image_cache.get(elem.url)
        if not image_bytes:
            logger.warning("Could not download image: %s", elem.url)
            return

        stream = io.BytesIO(image_bytes)
        pic = slide.shapes.add_picture(
            stream,
            Emu(region.left), Emu(region.top),
            Emu(region.width), Emu(region.height),
        )

        # object-fit: cover — crop the image to fill the box while preserving
        # aspect ratio (mirrors CSS `object-fit: cover`). python-pptx crop
        # properties are fractions of the source image dimensions.
        if elem.object_fit == "cover":
            self._apply_object_fit_cover(pic, image_bytes, region)
        elif elem.object_fit == "contain":
            self._apply_object_fit_contain(pic, image_bytes, region)

        # Border
        if elem.border and elem.border.width_pt > 0:
            pic.line.color.rgb = hex_to_rgbcolor(elem.border.color) or RGBColor(0, 0, 0)
            pic.line.width = Pt(float(elem.border.width_pt))

        # Rounded corners — modify prstGeom from rect to roundRect
        if elem.corner_radius_pct > 0:
            sp = pic._element
            prst_geom = sp.find('.//' + qn('a:prstGeom'))
            if prst_geom is not None:
                prst_geom.set('prst', 'roundRect')
                self._set_rounded_rect_adj(pic, elem.corner_radius_pct)

    def _render_shape(self, slide, elem: ShapeElement, slide_w: int, slide_h: int):
        """Render a geometric shape."""
        region = self._resolve_region(elem, slide_w, slide_h)

        shape_map = {
            "rectangle": MSO_SHAPE.RECTANGLE,
            "rounded_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
            "oval": MSO_SHAPE.OVAL,
            "circle": MSO_SHAPE.OVAL,
        }
        mso_type = shape_map.get(elem.shape_type, MSO_SHAPE.RECTANGLE)

        shape = slide.shapes.add_shape(
            mso_type,
            Emu(region.left), Emu(region.top),
            Emu(region.width), Emu(region.height),
        )

        # Apply corner radius for rounded rectangles
        if elem.shape_type == "rounded_rect" and elem.corner_radius_pct > 0:
            self._set_rounded_rect_adj(shape, elem.corner_radius_pct)

        # Fill
        if elem.background:
            self._apply_background(shape.fill, elem.background)
        else:
            shape.fill.background()  # no fill

        # Border
        if elem.border and elem.border.width_pt > 0:
            rgb = hex_to_rgbcolor(elem.border.color)
            if rgb:
                shape.line.color.rgb = rgb
            shape.line.width = Pt(float(elem.border.width_pt))
            self._apply_border_dash_style(shape, elem.border.style)
        else:
            shape.line.fill.background()

        # Rotation
        if elem.rotation_deg:
            shape.rotation = elem.rotation_deg

    def _render_gradient_overlay(self, slide, elem: GradientOverlayElement,
                                 slide_w: int, slide_h: int):
        """Render a semi-transparent gradient overlay."""
        region = self._resolve_region(elem, slide_w, slide_h)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(region.left), Emu(region.top),
            Emu(region.width), Emu(region.height),
        )

        if elem.gradient:
            self._apply_gradient_fill(shape.fill, elem.gradient, elem.opacity)
        else:
            shape.fill.solid()
            rgb = hex_to_rgbcolor("#133EFF")
            if rgb:
                shape.fill.fore_color.rgb = rgb
            self._set_shape_alpha(shape, elem.opacity)

        shape.line.fill.background()

    def _render_blur_glow(self, slide, elem: BlurGlowElement,
                          slide_w: int, slide_h: int):
        """Render a blur/glow element as an oval with soft edges + alpha."""
        region = self._resolve_region(elem, slide_w, slide_h)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Emu(region.left), Emu(region.top),
            Emu(region.width), Emu(region.height),
        )

        # Solid fill with transparency
        shape.fill.solid()
        rgb = hex_to_rgbcolor(elem.color)
        if rgb:
            shape.fill.fore_color.rgb = rgb
        self._set_shape_alpha(shape, elem.opacity)

        shape.line.fill.background()

        # Apply soft edge to simulate CSS blur
        if elem.blur_radius_pt > 0:
            self._set_soft_edge(shape, elem.blur_radius_pt)

    def _render_icon_text(self, slide, elem: IconTextElement,
                          slide_w: int, slide_h: int):
        """Render an icon + optional label side-by-side.

        Two cases:
        - Icon-only element (no ``text``): the whole region belongs to the
          icon. Backwards-compatible with the previous behaviour.
        - Icon + label (e.g. ``<h3><i>icon</i>Heading</h3>``): the icon
          takes a square on the left sized by ``icon_size_pt``; the text
          fills the remaining width, vertically centred.
        """
        region = self._resolve_region(elem, slide_w, slide_h)
        has_label = bool(elem.text and elem.text.strip())

        if has_label:
            # Reserve a square on the left for the icon.
            icon_size_emu = int(float(elem.icon_size_pt) * 12700)
            # Vertical-centre the icon square within the region.
            icon_top = region.top + (region.height - icon_size_emu) // 2
            icon_region = Region(
                left=region.left,
                top=icon_top,
                width=icon_size_emu,
                height=icon_size_emu,
            )
            # 6pt gap between icon and text.
            gap_emu = int(6 * 12700)
            text_left = region.left + icon_size_emu + gap_emu
            text_region = Region(
                left=text_left,
                top=region.top,
                width=max(region.width - icon_size_emu - gap_emu, 0),
                height=region.height,
            )
        else:
            icon_region = region
            text_region = None

        # Try to render as vector custom geometry from font
        shape = self._try_render_vector_icon(slide, elem, icon_region)
        if shape is None:
            # Fallback: simple text box with a Unicode glyph for the icon.
            # Vector rendering fails when the Material Icons font isn't bundled
            # locally; fall back to a Unicode equivalent (✓, ⚠, ?, etc.) so we
            # never render the raw icon *name* ("bug_report") as display text.
            logger.warning("Vector icon rendering failed for '%s', using text fallback",
                           elem.icon_name)
            self._render_icon_unicode_fallback(
                slide, elem.icon_name, elem.icon_color,
                float(elem.icon_size_pt), icon_region,
            )

        # Render the label text in the remaining region (if any).
        if has_label and text_region is not None:
            self._render_icon_text_label(slide, elem, text_region)

    def _render_icon_unicode_fallback(
        self, slide,
        icon_name: str,
        icon_color: Optional[str],
        icon_size_pt: float,
        region: Region,
    ) -> None:
        """Render a Unicode glyph (✓ ⚠ ? …) as the icon when vector rendering fails.

        Used by _render_icon_text (icon_text element) and _render_bullet_list
        (per-item LI icon). Both paths share the same fallback so a missing
        Material Icons font never produces the raw icon name as display text.
        """
        txBox = slide.shapes.add_textbox(
            Emu(region.left), Emu(region.top),
            Emu(region.width), Emu(region.height),
        )
        tf = txBox.text_frame
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.word_wrap = True
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER

        icon_run = para.add_run()
        icon_run.text = icon_to_unicode(icon_name)
        icon_run.font.size = Pt(float(icon_size_pt))
        if icon_color:
            rgb = hex_to_rgbcolor(icon_color)
            if rgb:
                icon_run.font.color.rgb = rgb

    def _render_icon_text_label(self, slide, elem: IconTextElement, region: Region):
        """Render the text label of an icon_text element in the given region."""
        txBox = slide.shapes.add_textbox(
            Emu(region.left), Emu(region.top),
            Emu(region.width), Emu(region.height),
        )
        tf = txBox.text_frame
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.word_wrap = True
        # Tight margins so the text aligns flush with the gap edge.
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.LEFT
        run = para.add_run()
        run.text = elem.text
        if elem.text_font_size_pt:
            run.font.size = Pt(float(elem.text_font_size_pt))
        if elem.text_color:
            rgb = hex_to_rgbcolor(elem.text_color)
            if rgb:
                run.font.color.rgb = rgb
        run.font.bold = elem.text_bold
        if elem.text_font_name:
            _apply_run_fonts(run, elem.text_font_name)

    def _try_render_vector_icon(self, slide, elem: IconTextElement, region):
        """Attempt to render icon as a vector custom geometry shape.

        Returns the shape if successful, None if fallback should be used.
        Thin wrapper over _render_vector_icon_primitives so callers that
        already hold an IconTextElement keep their signature. Other element
        types (BulletItem, future Badge-with-icon, …) call the primitives
        form directly to avoid forcing an IconTextElement.
        """
        return self._render_vector_icon_primitives(
            slide,
            icon_name=elem.icon_name,
            icon_font=elem.icon_font,
            icon_color=elem.icon_color,
            icon_size_pt=elem.icon_size_pt,
            icon_shadow=elem.icon_shadow,
            region=region,
        )

    def _render_vector_icon_primitives(
        self, slide, *,
        icon_name: str,
        icon_font: Optional[str],
        icon_color: Optional[str],
        icon_size_pt: float,
        icon_shadow: Optional[ShadowDef],
        region: Region,
    ):
        """Render a Material-Icons-style glyph as a vector custGeom shape.

        Central implementation shared by every element type that embeds
        a named icon (icon_text, badge, bullet_list, …). Returns the
        shape on success, None on any failure (missing font, unknown
        codepoint, empty outline) so the caller can fall back to a
        Unicode glyph.
        """
        from shuttleslide.html_to_pptx.fonts import get_font_bytes, icon_to_codepoint
        from shuttleslide.html_to_pptx.vector_icon import (
            extract_glyph_outline,
            glyph_outline_to_custgeom_xml,
        )

        icon_class = icon_font
        if not icon_class or not icon_name:
            logger.debug("Vector icon skip: icon_font=%s icon_name=%s", icon_class, icon_name)
            return None

        font_bytes = get_font_bytes(icon_class)
        if not font_bytes:
            logger.debug("Vector icon skip: no font bytes for '%s'", icon_class)
            return None

        codepoint = icon_to_codepoint(icon_class, icon_name)
        if codepoint is None:
            logger.debug("Vector icon skip: no codepoint for '%s' in '%s'", icon_name, icon_class)
            return None

        outline = extract_glyph_outline(font_bytes, codepoint)
        if outline is None or not outline.commands:
            logger.debug("Vector icon skip: no outline for '%s' (U+%04X)", icon_name, codepoint)
            return None

        cust_geom_xml = glyph_outline_to_custgeom_xml(outline)

        logger.info(
            "Vector icon '%s' U+%04X: %d commands, path %dx%d, region %s",
            icon_name, codepoint, len(outline.commands),
            outline.path_w, outline.path_h, region,
        )

        # Adjust shape to preserve custGeom path aspect ratio
        # so the icon doesn't get stretched
        shape_left, shape_top, shape_w, shape_h = region.left, region.top, region.width, region.height
        path_aspect = outline.path_w / outline.path_h
        shape_aspect = shape_w / shape_h if shape_h > 0 else 1.0
        if path_aspect > shape_aspect:
            # Path wider than shape → shrink height, center vertically
            new_h = round(shape_w / path_aspect)
            shape_top += (shape_h - new_h) // 2
            shape_h = new_h
        else:
            # Path taller than shape → shrink width, center horizontally
            new_w = round(shape_h * path_aspect)
            shape_left += (shape_w - new_w) // 2
            shape_w = new_w

        # Create placeholder shape
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(shape_left), Emu(shape_top),
            Emu(shape_w), Emu(shape_h),
        )

        # Replace <a:prstGeom> with <a:custGeom>
        # Use spPr.makeelement() to inherit parent's namespace context
        spPr = shape._element.spPr
        prst_geom = spPr.find(qn("a:prstGeom"))
        if prst_geom is not None:
            spPr.replace(prst_geom, cust_geom_xml)
        else:
            spPr.insert(0, cust_geom_xml)

        # Fill with icon color
        if icon_color:
            shape.fill.solid()
            rgb = hex_to_rgbcolor(icon_color)
            if rgb:
                shape.fill.fore_color.rgb = rgb
            else:
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(0, 0, 0)
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0, 0, 0)

        # No outline
        shape.line.fill.background()

        # ICON+GLOW: filter: drop-shadow on the icon → outerShdw.
        self._apply_shape_shadow(shape, icon_shadow)

        return shape

    # ------------------------------------------------------------------
    # SVG → native DrawingML via vendored svg_to_pptx
    # ------------------------------------------------------------------

    def _render_svg(self, slide, elem: SVGElement,
                    slide_w: int, slide_h: int):
        """Render an inline ``<svg>`` as a native editable PPT group shape.

        Parses ``elem.svg_markup``, computes scale from viewBox → target
        rect, runs ``convert_element`` per child via the vendored
        svg_to_pptx library, then wraps all emitted ``<p:sp>`` shapes in
        a single ``<p:grpSp>`` so the user can select / move / resize
        the whole SVG as one object in PowerPoint (Ctrl+U ungroup to
        get individual shapes).
        """
        if not elem.svg_markup or not elem.position:
            return

        # Lazy imports — vendor lib is only pulled in when a slide
        # actually contains SVG markup, keeping import-time cost zero
        # for slides that don't.
        from xml.etree import ElementTree as ET
        from pptx.oxml import parse_xml
        from pptx.oxml.ns import nsdecls
        from shuttleslide._vendored.svg_to_pptx import (
            convert_element, collect_defs, ConvertContext,
        )
        from shuttleslide._vendored.svg_to_pptx.drawingml_utils import EMU_PER_PX

        try:
            root = ET.fromstring(elem.svg_markup)
        except ET.ParseError as exc:
            logger.warning("SVG slot=%s parse failed: %s", elem.slot_id, exc)
            self._render_svg_fallback(slide, elem, slide_w, slide_h)
            return

        # Target rect on the slide, in EMU.
        target_left = pct_to_emu(elem.position.x_pct, slide_w)
        target_top = pct_to_emu(elem.position.y_pct, slide_h)
        target_w = pct_to_emu(elem.position.w_pct, slide_w)
        target_h = pct_to_emu(elem.position.h_pct, slide_h)

        # Resolve SVG-space dimensions from viewBox (cached attr or root
        # attribute); fall back to width/height attrs.
        vb = elem.viewBox or root.get("viewBox") or root.get("viewbox")
        vb_w = vb_h = 0.0
        if vb:
            parts = vb.split()
            if len(parts) == 4:
                try:
                    vb_w, vb_h = float(parts[2]), float(parts[3])
                except ValueError:
                    vb_w = vb_h = 0.0
        if vb_w <= 0 or vb_h <= 0:
            try:
                vb_w = float(root.get("width", "0") or 0)
                vb_h = float(root.get("height", "0") or 0)
            except ValueError:
                vb_w = vb_h = 0.0
        if vb_w <= 0 or vb_h <= 0:
            logger.warning(
                "SVG slot=%s has no usable viewBox/size — fallback",
                elem.slot_id,
            )
            self._render_svg_fallback(slide, elem, slide_w, slide_h)
            return

        # Vendor library's translate/scale operate in SVG pixel units
        # (1 px = EMU_PER_PX EMU). We compute the affine that maps the
        # SVG's viewBox rectangle onto the target rect in slide-EMU
        # space, then express the translation in SVG pixels.
        #
        # object-fit handling mirrors _apply_object_fit_cover / _contain
        # for raster images:
        #   * fill (default): stretch the viewBox independently along x
        #     and y to fill the target rect. Legacy behavior — preserves
        #     backward compat for SVGs that were authored/cropped to the
        #     exact target aspect.
        #   * cover: uniform scale = max(scale_x, scale_y); center the
        #     over-scaled content so it covers the whole target rect,
        #     with overflow going past the rect (cropped visually by an
        #     overlay rectangle above the group, mirroring the source
        #     HTML's <div style="overflow:hidden">).
        #   * contain: uniform scale = min(scale_x, scale_y); center the
        #     under-scaled content, leaving empty bands inside the rect.
        scale_x = (target_w / EMU_PER_PX) / vb_w
        scale_y = (target_h / EMU_PER_PX) / vb_h
        translate_x = target_left / EMU_PER_PX
        translate_y = target_top / EMU_PER_PX
        # ``group_off_*`` / ``group_ext_*`` describe the SVG's actual
        # rendered bounding box in slide-EMU coordinates — what the
        # wrapper ``<p:grpSp>`` advertises as its ``<a:off>/<a:ext>``
        # for PowerPoint hit-testing and selection. This must reflect
        # the SVG bbox, NOT the HTML container bbox (``target_*``):
        #   - cover: rendered > target on one axis, offset < target_left
        #     (SVG overflows the container on both sides). Using target
        #     here would make the group selection box smaller than the
        #     rendered children, so clicking SVG content outside the
        #     box bypasses the group and selects individual shapes —
        #     the "SVG looks like a pile of loose parts, not a group"
        #     bug (slide 1 with landscape SVG in portrait container).
        #   - contain: rendered < target on one axis, offset > target_left
        #     (SVG under-fills the container). Using target here makes
        #     the selection box include empty space around the SVG —
        #     imprecise but not user-visible; still wrong.
        #   - fill: rendered == target on both axes (independent stretch).
        #     Using target here is correct and equals the rendered bbox.
        group_off_x_emu = target_left
        group_off_y_emu = target_top
        group_ext_w_emu = target_w
        group_ext_h_emu = target_h
        if elem.object_fit == "cover" or elem.object_fit == "contain":
            s = max(scale_x, scale_y) if elem.object_fit == "cover" else min(scale_x, scale_y)
            scale_x = scale_y = s
            # Re-center the uniformly-scaled SVG inside the target rect.
            # In EMU: rendered_w = vb_w * s * EMU_PER_PX. Offset so the
            # visible region is centered (cover: rendered ≥ target on
            # both axes, offsets ≤ 0; contain: rendered ≤ target,
            # offsets ≥ 0). Convert back to SVG px for the translate.
            rendered_w_emu = vb_w * s * EMU_PER_PX
            rendered_h_emu = vb_h * s * EMU_PER_PX
            offset_x_emu = target_left + (target_w - rendered_w_emu) / 2
            offset_y_emu = target_top + (target_h - rendered_h_emu) / 2
            translate_x = offset_x_emu / EMU_PER_PX
            translate_y = offset_y_emu / EMU_PER_PX
            # Overwrite the default (target rect) with the SVG's actual
            # rendered bbox so the group selection box tracks the SVG,
            # not the HTML container.
            group_off_x_emu = int(offset_x_emu)
            group_off_y_emu = int(offset_y_emu)
            group_ext_w_emu = int(rendered_w_emu)
            group_ext_h_emu = int(rendered_h_emu)

        # Allocate shape IDs that don't collide with existing shapes on
        # the slide. Vendor library starts its id_counter at 2, but the
        # slide may already have shapes (text boxes, cards, etc.) that
        # used ids 2..N. Find the max existing id and bump past it.
        sp_tree = slide.shapes._spTree
        max_id = 1
        for node in sp_tree.iter():
            id_attr = node.get("id")
            if id_attr and id_attr.isdigit():
                max_id = max(max_id, int(id_attr))
        group_id = max_id + 1

        defs = collect_defs(root)
        # If a wrapping HTML element had CSS opacity < 1.0 (e.g.
        # <div style="opacity:0.25"><svg>…), the converter captured the
        # cumulative ancestor opacity on elem.opacity. The vendored
        # library already multiplies inherited_styles.opacity into every
        # fill/stroke alpha (drawingml_context.py:120-128), so seeding
        # it at the root makes the entire SVG render at the browser's
        # effective opacity without any further changes.
        inherited = {}
        if elem.opacity < 1.0:
            inherited["opacity"] = str(elem.opacity)
        ctx = ConvertContext(
            defs=defs,
            id_counter=max_id + 2,  # leave group_id for the grpSp wrapper
            scale_x=scale_x,
            scale_y=scale_y,
            translate_x=translate_x,
            translate_y=translate_y,
            inherited_styles=inherited,
            merge_paragraphs=True,
        )

        # Convert each <svg> child to a <p:sp> and collect into a list.
        # We'll wrap them all in a single <p:grpSp> at the end so the
        # user can manipulate the SVG as one object.
        #
        # Defense-in-depth: filter out any full-bleed background rect
        # that slipped past the validator (e.g. hand-edited or
        # pre-existing HTML like tmp/agent_gen_output/*.html produced
        # before this fix). Vendor lib would happily convert such a
        # rect into a slide-covering <p:sp> that masks the slide bg.
        def _is_full_bleed_bg(c) -> bool:
            ctag = c.tag.split("}", 1)[-1] if "}" in c.tag else c.tag
            if ctag != "rect":
                return False
            try:
                cw = float(c.get("width", "0"))
                ch = float(c.get("height", "0"))
                cx = float(c.get("x", "0"))
                cy = float(c.get("y", "0"))
            except (TypeError, ValueError):
                return False
            tol = 0.01
            return (abs(cw - vb_w) <= vb_w * tol
                    and abs(ch - vb_h) <= vb_h * tol
                    and abs(cx) <= vb_w * tol
                    and abs(cy) <= vb_h * tol)

        child_nodes = []
        for child in root:
            # Skip <defs> (already collected), <title>/<desc> (metadata).
            tag = child.tag.split("}", 1)[-1] if "}" in child.tag else child.tag
            if tag in ("defs", "title", "desc"):
                continue
            if _is_full_bleed_bg(child):
                logger.info(
                    "SVG slot=%s: skipping full-bleed background rect",
                    elem.slot_id,
                )
                continue
            try:
                result = convert_element(child, ctx)
            except Exception as exc:  # vendor lib is broad-exception-prone
                logger.warning(
                    "SVG child <%s> slot=%s convert failed: %s",
                    tag, elem.slot_id, exc,
                )
                continue
            if result is None or not result.xml:
                continue
            try:
                # Vendor output uses <p:> / <a:> prefixes but doesn't
                # declare them. Wrap with nsdecls so lxml can resolve,
                # then pull the parsed children out of the wrapper.
                wrapper_xml = (
                    f'<root {nsdecls("p", "a")}>{result.xml}</root>'
                )
                wrapper = parse_xml(wrapper_xml)
                child_nodes.extend(wrapper)
            except Exception as exc:
                logger.warning(
                    "SVG shape XML inject failed slot=%s: %s",
                    elem.slot_id, exc,
                )

        if not child_nodes:
            logger.info(
                "SVG slot=%s produced no shapes — fallback",
                elem.slot_id,
            )
            self._render_svg_fallback(slide, elem, slide_w, slide_h)
            return

        # Build the grpSp container. The transform is identity-like:
        # since vendor already emitted each child at absolute slide-EMU
        # coordinates (translate_x/y + scale_x/y folded the SVG viewBox
        # onto the target rect), set off/ext == chOff/chExt == SVG's
        # actual rendered bbox so the group transform preserves child
        # positions while advertising the SVG's bbox for hit-testing /
        # selection. For cover this bbox extends past the HTML container
        # (and possibly past the slide edge — PowerPoint's slide viewport
        # still clips the visual, the negative off is just for selection).
        grp_name = f"SVGGroup_{elem.slot_id}" if elem.slot_id else f"SVGGroup_{group_id}"
        grp_xml = (
            f'<p:grpSp {nsdecls("p", "a")}>'
            f'<p:nvGrpSpPr>'
            f'<p:cNvPr id="{group_id}" name="{grp_name}"/>'
            f'<p:cNvGrpSpPr/>'
            f'<p:nvPr/>'
            f'</p:nvGrpSpPr>'
            f'<p:grpSpPr>'
            f'<a:xfrm>'
            f'<a:off x="{group_off_x_emu}" y="{group_off_y_emu}"/>'
            f'<a:ext cx="{group_ext_w_emu}" cy="{group_ext_h_emu}"/>'
            f'<a:chOff x="{group_off_x_emu}" y="{group_off_y_emu}"/>'
            f'<a:chExt cx="{group_ext_w_emu}" cy="{group_ext_h_emu}"/>'
            f'</a:xfrm>'
            f'</p:grpSpPr>'
            f'</p:grpSp>'
        )
        try:
            group = parse_xml(grp_xml)
            group.extend(child_nodes)
            sp_tree.append(group)
        except Exception as exc:
            logger.warning(
                "SVG grpSp build failed slot=%s: %s — falling back to flat shapes",
                elem.slot_id, exc,
            )
            # Last resort: append children directly without grouping.
            sp_tree.extend(child_nodes)

    def _render_svg_fallback(self, slide, elem: SVGElement,
                             slide_w: int, slide_h: int):
        """Last-resort placeholder when vendor conversion yields nothing.

        Renders a light-grey rectangle labelled with the slot_id so the
        region is visible in PowerPoint rather than silently empty.
        """
        if not elem.position:
            return
        region = position_percent_to_region(
            elem.position.x_pct, elem.position.y_pct,
            elem.position.w_pct, elem.position.h_pct,
            slide_w, slide_h,
        )
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(region.left), Emu(region.top),
            Emu(region.width), Emu(region.height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
        shape.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        tf = shape.text_frame
        tf.text = f"[SVG: {elem.slot_id or 'untitled'}]"

    def _render_card(self, slide, elem: CardElement,
                     theme: ThemeDef, slide_w: int, slide_h: int):
        """Render a card (rounded rectangle with optional accent border)."""
        region = self._resolve_region(elem, slide_w, slide_h)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Emu(region.left), Emu(region.top),
            Emu(region.width), Emu(region.height),
        )

        # Apply corner radius
        self._set_rounded_rect_adj(shape, elem.corner_radius_pct)

        # Fill
        if elem.background:
            self._apply_background(shape.fill, elem.background)
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # Opacity (must come after fill so <a:srgbClr> exists to attach alpha to)
        if elem.opacity < 1.0:
            self._set_shape_alpha(shape, elem.opacity)

        # Border
        if elem.border and elem.border.width_pt > 0:
            rgb = hex_to_rgbcolor(elem.border.color)
            if rgb:
                shape.line.color.rgb = rgb
            shape.line.width = Pt(float(elem.border.width_pt))
            self._apply_border_dash_style(shape, elem.border.style)
        else:
            shape.line.fill.background()

        # Drop shadow (CSS box-shadow) — applied after fill so the effectLst
        # sibling ordering stays valid.
        self._apply_shape_shadow(shape, elem.shadow)

        # Per-side accent stripes (CSS `border-left/right/top/bottom: Npx solid #color`).
        # PowerPoint shapes only support uniform borders on all four sides, so
        # single-side borders are rendered as thin rectangles on the relevant edge.
        sides = (
            ("left", elem.border_left, region.left, region.top,
             None, region.height),  # width set per-stripe; height = full
            ("right", elem.border_right,
             region.left + region.width, region.top,
             None, region.height),
            ("top", elem.border_top, region.left, region.top,
             region.width, None),
            ("bottom", elem.border_bottom,
             region.left, region.top + region.height,
             region.width, None),
        )
        for side_name, border_def, base_x, base_y, full_w, full_h in sides:
            if border_def is None or border_def.width_pt <= 0:
                continue
            rgb = hex_to_rgbcolor(border_def.color)
            if rgb is None:
                continue
            stripe_w_emu = pt_to_emu(float(border_def.width_pt))
            if side_name == "left":
                x, y, w, h = base_x, base_y, stripe_w_emu, full_h
            elif side_name == "right":
                x, y, w, h = base_x - stripe_w_emu, base_y, stripe_w_emu, full_h
            elif side_name == "top":
                x, y, w, h = base_x, base_y, full_w, stripe_w_emu
            else:  # bottom
                x, y, w, h = base_x, base_y - stripe_w_emu, full_w, stripe_w_emu
            stripe = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Emu(x), Emu(y), Emu(w), Emu(h),
            )
            stripe.fill.solid()
            stripe.fill.fore_color.rgb = rgb
            stripe.line.fill.background()  # no outline on the stripe itself
            stripe.shadow.inherit = False

        # Render child text elements inside the card as a text frame
        if elem.children:
            tf = shape.text_frame
            tf.word_wrap = True
            first = True
            for child in elem.children:
                if hasattr(child, 'content'):  # TextBoxElement
                    for block in child.content:
                        if first:
                            para = tf.paragraphs[0]
                            first = False
                        else:
                            para = tf.add_paragraph()
                        para.alignment = map_alignment(getattr(block, 'alignment', 'left'))
                        if block.runs:
                            for r in block.runs:
                                run = para.add_run()
                                self._set_run_text(run, r.text)
                                if r.font_size_pt:
                                    run.font.size = Pt(float(r.font_size_pt))
                                if r.color:
                                    rgb = hex_to_rgbcolor(r.color)
                                    if rgb:
                                        run.font.color.rgb = rgb
                                run.font.bold = r.bold
                        else:
                            run = para.add_run()
                            self._set_run_text(run, block.text)

    def _render_numbered_step(self, slide, elem: NumberedStepElement,
                              slide_w: int, slide_h: int):
        """Render a numbered step with circle badge, title, and description."""
        region = self._resolve_region(elem, slide_w, slide_h)

        # Number circle
        circle_size = min(region.height, px_to_emu(64))

        # Circle shape
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Emu(region.left), Emu(region.top),
            Emu(circle_size), Emu(circle_size),
        )
        circle.fill.solid()
        rgb = hex_to_rgbcolor(elem.number_bg_color)
        if rgb:
            circle.fill.fore_color.rgb = rgb
        circle.line.fill.background()

        # Number text
        tf = circle.text_frame
        tf.word_wrap = False
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = str(elem.step_number)
        run.font.size = Pt(21)
        run.font.bold = True
        rgb = hex_to_rgbcolor(elem.number_text_color)
        if rgb:
            run.font.color.rgb = rgb

        # Title and description text box (to the right of the circle)
        text_left = region.left + circle_size + px_to_emu(16)
        text_w = region.width - circle_size - px_to_emu(16)
        if text_w > 0:
            txBox = slide.shapes.add_textbox(
                Emu(text_left), Emu(region.top),
                Emu(text_w), Emu(region.height),
            )
            tf2 = txBox.text_frame
            tf2.auto_size = MSO_AUTO_SIZE.NONE
            tf2.word_wrap = True

            # Title
            p_title = tf2.paragraphs[0]
            run_t = p_title.add_run()
            run_t.text = elem.title
            run_t.font.size = Pt(19.5)  # 26px * 0.75
            run_t.font.bold = True
            if elem.title_color:
                rgb = hex_to_rgbcolor(elem.title_color)
                if rgb:
                    run_t.font.color.rgb = rgb

            # Description
            if elem.description:
                p_desc = tf2.add_paragraph()
                run_d = p_desc.add_run()
                run_d.text = elem.description
                run_d.font.size = Pt(16.5)  # 22px * 0.75
                if elem.description_color:
                    rgb = hex_to_rgbcolor(elem.description_color)
                    if rgb:
                        run_d.font.color.rgb = rgb

        # Arrow (simple text arrow)
        if elem.show_arrow:
            arrow_left = region.left + circle_size // 2 - px_to_emu(16)
            arrow_top = region.top + region.height
            arrow_box = slide.shapes.add_textbox(
                Emu(arrow_left), Emu(arrow_top),
                Emu(px_to_emu(32)), Emu(px_to_emu(32)),
            )
            tf3 = arrow_box.text_frame
            tf3.auto_size = MSO_AUTO_SIZE.NONE
            run_a = tf3.paragraphs[0].add_run()
            run_a.text = "↓"
            run_a.font.size = Pt(24)
            if elem.arrow_color:
                rgb = hex_to_rgbcolor(elem.arrow_color)
                if rgb:
                    run_a.font.color.rgb = rgb

    def _render_divider_line(self, slide, elem: DividerLineElement,
                             slide_w: int, slide_h: int):
        """Render a horizontal divider line."""
        region = self._resolve_region(elem, slide_w, slide_h)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(region.left), Emu(region.top),
            Emu(region.width), Emu(max(px_to_emu(2), pt_to_emu(elem.height_pt))),
        )
        shape.fill.solid()
        rgb = hex_to_rgbcolor(elem.color)
        if rgb:
            shape.fill.fore_color.rgb = rgb
        shape.line.fill.background()

    def _render_badge(self, slide, elem: BadgeElement, slide_w: int, slide_h: int):
        """Render a pill/badge shape with optional icon and text."""
        region = self._resolve_region(elem, slide_w, slide_h)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Emu(region.left), Emu(region.top),
            Emu(region.width), Emu(region.height),
        )

        # Apply corner radius
        self._set_rounded_rect_adj(shape, elem.corner_radius_pct)

        # Fill
        if elem.background:
            self._apply_background(shape.fill, elem.background)
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0x13, 0x3E, 0xFF)

        # Opacity (must come after fill so <a:srgbClr> exists to attach alpha to)
        if elem.opacity < 1.0:
            self._set_shape_alpha(shape, elem.opacity)

        shape.line.fill.background()

        # Drop shadow (CSS box-shadow)
        self._apply_shape_shadow(shape, elem.shadow)

        # An icon-only badge (no text after icon-name filtering) is the common
        # case in slides 5/6: a coloured circle with a single Material Icon
        # centred via flexbox. In PPTX we centre the icon's custGeom shape in
        # the badge region and skip the text frame entirely — leaving the
        # default text frame in place would add an empty (but margin-padded)
        # text box that subtly offsets the badge's effective centre.
        has_text = bool(elem.text and elem.text.strip())
        icon_only = bool(elem.icon_name and elem.icon_font) and not has_text
        if icon_only:
            self._render_badge_icon(slide, elem, region, center=True)
            return

        # Text frame + (optional) left-aligned vector icon.
        tf = shape.text_frame
        tf.word_wrap = False
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER

        # Text run
        text_run = para.add_run()
        text_run.text = elem.text
        text_run.font.size = Pt(float(elem.font_size_pt))
        text_run.font.color.rgb = hex_to_rgbcolor(elem.font_color) or RGBColor(0xFF, 0xFF, 0xFF)

        if elem.icon_name and elem.icon_font:
            self._render_badge_icon(slide, elem, region, center=False)

    def _render_badge_icon(self, slide, elem: BadgeElement, badge_region,
                           center: bool = False):
        """Render a small vector icon shape inside the badge.

        If `center`, the icon is centred in the badge (icon-only badges —
        matches CSS `flex items-center justify-center`). Otherwise it sits
        on the left edge with 6pt padding (icon + text badges).
        """
        from shuttleslide.html_to_pptx.fonts import get_font_bytes, icon_to_codepoint
        from shuttleslide.html_to_pptx.vector_icon import (
            extract_glyph_outline,
            glyph_outline_to_custgeom_xml,
        )

        font_bytes = get_font_bytes(elem.icon_font)
        if not font_bytes:
            return
        codepoint = icon_to_codepoint(elem.icon_font, elem.icon_name)
        if codepoint is None:
            return
        outline = extract_glyph_outline(font_bytes, codepoint)
        if outline is None or not outline.commands:
            return

        cust_geom_xml = glyph_outline_to_custgeom_xml(outline)

        # Calculate icon size from font size (pt → EMU)
        icon_size_pt = float(elem.icon_size_pt or 14)
        icon_size_emu = Pt(icon_size_pt)

        # Adjust for custGeom aspect ratio
        path_aspect = outline.path_w / outline.path_h
        if path_aspect > 1.0:
            icon_w = icon_size_emu
            icon_h = round(icon_size_emu / path_aspect)
        else:
            icon_h = icon_size_emu
            icon_w = round(icon_size_emu * path_aspect)

        if center:
            icon_x = badge_region.left + (badge_region.width - icon_w) // 2
            icon_y = badge_region.top + (badge_region.height - icon_h) // 2
        else:
            # Left-aligned with padding (icon + text layout)
            padding = Pt(6)
            icon_x = badge_region.left + padding
            icon_y = badge_region.top + (badge_region.height - icon_h) // 2

        icon_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(icon_x), Emu(icon_y),
            Emu(icon_w), Emu(icon_h),
        )

        # Replace geometry with custGeom
        spPr = icon_shape._element.spPr
        prst_geom = spPr.find(qn("a:prstGeom"))
        if prst_geom is not None:
            spPr.replace(prst_geom, cust_geom_xml)
        else:
            spPr.insert(0, cust_geom_xml)

        # Fill with icon color
        if elem.icon_color:
            icon_shape.fill.solid()
            rgb = hex_to_rgbcolor(elem.icon_color)
            if rgb:
                icon_shape.fill.fore_color.rgb = rgb
            else:
                icon_shape.fill.solid()
                icon_shape.fill.fore_color.rgb = RGBColor(0, 0, 0)
        else:
            icon_shape.fill.solid()
            icon_shape.fill.fore_color.rgb = RGBColor(0, 0, 0)

        icon_shape.line.fill.background()

    def _render_bullet_list(self, slide, elem: BulletListElement,
                            slide_w: int, slide_h: int):
        """Render a bullet point list.

        Two rendering paths split on whether any item carries a custom
        marker (Material Icon, future custGeom marker, …):

        - **Custom-marker path** (``any_icon``): each item is decomposed
          into TWO independent shapes — the marker (icon shape or
          unicode fallback) and a dedicated textbox holding only that
          item's text. The Y cursor advances freely between items
          (``item_y += item_h + spacing``), so vertical rhythm is fully
          under our control. A single textbox with multi-paragraph
          content locks line spacing to PowerPoint's internal renderer
          and forces a hanging-indent model that overlaps icons ≥ the
          indent width — both failure modes disappear once each item
          owns its own box.

        - **Plain path** (no icons): one textbox, one paragraph per
          item, ``<a:buChar char="•">`` marker supplied by DrawingML.
          Zero regression from the pre-icon behaviour.

        INVARIANT: never inject the bullet as a ``"• "`` text run inside
        the paragraph. Doing so steals ~2% of frame width from the text
        area; combined with PPTX rendering text ~5-10% wider than the
        browser, this forces text that fit on one line in the browser
        (e.g. "Metal support enabled" in 3.html) to wrap in PPTX. The
        same invariant applies to any future marker-based element
        (numbered_step-as-paragraph, definition lists): use
        ``<a:buAutoNum>`` / etc., never text runs.

        SYSTEMIC DESIGN: the per-item-textbox architecture is the
        canonical layout for ANY custom bullet marker. Future marker
        types (numbered badges, custom shapes, emoji-as-marker) should
        reuse this path — render the marker via the primitives API,
        render the text in its own textbox, advance Y freely.
        """
        region = self._resolve_region(elem, slide_w, slide_h)
        line_h_emu = int(float(elem.font_size_pt) * 1.4 * 12700)
        spacing_emu = int(float(elem.spacing_pt) * 12700)
        any_icon = any(it.icon_name for it in elem.items)

        if any_icon:
            self._render_bullet_list_custom_markers(slide, elem, region,
                                                    line_h_emu, spacing_emu)
            return

        # ---- Plain buChar '•' path ------------------------------------
        # Single textbox, one paragraph per item, DrawingML supplies the
        # marker via <a:buChar>. Zero changes from pre-icon behaviour.
        txBox = slide.shapes.add_textbox(
            Emu(region.left), Emu(region.top),
            Emu(region.width), Emu(region.height),
        )
        tf = txBox.text_frame
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.word_wrap = True

        # ~1.875% of slide width, matching typical ``padding-left: 24px``
        # on a ``<ul>``. Sized for the '•' glyph (~half the font size).
        SLIDE_W_EMU = 9144000  # 10 inch * 914400 EMU/inch (16:9 default)
        HANGING_INDENT_EMU = int(SLIDE_W_EMU * 0.01875)
        bullet_indent_emu = HANGING_INDENT_EMU

        first = True
        for item in elem.items:
            if first:
                para = tf.paragraphs[0]
                first = False
            else:
                para = tf.add_paragraph()

            para.space_after = Pt(float(elem.spacing_pt))

            pPr = para._p.get_or_add_pPr()
            for tag in ('a:buNone', 'a:buChar', 'a:buAutoNum',
                        'a:buClr', 'a:buSzPct', 'a:buFont'):
                for existing in pPr.findall(qn(tag)):
                    pPr.remove(existing)

            pPr.set('marL', str(bullet_indent_emu))
            pPr.set('indent', str(-bullet_indent_emu))

            # CT_TextParagraphProperties sequence (ECMA-376):
            # buClr?, buSzPct?, (buNone | buChar | buAutoNum).
            if elem.bullet_color:
                buClr = etree.SubElement(pPr, qn('a:buClr'))
                srgb = etree.SubElement(buClr, qn('a:srgbClr'))
                srgb.set('val', elem.bullet_color.lstrip('#').upper())

            buSz = etree.SubElement(pPr, qn('a:buSzPct'))
            buSz.set('val', '100000')  # 100000 = 100%

            buChar = etree.SubElement(pPr, qn('a:buChar'))
            buChar.set('char', '•')

            text_run = para.add_run()
            text_run.text = item.text
            text_run.font.size = Pt(float(elem.font_size_pt))
            if elem.font_color:
                rgb = hex_to_rgbcolor(elem.font_color)
                if rgb:
                    text_run.font.color.rgb = rgb

    def _render_bullet_list_custom_markers(
        self, slide, elem: BulletListElement, region: Region,
        line_h_emu: int, spacing_emu: int,
    ) -> None:
        """Render a bullet list where ≥1 item carries a custom marker.

        Each item becomes two independent shapes:
          1. marker shape (custGeom vector icon, or unicode fallback
             textbox) at ``(region.left, item_y + vcenter_on_line_1)``
          2. text textbox at ``(region.left + icon_col_w, item_y,
             text_w, item_h)`` with buNone so PowerPoint draws no
             marker of its own.

        The Y cursor advances per item by ``item_h + spacing_emu``,
        where ``item_h`` accounts for estimated wrap lines — so a
        wrapping item grows its box instead of overflowing into the
        next item's region (which is what happens with a single shared
        textbox locked to PowerPoint's internal line height).

        Wrap estimate uses an average char width of ``font_size * 0.55``
        — rough, but the only way to size the per-item textbox without
        reading PowerPoint's post-render geometry back. Degrades
        gracefully: overestimates lines for narrow text (small
        vertical gap), underestimates only for pathologically wide
        wrap (rare in Slide HTML subset).
        """
        # Icon column width = widest icon + 6pt gap to text.
        max_icon_w_emu = max(
            (int(float(it.icon_size_pt or elem.font_size_pt) * 12700)
             for it in elem.items if it.icon_name),
            default=0,
        )
        icon_gap_emu = int(6 * 12700)
        icon_col_w = max_icon_w_emu + icon_gap_emu
        text_x = region.left + icon_col_w
        text_w = max(region.width - icon_col_w, int(50 * 12700))

        # Average char width for wrap estimate. PPTX renders ~5-10%
        # wider than the browser, so 0.55× font_size is conservative.
        font_size_pt = float(elem.font_size_pt)
        approx_char_w_emu = max(1, int(font_size_pt * 0.55 * 12700))
        chars_per_line = max(1, text_w // approx_char_w_emu)

        def _est_wrap_lines(text: str) -> int:
            if not text:
                return 1
            return max(1, (len(text) + chars_per_line - 1) // chars_per_line)

        item_y = region.top
        for item in elem.items:
            n_lines = _est_wrap_lines(item.text)
            item_h = n_lines * line_h_emu

            # ---- Marker (icon column) ------------------------------
            if item.icon_name:
                icon_size_pt = float(item.icon_size_pt or elem.font_size_pt)
                icon_box_emu = int(icon_size_pt * 12700)
                # Vertically centre on the FIRST line of text — the
                # marker visually belongs to line 1, even if the item
                # wraps to multiple lines.
                icon_y = item_y + (line_h_emu - icon_box_emu) // 2
                icon_region = Region(
                    left=region.left,
                    top=icon_y,
                    width=icon_box_emu,
                    height=icon_box_emu,
                )
                shape = self._render_vector_icon_primitives(
                    slide,
                    icon_name=item.icon_name,
                    icon_font=item.icon_font,
                    icon_color=item.icon_color,
                    icon_size_pt=icon_size_pt,
                    icon_shadow=None,
                    region=icon_region,
                )
                if shape is None:
                    # Unicode fallback (✓ ⚠ …) — never fall through to
                    # the raw icon name as display text.
                    self._render_icon_unicode_fallback(
                        slide, item.icon_name, item.icon_color,
                        icon_size_pt, icon_region,
                    )

            # ---- Text (independent textbox) ------------------------
            txBox = slide.shapes.add_textbox(
                Emu(text_x), Emu(item_y),
                Emu(text_w), Emu(item_h),
            )
            tf = txBox.text_frame
            tf.auto_size = MSO_AUTO_SIZE.NONE
            tf.word_wrap = True
            # Zero margins: textbox bounds == text column bounds. The
            # icon_col_w already encodes the visual gap, so any inner
            # margin would compound it.
            tf.margin_left = Emu(0)
            tf.margin_right = Emu(0)
            tf.margin_top = Emu(0)
            tf.margin_bottom = Emu(0)
            tf.vertical_anchor = MSO_ANCHOR.TOP

            para = tf.paragraphs[0]
            # buNone on every paragraph so PowerPoint draws no marker;
            # the icon shape above is the visual marker.
            pPr = para._p.get_or_add_pPr()
            for tag in ('a:buNone', 'a:buChar', 'a:buAutoNum',
                        'a:buClr', 'a:buSzPct', 'a:buFont'):
                for existing in pPr.findall(qn(tag)):
                    pPr.remove(existing)
            etree.SubElement(pPr, qn('a:buNone'))

            text_run = para.add_run()
            text_run.text = item.text
            text_run.font.size = Pt(float(elem.font_size_pt))
            if elem.font_color:
                rgb = hex_to_rgbcolor(elem.font_color)
                if rgb:
                    text_run.font.color.rgb = rgb

            item_y += item_h + spacing_emu

    def _render_table(self, slide, elem: TableElement, theme: ThemeDef,
                      slide_w: int, slide_h: int):
        """Render a TableElement as a native PPTX table.

        Mirrors the cell text styling of _render_text_box (per-run color /
        bold / size / font_name). Cell fills come from each cell's
        `background` (usually a row tint from the row div). Borders are
        injected via DrawingML XML because python-pptx has no cell-border
        API.
        """
        if not elem.rows:
            return

        region = self._resolve_region(elem, slide_w, slide_h)
        n_rows = len(elem.rows)
        n_cols = len(elem.rows[0]) if elem.rows else 0
        if n_cols == 0:
            return

        graphic_frame = slide.shapes.add_table(
            n_rows, n_cols,
            Emu(region.left), Emu(region.top),
            Emu(region.width), Emu(region.height),
        )
        table = graphic_frame.table

        # Apply "No Style, No Grid" so the default blue banding doesn't
        # clash with our explicit fills/borders. GUID from
        # python-pptx table styles; firstRow/bandRow disabled too.
        self._apply_plain_table_style(table)

        # Column widths — distribute table width across columns per their
        # relative weights. Falls back to equal widths if unspecified.
        if elem.col_widths_pct and len(elem.col_widths_pct) == n_cols:
            total = sum(elem.col_widths_pct) or 1.0
            for j in range(n_cols):
                table.columns[j].width = Emu(
                    int(region.width * (elem.col_widths_pct[j] / total))
                )
        else:
            col_w = region.width // n_cols
            for j in range(n_cols):
                table.columns[j].width = Emu(col_w)

        # Row heights
        if elem.row_heights_pct and len(elem.row_heights_pct) == n_rows:
            total_h = sum(elem.row_heights_pct) or 1.0
            for i in range(n_rows):
                table.rows[i].height = Emu(
                    int(region.height * (elem.row_heights_pct[i] / total_h))
                )

        # Body cells
        for i, row in enumerate(elem.rows):
            for j, cell_dsl in enumerate(row):
                if j >= n_cols:
                    break
                cell = table.cell(i, j)
                self._write_table_cell(cell, cell_dsl)

        # Borders — apply AFTER cell fills so XML ordering stays correct.
        self._apply_table_borders(table, elem, n_rows, n_cols)

    def _apply_plain_table_style(self, table):
        """Force a 'No Style, No Grid' look on the table.

        Sets the tableStyleId to the No Style / No Grid GUID and disables
        firstRow / bandRow special-format flags so they don't override our
        explicit cell fills.
        """
        tbl = table._tbl
        tblPr = tbl.find(qn('a:tblPr'))
        if tblPr is None:
            tblPr = etree.SubElement(tbl, qn('a:tblPr'))
            tbl.insert(0, tblPr)
        tblPr.set('firstRow', '0')
        tblPr.set('bandRow', '0')

        # Replace (or add) the tableStyleId element with the No Style,
        # No Grid GUID: {2D5ABB26-0587-4C30-8999-92F81FD0307C}.
        for style_id in tblPr.findall(qn('a:tableStyleId')):
            tblPr.remove(style_id)
        style_id = etree.SubElement(tblPr, qn('a:tableStyleId'))
        style_id.text = '{2D5ABB26-0587-4C30-8999-92F81FD0307C}'

    def _write_table_cell(self, cell, cell_dsl):
        """Style one python-pptx cell: fill, margins, text + runs."""
        # Cell fill
        if cell_dsl.background and cell_dsl.background.color:
            cell.fill.solid()
            rgb = hex_to_rgbcolor(cell_dsl.background.color)
            if rgb:
                cell.fill.fore_color.rgb = rgb
            opacity = color_opacity(cell_dsl.background.color)
            if opacity < 1.0:
                # python-pptx has no cell-fill alpha API; inject <a:alpha>.
                self._apply_cell_fill_alpha(cell, opacity)
        else:
            cell.fill.background()  # transparent

        # Tight margins to mirror a CSS table-cell
        cell.margin_left = Emu(45720)    # 0.05 inch
        cell.margin_right = Emu(45720)
        cell.margin_top = Emu(27432)     # 0.03 inch
        cell.margin_bottom = Emu(27432)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        tf = cell.text_frame
        tf.word_wrap = True
        # Clear the default empty paragraph's runs before writing
        para = tf.paragraphs[0]
        para.alignment = map_alignment(cell_dsl.alignment)
        para.space_before = Pt(0)
        para.space_after = Pt(0)

        first = True
        if cell_dsl.runs:
            for r in cell_dsl.runs:
                if first:
                    run = para.add_run()
                    first = False
                else:
                    # Reuse the same paragraph — multiple runs are inline
                    # siblings (e.g. <strong>bold</strong> inside a cell).
                    run = para.add_run()
                self._set_run_text(run, r.text)
                if r.font_size_pt:
                    run.font.size = Pt(float(r.font_size_pt))
                if r.color:
                    rgb = hex_to_rgbcolor(r.color)
                    if rgb:
                        run.font.color.rgb = rgb
                run.font.bold = r.bold
                run.font.italic = r.italic
                if r.font_name:
                    _apply_run_fonts(run, r.font_name)
        elif cell_dsl.text:
            run = para.add_run()
            run.text = cell_dsl.text

    @staticmethod
    def _apply_cell_fill_alpha(cell, opacity: float):
        """Inject <a:alpha> into a table cell's solidFill srgbClr."""
        tcPr = cell._tc.find(qn('a:tcPr'))
        if tcPr is None:
            return
        srgbClr = tcPr.find(qn('a:solidFill') + '/' + qn('a:srgbClr'))
        if srgbClr is None:
            return
        for a in srgbClr.findall(qn('a:alpha')):
            srgbClr.remove(a)
        alpha_elem = etree.SubElement(srgbClr, qn('a:alpha'))
        alpha_elem.set('val', str(int(opacity * 100000)))

    def _apply_table_borders(self, table, elem: TableElement,
                             n_rows: int, n_cols: int):
        """Apply table outline + row-separator borders via DrawingML XML.

        - `border` (table outline): perimeter — left on col 0, right on
          last col, top on row 0, bottom on last row.
        - `row_separator`: bottom edge on every row except the last
          (the last row's bottom is part of the perimeter outline).
        """
        outline = elem.border
        sep = elem.row_separator
        if not outline and not sep:
            return

        for i in range(n_rows):
            for j in range(n_cols):
                cell = table.cell(i, j)
                if outline and outline.width_pt > 0:
                    if j == 0:
                        self._apply_cell_edge(cell, 'L', outline)
                    if j == n_cols - 1:
                        self._apply_cell_edge(cell, 'R', outline)
                    if i == 0:
                        self._apply_cell_edge(cell, 'T', outline)
                    if i == n_rows - 1:
                        self._apply_cell_edge(cell, 'B', outline)
                if sep and sep.width_pt > 0 and i < n_rows - 1:
                    self._apply_cell_edge(cell, 'B', sep)

    @staticmethod
    def _apply_cell_edge(cell, edge: str, border: BorderDef):
        """Inject <a:ln{edge}> into the cell's <a:tcPr>.

        `edge` is one of 'L', 'R', 'T', 'B'. Per OpenXML schema the line
        element order inside tcPr is lnL, lnR, lnT, lnB, followed by fill
        children — but python-pptx inserts solidFill when cell.fill.solid()
        was called, so we add the line element to tcPr manually without
        worrying about strict schema ordering (PowerPoint is tolerant).
        """
        edge_to_tag = {'L': 'a:lnL', 'R': 'a:lnR', 'T': 'a:lnT', 'B': 'a:lnB'}
        tag = edge_to_tag.get(edge)
        if tag is None:
            return
        tcPr = cell._tc.find(qn('a:tcPr'))
        if tcPr is None:
            tcPr = etree.SubElement(cell._tc, qn('a:tcPr'))
            cell._tc.append(tcPr)

        # Remove any existing line element for this edge
        for el in tcPr.findall(qn(tag)):
            tcPr.remove(el)

        ln = etree.SubElement(tcPr, qn(tag))
        # width is in EMU; 1pt = 12700 EMU
        ln.set('w', str(int(float(border.width_pt) * 12700)))
        # capStyle / cmpdStyle are optional but PowerPoint tolerates omission.

        # Fill inside the line — solidFill srgbClr
        rgb = hex_to_rgbcolor(border.color)
        if rgb is None:
            rgb = RGBColor(0, 0, 0)
        solid_fill = etree.SubElement(ln, qn('a:solidFill'))
        srgb = etree.SubElement(solid_fill, qn('a:srgbClr'))
        srgb.set('val', str(rgb))
        opacity = color_opacity(border.color)
        if opacity < 1.0:
            alpha = etree.SubElement(srgb, qn('a:alpha'))
            alpha.set('val', str(int(opacity * 100000)))

        # Dash style for dashed/dotted CSS borders
        if border.style in ("dashed", "dotted"):
            val = "dash" if border.style == "dashed" else "dot"
            prst = etree.SubElement(ln, qn('a:prstDash'))
            prst.set('val', val)

    # ------------------------------------------------------------------
    # Fill helpers
    # ------------------------------------------------------------------

    def _apply_background(self, fill, bg: BackgroundDef):
        """Apply a BackgroundDef to a shape fill."""
        if bg.type == "solid":
            fill.solid()
            rgb = hex_to_rgbcolor(bg.color or "#FFFFFF")
            if rgb:
                fill.fore_color.rgb = rgb
            # CSS rgba() backgrounds are encoded as #RRGGBBAA in extract_layout.js.
            # python-pptx's RGBColor drops the alpha channel, so inject <a:alpha>
            # into the solidFill XML when the source color was semi-transparent.
            if bg.color:
                opacity = color_opacity(bg.color)
                if opacity < 1.0:
                    self._apply_solid_fill_alpha(fill, opacity)
        elif bg.type == "gradient" and bg.gradient:
            self._apply_gradient_fill(fill, bg.gradient)
        else:
            fill.background()

    @staticmethod
    def _apply_solid_fill_alpha(fill, opacity: float):
        """Inject <a:alpha val="..."/> into a solid fill's <a:srgbClr>.

        Must be called after fill.solid() + fore_color.rgb assignment so the
        <a:solidFill><a:srgbClr/></a:solidFill> subtree already exists.
        """
        # python-pptx wraps solid fill as _SolidFill which has no _element.
        # Access the underlying <a:solidFill> via fore_color._xFill instead.
        solidFill = getattr(getattr(fill, 'fore_color', None), '_xFill', None)
        if solidFill is None:
            return
        srgbClr = solidFill.find(qn('a:srgbClr'))
        if srgbClr is None:
            return
        for a in srgbClr.findall(qn('a:alpha')):
            srgbClr.remove(a)
        alpha_elem = etree.SubElement(srgbClr, qn('a:alpha'))
        alpha_elem.set('val', str(int(opacity * 100000)))

    def _apply_gradient_fill(self, fill, gradient: GradientDef,
                             global_opacity: float = 1.0):
        """Apply a gradient fill to a shape, supporting any number of stops."""
        fill.gradient()
        stops = gradient.stops
        if not stops:
            return

        # Access the <a:gradFill> XML element
        grad_fill_elem = fill._fill._element

        # Clear all existing children to ensure correct schema order
        for child in list(grad_fill_elem):
            grad_fill_elem.remove(child)

        # 1) Rebuild <a:gsLst> FIRST — OpenXML schema: gsLst, lin, path
        gsLst = etree.SubElement(grad_fill_elem, qn('a:gsLst'))

        for stop in stops:
            pos_val = str(int(stop.position * 100000))
            gs_elem = etree.SubElement(gsLst, qn('a:gs'))
            gs_elem.set('pos', pos_val)

            rgb = hex_to_rgbcolor(stop.color)
            if rgb is None:
                rgb = RGBColor(0, 0, 0)
            srgbClr = etree.SubElement(gs_elem, qn('a:srgbClr'))
            srgbClr.set('val', str(rgb))

            opacity = stop.opacity * global_opacity
            if opacity < 1.0:
                alpha = etree.SubElement(srgbClr, qn('a:alpha'))
                alpha.set('val', str(int(opacity * 100000)))

        # 2) Add <a:lin> after <a:gsLst>
        if gradient.direction:
            angle = gradient_angle_deg(gradient.direction)
            lin_elem = etree.SubElement(grad_fill_elem, qn('a:lin'))
            lin_elem.set('ang', str(int(angle * 60000)))
            lin_elem.set('scaled', '1')

    # ------------------------------------------------------------------
    # Alpha / transparency helpers
    # ------------------------------------------------------------------

    def _set_shape_alpha(self, shape, opacity: float):
        """Set overall opacity on a shape's fill (0.0=invisible, 1.0=opaque)."""
        sp = shape._element
        srgbClr = sp.find('.//' + qn('a:srgbClr'))
        if srgbClr is not None:
            # Remove existing alpha
            for a in srgbClr.findall(qn('a:alpha')):
                srgbClr.remove(a)
            alpha_elem = etree.SubElement(srgbClr, qn('a:alpha'))
            alpha_elem.set('val', str(int(opacity * 100000)))

    def _apply_text_run_alpha(self, run, opacity: float):
        """Inject <a:alpha> into a text run's solidFill color.

        Mirrors _set_shape_alpha but targets run text color instead of shape
        fill. Must be called AFTER run.font.color.rgb is set, so
        <a:solidFill>/<a:srgbClr> already exist in the run's <a:rPr>.
        Used for semi-transparent text (watermarks, ghosted section numbers,
        hint text) — value is in thousandths of a percent (8000 = 8%).
        """
        if opacity >= 1.0:
            return
        rPr = run._r.get_or_add_rPr()
        solidFill = rPr.find(qn('a:solidFill'))
        if solidFill is None:
            return
        srgbClr = solidFill.find(qn('a:srgbClr'))
        if srgbClr is None:
            return
        for a in srgbClr.findall(qn('a:alpha')):
            srgbClr.remove(a)
        alpha_elem = etree.SubElement(srgbClr, qn('a:alpha'))
        alpha_elem.set('val', str(int(opacity * 100000)))

    def _set_run_text(self, run, text: str):
        """Set a text run's text, preserving leading/trailing whitespace.

        PowerPoint (DrawingML) follows the XML spec's default whitespace
        handling for element content: leading and trailing whitespace on
        ``<a:t>`` is trimmed unless ``xml:space="preserve"`` is set. This
        matters for inline runs split across HTML spans — the text node
        between ``</span>`` and following text starts with a meaningful
        space (``"--editable"`` + ``" flag enables..."``). Without
        ``xml:space="preserve"`` PowerPoint collapses them to
        ``"--editableflag"``.
        """
        run.text = text
        if text and (text[0] in ' \t\n' or text[-1] in ' \t\n'):
            t_elem = run._r.find(qn('a:t'))
            if t_elem is not None:
                t_elem.set('{http://www.w3.org/XML/1998/namespace}space',
                           'preserve')

    def _set_gradient_stop_alpha(self, gs_stop, opacity: float):
        """Set alpha on a gradient stop."""
        # gs_stop._gs is the <a:gs> XML element; <a:srgbClr> is a direct child
        gs_elem = gs_stop._gs if hasattr(gs_stop, '_gs') else gs_stop._element
        srgbClr = gs_elem.find(qn('a:srgbClr'))
        if srgbClr is None:
            return
        for a in srgbClr.findall(qn('a:alpha')):
            srgbClr.remove(a)
        alpha_elem = etree.SubElement(srgbClr, qn('a:alpha'))
        alpha_elem.set('val', str(int(opacity * 100000)))

    @staticmethod
    def _set_soft_edge(shape, radius_pt: float):
        """Apply soft edge effect to simulate CSS blur.

        Args:
            shape: python-pptx shape.
            radius_pt: soft edge radius in points.
        """
        if radius_pt <= 0:
            return
        sp = shape._element
        spPr = sp.find(qn('p:spPr'))
        if spPr is None:
            return
        # Find or create effectLst
        effectLst = spPr.find(qn('a:effectLst'))
        if effectLst is None:
            effectLst = etree.SubElement(spPr, qn('a:effectLst'))
        # Remove existing softEdge
        for se in effectLst.findall(qn('a:softEdge')):
            effectLst.remove(se)
        softEdge = etree.SubElement(effectLst, qn('a:softEdge'))
        # rad is in EMUs; 1pt = 12700 EMU
        rad_emu = int(radius_pt * 12700)
        softEdge.set('rad', str(rad_emu))

    @staticmethod
    def _apply_shape_shadow(shape, shadow):
        """Apply an outer shadow to a shape from a ShadowDef.

        Mirrors CSS `box-shadow: Xpx Ypx BlurPx Color` by emitting an
        ``<a:outerShdw>`` element under the shape's ``<a:effectLst>``.
        DrawingML angles: dir is measured clockwise from +X in 1/60000
        degree units, so ``atan2(offset_y, offset_x)`` converts the
        CSS offset pair directly. blurRad is in EMU (1pt = 12700).
        The colour's alpha channel becomes the ``<a:alpha>`` child so
        rgba() shadows carry their opacity through.

        No-op when ``shadow`` is None — callers can pass the element's
        shadow field unconditionally.
        """
        if shadow is None:
            return
        from shuttleslide.html_to_pptx.style_mapper import color_opacity, parse_hex_color

        sp = shape._element
        spPr = sp.find(qn('p:spPr'))
        if spPr is None:
            return
        effectLst = spPr.find(qn('a:effectLst'))
        if effectLst is None:
            effectLst = etree.SubElement(spPr, qn('a:effectLst'))
        # Remove any existing outerShdw to avoid stacking on re-render.
        for os_ in effectLst.findall(qn('a:outerShdw')):
            effectLst.remove(os_)

        outerShdw = etree.SubElement(effectLst, qn('a:outerShdw'))
        outerShdw.set('blurRad', str(int(max(shadow.blur_pt, 0) * 12700)))
        outerShdw.set('rotWithShape', '0')

        # Direction: atan2(y, x) gives degrees clockwise from +X.
        import math
        dir_deg = math.degrees(math.atan2(shadow.offset_y_pt, shadow.offset_x_pt))
        outerShdw.set('dir', str(int(dir_deg * 60000)))

        # Distance: Euclidean magnitude of the offset, in EMU.
        dist_pt = math.hypot(shadow.offset_x_pt, shadow.offset_y_pt)
        outerShdw.set('dist', str(int(dist_pt * 12700)))

        # Colour with alpha. parse_hex_color returns (r, g, b, a).
        parsed = parse_hex_color(shadow.color)
        if parsed is None:
            hex_triplet = "000000"
            alpha = 1.0
        else:
            r, g, b, alpha = parsed
            hex_triplet = f"{r:02X}{g:02X}{b:02X}"
        srgbClr = etree.SubElement(outerShdw, qn('a:srgbClr'))
        srgbClr.set('val', hex_triplet)
        alpha_elem = etree.SubElement(srgbClr, qn('a:alpha'))
        alpha_elem.set('val', str(int(max(alpha, 0.0) * 100000)))

    @staticmethod
    def _apply_object_fit_cover(pic, image_bytes: bytes, region: Region):
        """Crop a picture to fill its box, preserving aspect ratio.

        Mirrors CSS `object-fit: cover`: the image is scaled so the shorter
        dimension fills the box and the overflow on the longer dimension is
        cropped evenly from both sides.

        Args:
            pic: python-pptx Picture.
            image_bytes: raw image bytes (used to read native dimensions).
            region: target rendering region in EMUs.
        """
        try:
            from PIL import Image as _PILImage
        except ImportError:
            return
        try:
            with _PILImage.open(io.BytesIO(image_bytes)) as im:
                img_w, img_h = im.size
        except Exception:
            return
        if img_w <= 0 or img_h <= 0 or region.width <= 0 or region.height <= 0:
            return

        img_aspect = img_w / img_h
        box_aspect = region.width / region.height

        if img_aspect > box_aspect:
            # Image is wider than the box — crop horizontally.
            crop = (1 - box_aspect / img_aspect) / 2
            pic.crop_left = crop
            pic.crop_right = crop
        elif img_aspect < box_aspect:
            # Image is taller than the box — crop vertically.
            crop = (1 - img_aspect / box_aspect) / 2
            pic.crop_top = crop
            pic.crop_bottom = crop

    @staticmethod
    def _apply_object_fit_contain(pic, image_bytes: bytes, region: Region):
        """Letterbox a picture inside its box, preserving aspect ratio.

        Mirrors CSS `object-fit: contain`: the image is scaled so the longer
        dimension fits inside the box and the shorter dimension is centred
        with empty space on either side. Implemented by overriding the
        picture's position and size after creation (python-pptx has no
        native letterbox mode).
        """
        try:
            from PIL import Image as _PILImage
        except ImportError:
            return
        try:
            with _PILImage.open(io.BytesIO(image_bytes)) as im:
                img_w, img_h = im.size
        except Exception:
            return
        if img_w <= 0 or img_h <= 0 or region.width <= 0 or region.height <= 0:
            return

        scale = min(region.width / img_w, region.height / img_h)
        new_w = int(img_w * scale)
        new_h = int(img_h * scale)
        # Centre inside the region.
        pic.left = Emu(region.left + (region.width - new_w) // 2)
        pic.top = Emu(region.top + (region.height - new_h) // 2)
        pic.width = Emu(new_w)
        pic.height = Emu(new_h)

    # ------------------------------------------------------------------
    # ------------------------------------------------------------------
    # Shape helpers
    # ------------------------------------------------------------------

    @staticmethod
    def _apply_border_dash_style(shape, style: str):
        """Inject <a:prstDash val="..."/> into the shape's <a:ln> element.

        Per OpenXML schema (CT_LineProperties), prstDash comes after the fill
        child inside <a:ln>. Maps CSS border-style 'dashed'/'dotted' to
        preset values 'dash'/'dot'. Solid (default) leaves <a:ln> untouched.
        """
        if style not in ("dashed", "dotted"):
            return
        sp = shape._element
        ln = sp.find('.//' + qn('a:ln'))
        if ln is None:
            return
        # Remove any existing dash declarations
        for tag in ('a:prstDash', 'a:custDash'):
            for el in ln.findall(qn(tag)):
                ln.remove(el)
        val = "dash" if style == "dashed" else "dot"
        prst = etree.SubElement(ln, qn('a:prstDash'))
        prst.set('val', val)

    @staticmethod
    def _set_rounded_rect_adj(shape, corner_radius_pct: float):
        """Set the adjustment value (corner radius) of a ROUNDED_RECTANGLE.

        Args:
            shape: python-pptx shape (must have prstGeom).
            corner_radius_pct: 0.0 – 0.5 (fraction of shorter dimension).
        """
        if corner_radius_pct <= 0:
            return
        adj = int(min(corner_radius_pct, 0.5) * 100000)
        sp = shape._element
        # Look for prstGeom — may be under spPr or directly
        prst_geom = sp.find('.//' + qn('a:prstGeom'))
        if prst_geom is None:
            return
        avLst = prst_geom.find(qn('a:avLst'))
        if avLst is None:
            avLst = etree.SubElement(prst_geom, qn('a:avLst'))
        # Remove existing adj
        for gd in avLst.findall(qn('a:gd')):
            if gd.get('name') == 'adj':
                avLst.remove(gd)
        gd = etree.SubElement(avLst, qn('a:gd'))
        gd.set('name', 'adj')
        gd.set('fmla', f'val {adj}')

    # ------------------------------------------------------------------
    # Region resolution
    # ------------------------------------------------------------------

    def _resolve_region(self, elem, slide_w: int, slide_h: int) -> Region:
        """Resolve an element to its EMU Region using its position percentages.

        Falls back to the full slide if the element has no explicit position.
        """
        if elem.position is not None:
            p = elem.position
            return position_percent_to_region(
                p.x_pct, p.y_pct, p.w_pct, p.h_pct, slide_w, slide_h,
            )
        return Region(0, 0, slide_w, slide_h)
