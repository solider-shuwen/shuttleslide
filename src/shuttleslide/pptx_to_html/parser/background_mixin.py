"""
Background parsing mixin for PPTXParser.

Handles slide background extraction following the OpenXML inheritance chain:
slide -> layout -> master -> default white.
"""

from typing import Optional

from pptx.enum.shapes import MSO_SHAPE_TYPE

from shuttleslide.pptx_to_html.models import SlideBackground
from shuttleslide.pptx_to_html.utils.namespaces import NS_R_CLARK
from shuttleslide.pptx_to_html.utils.units import angle_to_degrees


class BackgroundMixin:
    """Background extraction methods."""

    def _parse_bg_element(self, bg_elem) -> Optional[SlideBackground]:
        """
        Parse a <p:bg> XML element and return a SlideBackground.

        Handles bgPr (background properties) and bgRef (background reference).

        Args:
            bg_elem: <p:bg> XML element

        Returns:
            SlideBackground or None if background is empty/transparent
        """
        ns = self._ns

        # Check for <p:bgPr> (direct background properties)
        bg_pr = bg_elem.find('p:bgPr', ns)
        if bg_pr is not None:
            return self._parse_fill_from_element(bg_pr)

        # Check for <p:bgRef> (reference to theme background)
        bg_ref = bg_elem.find('p:bgRef', ns)
        if bg_ref is not None:
            idx = bg_ref.get('idx')
            if idx == '0':
                # idx=0 means empty/no fill
                return None

            # bgRef contains a color directly (common pattern)
            # e.g., <p:bgRef idx="1001"><a:schemeClr val="bg1"/></p:bgRef>
            color = self._resolve_xml_color(bg_ref)
            if color:
                return SlideBackground(bg_type='solid', color=color)

        return None

    def _parse_fill_from_element(self, elem) -> Optional[SlideBackground]:
        """
        Parse fill properties from an XML element (bgPr, etc.).

        Handles solidFill, gradFill, and blipFill.

        Args:
            elem: XML element containing fill properties

        Returns:
            SlideBackground or None
        """
        ns = self._ns

        # Solid fill
        solid_fill = elem.find('a:solidFill', ns)
        if solid_fill is not None:
            color = self._resolve_xml_color(solid_fill)
            if color:
                return SlideBackground(bg_type='solid', color=color)

        # Gradient fill
        grad_fill = elem.find('a:gradFill', ns)
        if grad_fill is not None:
            return self._parse_gradient_fill(grad_fill)

        # Image/blip fill
        blip_fill = elem.find('a:blipFill', ns)
        if blip_fill is not None:
            return self._parse_blip_fill_background(blip_fill)

        return None

    def _parse_gradient_fill(self, grad_fill) -> Optional[SlideBackground]:
        """
        Parse <a:gradFill> and convert to CSS linear-gradient.

        Args:
            grad_fill: <a:gradFill> XML element

        Returns:
            SlideBackground with gradient CSS
        """
        ns = self._ns

        # Parse gradient stops
        gs_lst = grad_fill.find('a:gsLst', ns)
        if gs_lst is None:
            return None

        stops = []
        for gs in gs_lst.findall('a:gs', ns):
            pos = gs.get('pos', '0')
            pos_pct = int(pos) / 1000  # XML uses 0-100000, CSS uses 0-100
            color = self._resolve_xml_color(gs)
            if color:
                stops.append(f"{color} {pos_pct}%")

        if not stops:
            return None

        # Parse angle from <a:lin> element
        # OpenXML angle: 60000ths of a degree, counter-clockwise from positive X-axis
        # CSS angle: degrees, clockwise from positive Y-axis (north)
        css_angle = 180  # default: bottom-to-top
        lin = grad_fill.find('a:lin', ns)
        if lin is not None:
            ang = lin.get('ang')
            if ang:
                xml_angle = angle_to_degrees(int(ang))  # Convert to degrees
                # OpenXML: 0 = left-to-right, CSS: 0 = bottom-to-top
                # Conversion: css_angle = (xml_angle - 90) % 360 ... actually:
                # OpenXML measures from X-axis (East) counter-clockwise
                # CSS measures from Y-axis (North) clockwise
                # css_angle = 90 - xml_angle, normalized to 0-360
                css_angle = (90 - xml_angle) % 360

        gradient_css = f"linear-gradient({css_angle:.1f}deg, {', '.join(stops)})"
        return SlideBackground(bg_type='gradient', gradient_css=gradient_css)

    def _parse_blip_fill_background(self, blip_fill) -> Optional[SlideBackground]:
        """
        Parse <a:blipFill> for background image.

        Args:
            blip_fill: <a:blipFill> XML element

        Returns:
            SlideBackground with image data or None
        """
        # This is handled at the slide level where we have access to relationships
        # For now, store a marker; the actual extraction happens in _extract_slide_background
        return None

    def _extract_slide_background(self, slide) -> Optional[SlideBackground]:
        """
        Extract background for a slide following the OpenXML inheritance chain:
        slide -> layout -> master -> default white.

        Args:
            slide: python-pptx slide object

        Returns:
            SlideBackground or None (defaults to white in layout)
        """
        ns = self._ns

        # Check if slide hides master background/shapes
        slide_elem = slide._element
        show_master_sp = slide_elem.get('showMasterSp')

        # Priority 1: Full-slide picture shapes on layout/master
        # Many presentations use a <p:pic> element covering the entire slide
        # as the visual background. These are drawn ON TOP of <p:bg>, so they
        # are what the user actually sees.
        if show_master_sp != '0':
            image_bg = self._find_background_image_shape(slide)
            if image_bg:
                return image_bg

        # Priority 2: <p:bg> element (background fill definition)
        # Check slide's own background
        bg = slide_elem.find('.//p:bg', ns)
        if bg is not None:
            # Check for blipFill (image background) - needs relationship access
            bg_pr = bg.find('p:bgPr', ns)
            if bg_pr is not None:
                blip_fill = bg_pr.find('a:blipFill', ns)
                if blip_fill is not None:
                    image_bg = self._extract_bg_blip_fill(slide, blip_fill)
                    if image_bg:
                        return image_bg

            result = self._parse_bg_element(bg)
            if result:
                return result

        # Check layout/master backgrounds
        if show_master_sp != '0':
            try:
                layout = slide.slide_layout
                if layout is not None:
                    layout_elem = layout._element
                    layout_bg = layout_elem.find('.//p:bg', ns)
                    if layout_bg is not None:
                        result = self._parse_bg_element(layout_bg)
                        if result:
                            return result

                    try:
                        master = layout.slide_master
                        master_elem = master._element
                        master_bg = master_elem.find('.//p:bg', ns)
                        if master_bg is not None:
                            result = self._parse_bg_element(master_bg)
                            if result:
                                return result
                    except Exception:
                        pass
            except Exception:
                pass

        return None

    def _extract_bg_blip_fill(self, slide, blip_fill) -> Optional[SlideBackground]:
        """
        Extract image data from a background blipFill element.

        Args:
            slide: python-pptx slide object
            blip_fill: <a:blipFill> XML element

        Returns:
            SlideBackground with image data or None
        """
        ns = self._ns
        try:
            blip = blip_fill.find('.//a:blip', ns)
            if blip is None:
                return None

            embed_id = blip.get(f'{NS_R_CLARK}embed')
            if not embed_id:
                return None

            rel = slide.part.rels[embed_id]
            image_part = rel.target_part
            image_bytes = image_part.blob
            content_type = image_part.content_type
            image_type = content_type.split('/')[-1] if '/' in content_type else 'png'
            if image_type == 'jpeg':
                image_type = 'jpg'

            return SlideBackground(
                bg_type='image',
                image_data={
                    'image_bytes': image_bytes,
                    'image_type': image_type,
                }
            )
        except Exception:
            return None

    def _find_background_image_shape(self, slide) -> Optional[SlideBackground]:
        """
        Find full-slide picture shapes and overlay shapes on the layout/master
        that serve as the visual background.

        Many presentations place background images as <p:pic> elements in the
        master's spTree rather than using <p:bg>. A semi-transparent rectangle
        overlay may be placed on top to tint/dim the image.

        Args:
            slide: python-pptx slide object

        Returns:
            SlideBackground with image data and optional overlay, or None
        """
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        # Get slide dimensions for size comparison
        slide_w = self.presentation.slide_width if self.presentation else 12192000
        slide_h = self.presentation.slide_height if self.presentation else 6858000
        min_area_ratio = 0.9  # Shape must cover at least 90% of slide
        slide_area = slide_w * slide_h
        threshold = slide_area * min_area_ratio

        # Check layout shapes, then master shapes
        sources = []
        try:
            layout = slide.slide_layout
            if layout is not None:
                sources.append(layout)
                try:
                    master = layout.slide_master
                    if master is not None:
                        sources.append(master)
                except Exception:
                    pass
        except Exception:
            pass

        for source in sources:
            bg_image = None
            overlay_color = None
            overlay_opacity = None

            try:
                for shape in source.shapes:
                    # Check if shape covers most of the slide
                    shape_area = shape.width * shape.height
                    if shape_area < threshold:
                        continue

                    # Position should be near (0,0)
                    if shape.left > slide_w * 0.1 or shape.top > slide_h * 0.1:
                        continue

                    # Picture shape -> background image
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            image = shape.image
                            bg_image = {
                                'image_bytes': image.blob,
                                'image_type': image.ext,
                            }
                        except (ValueError, AttributeError):
                            continue

                    # Non-picture shape with blipFill (image fill) -> also background image
                    elif bg_image is None and hasattr(shape, '_element'):
                        blip = self._extract_blip_fill(shape)
                        if blip is not None:
                            bg_image = blip

                    # Auto-shape with solidFill -> potential overlay
                    # (semi-transparent rectangle on top of background image)
                    elif hasattr(shape, 'fill') and shape.fill:
                        try:
                            if hasattr(shape.fill, 'fore_color'):
                                fc = shape.fill.fore_color
                                color = None
                                if hasattr(fc, 'type') and fc.type is not None:
                                    if fc.type == 1 and hasattr(fc, 'rgb') and fc.rgb:
                                        color = f"#{fc.rgb}"
                                    elif fc.type == 2 and self.theme_color_extractor:
                                        color = self.theme_color_extractor.get_theme_color(int(fc.theme_color))

                                if color:
                                    # Check for alpha in XML
                                    alpha = self._extract_shape_fill_alpha(shape)
                                    if alpha is not None and alpha < 1.0:
                                        overlay_color = color
                                        overlay_opacity = alpha
                        except (AttributeError, TypeError):
                            pass
            except Exception:
                continue

            if bg_image:
                return SlideBackground(
                    bg_type='image',
                    image_data=bg_image,
                    overlay_color=overlay_color,
                    overlay_opacity=overlay_opacity,
                )

        return None
