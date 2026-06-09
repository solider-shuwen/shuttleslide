"""
Shape parsing mixin for PPTXParser.

Handles generic shape extraction, blip fills, geometry data, and shape fill alpha.
"""

import re
from typing import Optional, Dict, Any

from shuttleslide.pptx_to_html.models import ShapeElement
from shuttleslide.pptx_to_html.utils.namespaces import NS_A, NAMESPACES, NS_R_CLARK
from shuttleslide.pptx_to_html.utils.units import emu_to_px, angle_to_degrees


class ShapeMixin:
    """Generic shape parsing methods."""

    def _parse_generic_shape(
        self, shape, left: float, top: float, width: float, height: float, z_order: int
    ) -> Optional[ShapeElement]:
        """Parse a generic shape (rectangles, circles, lines, etc.)."""
        # Get shape type name
        # python-pptx EnumValue doesn't have .name; str() gives e.g. "AUTO_SHAPE (1)"
        shape_type_str = str(shape.shape_type)
        # Strip the enum value suffix like " (1)" for cleaner type names
        shape_type = re.sub(r'\s*\(\d+\)$', '', shape_type_str)

        # For AUTO_SHAPE, extract the actual preset geometry name from XML
        # (e.g., "cloudCallout", "roundedRectCallout") which is more specific
        preset_name = None
        if shape_type == "AUTO_SHAPE" and hasattr(shape, "_element"):
            try:
                ns = NS_A
                prst_geom = shape._element.find('.//a:prstGeom', ns)
                if prst_geom is not None:
                    preset_name = prst_geom.get('prst')
                    if preset_name:
                        shape_type = preset_name  # Use specific preset name instead of generic "AUTO_SHAPE"
            except Exception:
                pass

        # Similarly for FREEFORM shapes, try to detect if they have custom geometry
        if shape_type == "FREEFORM" and hasattr(shape, "_element"):
            # FREEFORM shapes already have custGeom which will be handled by geometry extraction
            pass

        # Extract fill and line colors
        fill_color = None
        line_color = None

        # Detect explicit noFill (e.g., FREEFORM shapes with stroke-only outlines)
        # When fill.type is BACKGROUND, it means "no fill, show background" (<a:noFill/>)
        if hasattr(shape, "fill") and shape.fill:
            try:
                if shape.fill.type is not None and str(shape.fill.type) == 'BACKGROUND (5)':
                    fill_color = "none"
            except (AttributeError, TypeError):
                pass

        # Extract fill color (handle different fill types, theme colors, and RGBColor objects)
        if fill_color is None and hasattr(shape, "fill") and shape.fill:
            try:
                # Check if fill has fore_color attribute
                if hasattr(shape.fill, "fore_color"):
                    fore_color = shape.fill.fore_color

                    # Try to get theme color first (if it's a theme color reference)
                    if hasattr(fore_color, "theme_color") and fore_color.theme_color is not None:
                        theme_color = fore_color.theme_color
                        # Try to get theme color value
                        if self.theme_color_extractor:
                            # Convert theme color enum to RGB
                            theme_rgb = self.theme_color_extractor.get_theme_color(theme_color)
                            if theme_rgb:
                                fill_color = theme_rgb

                    # If no theme color, try direct RGB
                    if not fill_color and hasattr(fore_color, "rgb") and fore_color.rgb:
                        rgb_obj = fore_color.rgb
                        # Handle RGBColor objects (pptx.dml.color.RGBColor)
                        if hasattr(rgb_obj, '__class__') and 'RGBColor' in str(rgb_obj.__class__):
                            # RGBColor object - convert to string representation
                            rgb_str = str(rgb_obj).strip()
                            if len(rgb_str) == 6 and rgb_str.isalnum():
                                fill_color = f"#{rgb_str}"
                        # Handle string RGB values
                        elif isinstance(rgb_obj, str):
                            rgb_str = rgb_obj.strip()
                            if len(rgb_str) >= 6:
                                fill_color = f"#{rgb_str[:6]}"
                        # Handle integer RGB values
                        elif isinstance(rgb_obj, int) and rgb_obj > 0:
                            fill_color = f"#{rgb_obj:06x}"
            except (AttributeError, TypeError, ValueError) as e:
                fill_color = None

        # Check for noFill in XML BEFORE accessing python-pptx line properties
        # Accessing shape.line.color triggers python-pptx to resolve style references,
        # which can modify the XML and replace <noFill/> with <solidFill/>
        _line_has_noFill = False
        if hasattr(shape, '_element'):
            try:
                _ns = {'a': NAMESPACES['a'], 'p': NAMESPACES['p']}
                _spPr = shape._element.find('./p:spPr', _ns)
                if _spPr is not None:
                    _ln = _spPr.find('./a:ln', _ns)
                    if _ln is not None:
                        _noFill = _ln.find('./a:noFill', _ns)
                        if _noFill is not None:
                            _line_has_noFill = True
            except Exception:
                pass

        # Extract line color (handle theme colors, different color types, and RGBColor objects)
        if not _line_has_noFill and hasattr(shape, "line") and shape.line:
            try:
                # Try color attribute first
                if hasattr(shape.line, "color") and shape.line.color:
                    line_color_obj = shape.line.color

                    # Try to get theme color first (if it's a theme color reference)
                    if hasattr(line_color_obj, "theme_color") and line_color_obj.theme_color is not None:
                        theme_color = line_color_obj.theme_color
                        # Try to get theme color value
                        if self.theme_color_extractor:
                            # Convert theme color enum to RGB
                            theme_rgb = self.theme_color_extractor.get_theme_color(theme_color)
                            if theme_rgb:
                                line_color = theme_rgb

                    # If no theme color, try direct RGB
                    if not line_color and hasattr(line_color_obj, "rgb") and line_color_obj.rgb:
                        rgb_obj = line_color_obj.rgb
                        # Handle RGBColor objects
                        if hasattr(rgb_obj, '__class__') and 'RGBColor' in str(rgb_obj.__class__):
                            rgb_str = str(rgb_obj).strip()
                            if len(rgb_str) == 6 and rgb_str.isalnum():
                                line_color = f"#{rgb_str}"
                        # Handle string RGB values
                        elif isinstance(rgb_obj, str):
                            rgb_str = rgb_obj.strip()
                            if len(rgb_str) >= 6:
                                line_color = f"#{rgb_str[:6]}"
                        # Handle integer RGB values
                        elif isinstance(rgb_obj, int) and rgb_obj > 0:
                            line_color = f"#{rgb_obj:06x}"
            except (AttributeError, TypeError, ValueError) as e:
                line_color = None

        # If no direct line color, try to extract from style/lnRef (theme style reference)
        if not _line_has_noFill and not line_color and hasattr(shape, '_element'):
            try:
                ns = {'a': NAMESPACES['a'], 'p': NAMESPACES['p']}

                # Find style element first, then lnRef child
                style = shape._element.find('.//p:style', ns)
                if style is not None:
                    # Find style element first, then lnRef child
                    style = shape._element.find('.//p:style', ns)
                    if style is not None:
                        ln_ref = style.find('.//a:lnRef', ns)
                        if ln_ref is not None:
                            # Check for schemeClr (theme color reference)
                            scheme_clr = ln_ref.find('.//a:schemeClr', ns)
                            if scheme_clr is not None:
                                val = scheme_clr.get('val')  # e.g., "accent1", "accent2", etc.
                                if val and self.theme_color_extractor:
                                    # Convert scheme color name to theme color enum
                                    # val is like "accent3", we need to convert to MSO_THEME_COLOR.ACCENT_3
                                    # The scheme color name uses lowercase without underscore, but the enum uses uppercase WITH underscore
                                    theme_color_name = val.upper()
                                    # Add underscore between letter and number (e.g., "ACCENT3" -> "ACCENT_3")
                                    theme_color_name = re.sub(r'([A-Z]+)(\d+)', r'\1_\2', theme_color_name)
                                    try:
                                        from pptx.enum.dml import MSO_THEME_COLOR
                                        theme_color_enum = getattr(MSO_THEME_COLOR, theme_color_name)
                                        theme_rgb = self.theme_color_extractor.get_theme_color(theme_color_enum)
                                        if theme_rgb:
                                            line_color = theme_rgb
                                    except (AttributeError, ValueError):
                                        pass
            except Exception:
                pass

        # Extract line dash style from XML (<a:ln><a:prstDash val="..."/>)
        dash_style = None
        if hasattr(shape, '_element'):
            try:
                ns = NS_A
                prst_dash = shape._element.find('.//a:ln/a:prstDash', ns)
                if prst_dash is not None:
                    dash_style = prst_dash.get('val')
            except Exception:
                pass

        # Check if shape has text
        text = None
        if hasattr(shape, "text_frame") and shape.text_frame:
            text = shape.text_frame.text

        # Extract blipFill (image fill) - some shapes have images as fill
        blip_fill = self._extract_blip_fill(shape)

        # Extract geometry data for SVG generation
        geometry_data = self._extract_shape_geometry(shape)

        # Build metadata
        metadata = {}
        if geometry_data:
            metadata["geometry"] = geometry_data
        if preset_name:
            metadata["preset_name"] = preset_name

        # Extract flipH/flipV, rotation, scene3d, and line width from XML
        flip_h = False
        flip_v = False
        rotation = None
        if hasattr(shape, '_element'):
            try:
                ns_xml = self._ns
                xfrm = shape._element.find('.//a:xfrm', ns_xml)
                if xfrm is not None:
                    flip_h = xfrm.get('flipH') == '1'
                    flip_v = xfrm.get('flipV') == '1'
                    rot_val = xfrm.get('rot')
                    if rot_val:
                        rotation = angle_to_degrees(float(rot_val))

                # Extract scene3d camera preset
                camera = shape._element.find('.//a:scene3d/a:camera', ns_xml)
                if camera is not None:
                    prst = camera.get('prst')
                    if prst:
                        metadata['scene3d_camera'] = prst

                # Extract explicit line width (<a:ln w="...">)
                ln_elem = shape._element.find('.//a:ln', ns_xml)
                if ln_elem is not None:
                    w_attr = ln_elem.get('w')
                    if w_attr:
                        metadata['line_width'] = emu_to_px(int(w_attr))
            except Exception:
                pass

        return ShapeElement(
            element_type="shape",
            left=left,
            top=top,
            width=width,
            height=height,
            z_order=z_order,
            shape_type=shape_type,
            fill_color=fill_color,
            line_color=line_color,
            dash_style=dash_style,
            text=text,
            blip_fill=blip_fill,
            flip_h=flip_h,
            flip_v=flip_v,
            rotation=rotation,
            metadata=metadata,
        )

    def _extract_blip_fill(self, shape) -> Optional[Dict[str, Any]]:
        """
        Extract blipFill (image fill) data from a shape.

        Some shapes use an image as their fill (e.g., callout shapes with
        handwritten-style content). The image is referenced via relationship ID.

        Args:
            shape: python-pptx shape object

        Returns:
            Dictionary with {image_bytes, image_type} or None
        """
        if not hasattr(shape, '_element'):
            return None

        try:
            ns = {'a': NAMESPACES['a'], 'r': NAMESPACES['r']}
            elem = shape._element

            # Look for blipFill with embedded blip
            blip_fill = elem.find('.//a:blipFill', ns)
            if blip_fill is None:
                return None

            blip = blip_fill.find('.//a:blip', ns)
            if blip is None:
                return None

            # Check for embedded image (r:embed) or linked image (r:link)
            embed_id = blip.get(f'{NS_R_CLARK}embed')
            if not embed_id:
                return None

            # Get the image from the relationship
            try:
                rel = shape.part.rels[embed_id]
                image_part = rel.target_part
                image_bytes = image_part.blob
                # Determine image type from content type
                content_type = image_part.content_type
                image_type = content_type.split('/')[-1] if '/' in content_type else 'png'
                # Normalize common types
                if image_type == 'jpeg':
                    image_type = 'jpg'
                elif image_type == 'x-emf':
                    image_type = 'emf'
                elif image_type == 'x-wmf':
                    image_type = 'wmf'

                return {
                    'image_bytes': image_bytes,
                    'image_type': image_type,
                }
            except (KeyError, AttributeError):
                return None

        except Exception:
            return None

    def _extract_shape_geometry(self, shape) -> Optional[Dict[str, Any]]:
        """
        Extract geometry data from PowerPoint shape for SVG generation.

        Args:
            shape: python-pptx shape object

        Returns:
            Dictionary with geometry data or None
        """
        if not hasattr(shape, '_element'):
            return None

        try:
            elem = shape._element
            namespaces = NS_A

            geometry_data = {}

            # Check for preset geometry
            prst_geom = elem.find('.//a:prstGeom', namespaces)
            if prst_geom is not None:
                geometry_data['type'] = 'preset'
                geometry_data['prst'] = prst_geom.get('prst')

                # Extract adjustment values from avLst
                av_lst = prst_geom.find('a:avLst', namespaces)
                if av_lst is not None:
                    adjustments = []
                    for gd in av_lst.findall('a:gd', namespaces):
                        name = gd.get('name')
                        fmla = gd.get('fmla', '')
                        if fmla.startswith('val '):
                            try:
                                value = int(fmla[4:])
                                adjustments.append({'name': name, 'value': value})
                            except ValueError:
                                pass
                    if adjustments:
                        geometry_data['adjustments'] = adjustments

            # Check for custom geometry (freeform shapes)
            custgeom = elem.find('.//a:custGeom', namespaces)
            if custgeom is not None:
                geometry_data['type'] = 'custom'

                # Extract path data
                path_lst = custgeom.find('a:pathLst', namespaces)
                if path_lst is not None:
                    paths = []
                    for path in path_lst.findall('a:path', namespaces):
                        path_data = {
                            'w': path.get('w'),
                            'h': path.get('h'),
                            'commands': []
                        }

                        # Parse path commands
                        for child in path:
                            tag_name = child.tag.split('}')[-1]

                            if tag_name == 'moveTo':
                                pt = child.find('a:pt', namespaces)
                                if pt is not None:
                                    path_data['commands'].append({
                                        'type': 'M',
                                        'x': pt.get('x'),
                                        'y': pt.get('y')
                                    })

                            elif tag_name == 'lnTo':
                                pt = child.find('a:pt', namespaces)
                                if pt is not None:
                                    path_data['commands'].append({
                                        'type': 'L',
                                        'x': pt.get('x'),
                                        'y': pt.get('y')
                                    })

                            elif tag_name == 'cubicBezTo':
                                points = child.findall('a:pt', namespaces)
                                if len(points) >= 3:
                                    path_data['commands'].append({
                                        'type': 'C',
                                        'x1': points[0].get('x'),
                                        'y1': points[0].get('y'),
                                        'x2': points[1].get('x'),
                                        'y2': points[1].get('y'),
                                        'x': points[2].get('x'),
                                        'y': points[2].get('y')
                                    })

                            elif tag_name == 'quadBezTo':
                                points = child.findall('a:pt', namespaces)
                                if len(points) >= 2:
                                    path_data['commands'].append({
                                        'type': 'Q',
                                        'x1': points[0].get('x'),
                                        'y1': points[0].get('y'),
                                        'x': points[1].get('x'),
                                        'y': points[1].get('y')
                                    })

                            elif tag_name == 'arcTo':
                                path_data['commands'].append({
                                    'type': 'A',
                                    'wR': child.get('wR'),
                                    'hR': child.get('hR'),
                                    'stAng': child.get('stAng'),
                                    'swAng': child.get('swAng'),
                                    'x': child.get('x'),
                                    'y': child.get('y')
                                })

                            elif tag_name == 'close':
                                path_data['commands'].append({'type': 'Z'})

                        if path_data['commands']:
                            paths.append(path_data)

                    geometry_data['paths'] = paths

            return geometry_data if geometry_data else None

        except Exception:
            # If geometry extraction fails, return None
            return None

    def _extract_shape_fill_alpha(self, shape) -> Optional[float]:
        """
        Extract alpha value from a shape's solidFill.

        Args:
            shape: python-pptx shape object

        Returns:
            Alpha as 0.0-1.0, or None if no alpha specified
        """
        if not hasattr(shape, '_element'):
            return None

        ns = self._ns
        try:
            # Look for <a:solidFill><a:srgbClr><a:alpha val="..."/>
            # or <a:solidFill><a:schemeClr><a:alpha val="..."/>
            alpha_elem = shape._element.find('.//a:solidFill/*/a:alpha', ns)
            if alpha_elem is not None:
                val = alpha_elem.get('val')
                if val:
                    # OpenXML alpha: 0-100000 where 100000 = 100%
                    return int(val) / 100000.0
        except Exception:
            pass
        return None
