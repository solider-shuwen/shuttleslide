"""
PPTX Parser - core orchestration module.

PPTXParser coordinates the parsing of PowerPoint files by delegating
to specialized mixins for each element type.
"""

from pathlib import Path
import re
from typing import List, Dict, Any, Optional, Union
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from shuttleslide.pptx_to_html.models import (
    SlideElement, ShapeElement, ParsedSlide, MasterTextStyle,
)
from shuttleslide.pptx_to_html.utils.namespaces import NAMESPACES, NS_P_CLARK
from shuttleslide.pptx_to_html.utils.units import emu_to_px
from shuttleslide.pptx_to_html.theme_colors import ThemeColorExtractor

from shuttleslide.pptx_to_html.parser.color_mixin import ColorMixin
from shuttleslide.pptx_to_html.parser.master_styles_mixin import MasterStylesMixin
from shuttleslide.pptx_to_html.parser.bullets_mixin import BulletsMixin
from shuttleslide.pptx_to_html.parser.text_mixin import TextMixin
from shuttleslide.pptx_to_html.parser.table_mixin import TableMixin
from shuttleslide.pptx_to_html.parser.image_mixin import ImageMixin
from shuttleslide.pptx_to_html.parser.shape_mixin import ShapeMixin
from shuttleslide.pptx_to_html.parser.group_mixin import GroupMixin
from shuttleslide.pptx_to_html.parser.background_mixin import BackgroundMixin


class PPTXParser(
    ColorMixin,
    MasterStylesMixin,
    BulletsMixin,
    TextMixin,
    TableMixin,
    ImageMixin,
    ShapeMixin,
    GroupMixin,
    BackgroundMixin,
):
    """
    Parser for PPTX files that extracts slide structure and elements.

    Delegates to specialized mixins:
    - ColorMixin: color extraction and resolution
    - MasterStylesMixin: slide master spacing/style extraction
    - BulletsMixin: bullet property parsing and inheritance
    - TextMixin: text box and placeholder parsing
    - TableMixin: table element parsing
    - ImageMixin: image extraction, cropping, color replacement
    - ShapeMixin: generic shape parsing, geometry, blip fill
    - GroupMixin: group shape parsing with coordinate transforms
    - BackgroundMixin: background extraction with inheritance
    """

    def __init__(self, pptx_path: Union[str, Path]):
        """
        Initialize the parser with a PPTX file path.

        Args:
            pptx_path: Path to the PowerPoint file
        """
        self.pptx_path = Path(pptx_path)
        self.presentation: Optional[Presentation] = None
        self.slides: List[ParsedSlide] = []

        # Initialize theme color extractor
        try:
            self.theme_color_extractor = ThemeColorExtractor(str(self.pptx_path))
        except Exception as e:
            print(f"Warning: Could not initialize theme color extractor: {e}")
            self.theme_color_extractor = None

        # XML namespaces used for background and style extraction
        self._ns = NAMESPACES

        # Extract default paragraph spacing from slide master (per-master cache)
        self.default_line_spacing: Optional[float] = None
        self._master_spacing_cache: Dict[int, Optional[float]] = {}  # keyed by id(master)

        # Extract master-level text styles (per-master cache)
        self.master_title_styles: Dict[int, MasterTextStyle] = {}
        self.master_body_styles: Dict[int, MasterTextStyle] = {}
        self._master_text_styles_cache: Dict[int, tuple] = {}  # keyed by id(master)

        # Current slide's master reference (set per-slide during parsing)
        self._current_master = None
        self._current_layout = None

        self._extract_master_default_spacing()
        self._extract_master_text_styles()

    def parse(self) -> List[ParsedSlide]:
        """
        Parse the PPTX file and return list of parsed slides.

        Returns:
            List of parsed slides with their elements
        """
        self.presentation = Presentation(str(self.pptx_path))
        self.slides = []

        for slide_idx, slide in enumerate(self.presentation.slides, start=1):
            parsed_slide = self._parse_slide(slide, slide_idx)
            # Check if slide is hidden
            parsed_slide.hidden = self._is_slide_hidden(slide)
            # Detect animations
            parsed_slide.has_animations = self._detect_slide_animations(slide)
            self.slides.append(parsed_slide)

        return self.slides

    def _parse_slide(self, slide, slide_number: int) -> ParsedSlide:
        """
        Parse a single slide and extract all elements.

        Args:
            slide: python-pptx slide object
            slide_number: Slide number (1-indexed)

        Returns:
            ParsedSlide object with all elements
        """
        # Get slide dimensions from presentation (convert EMU to pixels)
        # PPTX uses EMU (English Metric Units): 1 inch = 914400 EMU
        # Standard screen: 1 inch = 96 pixels
        # Conversion: pixels = EMU / 914400 * 96 = EMU / 9525
        slide_width = emu_to_px(self.presentation.slide_width)
        slide_height = emu_to_px(self.presentation.slide_height)

        # Get layout name
        layout_name = slide.slide_layout.name if slide.slide_layout else "Blank"

        parsed_slide = ParsedSlide(
            slide_number=slide_number,
            layout_name=layout_name,
            width=slide_width,
            height=slide_height,
        )

        # Extract background following inheritance chain (slide -> layout -> master)
        parsed_slide.background = self._extract_slide_background(slide)

        # Set current master for per-master style lookups during shape parsing
        self._current_master = None
        self._current_layout = None
        try:
            layout = slide.slide_layout
            if layout is not None:
                self._current_master = layout.slide_master
                self._current_layout = layout
        except Exception:
            pass

        # Extract all shapes from the slide
        for z_order, shape in enumerate(slide.shapes):
            result = self._parse_shape(shape, z_order)
            if result is not None:
                if isinstance(result, list):
                    parsed_slide.elements.extend(result)
                else:
                    parsed_slide.elements.append(result)

        # Extract non-background shapes from layout and master
        # (logos, decorative elements, etc. that appear on every slide)
        self._add_inherited_shapes(slide, parsed_slide)

        return parsed_slide

    def _add_inherited_shapes(self, slide, parsed_slide: ParsedSlide):
        """
        Add non-background shapes from the slide's layout and master.

        PowerPoint renders layout/master shapes (logos, decorative elements)
        behind slide content. We skip shapes that are part of the background
        (full-slide images/overlays) since those are handled by background
        extraction.

        Args:
            slide: python-pptx slide object
            parsed_slide: ParsedSlide to add inherited shapes to
        """
        slide_w = self.presentation.slide_width if self.presentation else 12192000
        slide_h = self.presentation.slide_height if self.presentation else 6858000
        slide_area = slide_w * slide_h

        # Check if master shapes are hidden
        show_master_sp = slide._element.get('showMasterSp')
        if show_master_sp == '0':
            return

        sources = []
        try:
            layout = slide.slide_layout
            if layout is not None:
                sources.append(layout)
                try:
                    master = layout.slide_master
                    if master is not None:
                        sources.append(master)
                except Exception:
                    pass
        except Exception:
            return

        # Get the starting z-order for inherited shapes (behind slide shapes)
        base_z = 0  # Inherited shapes go behind (lower z-order)

        for source in sources:
            try:
                for shape in source.shapes:
                    shape_area = shape.width * shape.height
                    area_ratio = shape_area / slide_area if slide_area > 0 else 0

                    # Skip full-slide background shapes (handled by background extraction)
                    if area_ratio > 0.9 and shape.left < slide_w * 0.1 and shape.top < slide_h * 0.1:
                        continue

                    # Skip placeholder shapes (title, body, etc.)
                    if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                        continue

                    # Parse the shape
                    result = self._parse_shape(shape, base_z)
                    if result is not None:
                        if isinstance(result, list):
                            parsed_slide.elements.extend(result)
                            base_z -= len(result)
                        else:
                            parsed_slide.elements.append(result)
                            base_z -= 1  # Each inherited shape gets lower z-order
            except Exception:
                continue

    def _shape_has_text(self, shape) -> bool:
        """Check if a shape has actual text content in its text_frame."""
        if not (hasattr(shape, "text_frame") and shape.text_frame):
            return False
        if hasattr(shape.text_frame, "text") and shape.text_frame.text.strip():
            return True
        for para in shape.text_frame.paragraphs:
            if hasattr(para, "text") and para.text.strip():
                return True
        return False

    def _parse_shape(self, shape, z_order: int) -> Optional[Union[SlideElement, List[SlideElement]]]:
        """
        Parse a single shape and return the appropriate element.

        For geometric shapes (FREEFORM, AUTO_SHAPE, etc.) that also contain text,
        returns a list of [ShapeElement, TextElement] so both the visual geometry
        and the rich text content are rendered.

        Args:
            shape: python-pptx shape object
            z_order: Z-order of the shape

        Returns:
            SlideElement, list of SlideElements, or None if shape type is not supported
        """
        # Convert EMU to pixels for element positions
        left = emu_to_px(shape.left)
        top = emu_to_px(shape.top)
        width = emu_to_px(shape.width)
        height = emu_to_px(shape.height)

        # Table (check BEFORE placeholder - tables inside placeholders must be handled as tables)
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return self._parse_table(shape, left, top, width, height, z_order)

        # Placeholder shape (non-table placeholders)
        if shape.is_placeholder:
            return self._parse_placeholder(shape, left, top, width, height, z_order)

        # Picture/Image
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_element = self._parse_image(shape, left, top, width, height, z_order)
            if image_element is None:
                # If image parsing fails, create a placeholder shape
                return ShapeElement(
                    element_type="shape",
                    left=left,
                    top=top,
                    width=width,
                    height=height,
                    z_order=z_order,
                    shape_type="PICTURE_PLACEHOLDER",
                    fill_color="#CCCCCC",
                    line_color="#666666",
                    text="[Image]",
                )
            return image_element

        # Group shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            group_element = self._parse_group(shape, z_order)
            if group_element is None:
                # If group parsing fails, create a placeholder shape
                return ShapeElement(
                    element_type="shape",
                    left=left,
                    top=top,
                    width=width,
                    height=height,
                    z_order=z_order,
                    shape_type="GROUP_PLACEHOLDER",
                    fill_color="#DDDDDD",
                    line_color="#999999",
                    text=f"[Group with {len(shape.shapes)} items]" if hasattr(shape, 'shapes') else "[Group]",
                )
            return group_element

        # Pure text box shapes - no visual geometry to render
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            if self._shape_has_text(shape):
                return self._parse_text_box(shape, left, top, width, height, z_order)
            return None

        # All other shapes (FREEFORM, AUTO_SHAPE, LINE, CALLOUT, etc.)
        # These have visual geometry that must be rendered
        shape_element = self._parse_generic_shape(shape, left, top, width, height, z_order)

        # If the shape also has text, render both geometry AND rich text
        if self._shape_has_text(shape) and shape_element:
            text_element = self._parse_text_box(shape, left, top, width, height, z_order + 1)
            # Clear text on shape element to avoid duplicate text rendering
            shape_element.text = None

            # Propagate scene3d camera metadata so text gets the same 3D transform
            if shape_element.metadata and shape_element.metadata.get('scene3d_camera'):
                if text_element.metadata is None:
                    text_element.metadata = {}
                text_element.metadata['scene3d_camera'] = shape_element.metadata['scene3d_camera']

            # When a shape has visual geometry (SVG), the outline/border is rendered
            # by the shape SVG, not as CSS border on the text element.
            # Clear line_color on the text overlay to avoid a rectangular border
            # that incorrectly overlaps the shape geometry.
            text_element.line_color = None
            text_element.line_width = None

            # Only include shape element if it has visible geometry
            # (fill, stroke, image fill, gradient fill, or custom geometry like FREEFORM paths)
            has_visible_geometry = (
                shape_element.fill_color is not None
                or (hasattr(shape_element, 'fill_gradient') and shape_element.fill_gradient)
                or shape_element.line_color is not None
                or (hasattr(shape_element, 'blip_fill') and shape_element.blip_fill)
                or (shape_element.metadata and shape_element.metadata.get('geometry'))
            )
            if has_visible_geometry:
                return [shape_element, text_element]
            else:
                # No visible geometry — just return the text element
                return text_element

        return shape_element

    def _is_slide_hidden(self, slide) -> bool:
        """
        Detect if a slide is marked as hidden in PowerPoint.

        Args:
            slide: python-pptx slide object

        Returns:
            True if slide is hidden, False otherwise
        """
        try:
            # Try to access the slide's XML element directly
            elem = slide._element

            # Check for common hidden attributes
            # PowerPoint uses different attributes to mark slides as hidden
            for attr in ("show", "hidden"):
                val = elem.get(attr)
                if val is not None:
                    # "show" attribute: 0 means hidden, 1 means visible
                    # "hidden" attribute: 1 means hidden, 0/none means visible
                    if attr == "show":
                        return str(val) == "0"
                    else:  # "hidden" attribute
                        return str(val) == "1"

            # Try namespace-specific attributes
            # PowerPoint uses the PresentationML namespace
            ns = NS_P_CLARK
            for attr in ("show", "hidden"):
                val = elem.get(f"{ns}{attr}")
                if val is not None:
                    if attr == "show":
                        return str(val) == "0"
                    else:  # "hidden" attribute
                        return str(val) == "1"

        except Exception:
            # If any error occurs, assume slide is not hidden
            pass

        return False

    def _detect_slide_animations(self, slide) -> bool:
        """
        Detect if a slide has animations.

        Args:
            slide: python-pptx slide object

        Returns:
            True if slide has animations, False otherwise
        """
        try:
            # Try to access the slide's XML element directly
            elem = slide._element

            # Define the PresentationML namespace
            ns = {"p": NAMESPACES['p']}

            # Look for timing elements which indicate animations
            # In PowerPoint XML, animations are stored in timing elements
            timing = elem.find('.//p:timing', ns)
            if timing is not None:
                return True

            # Also check for animation elements directly
            anim = elem.find('.//p:anim', ns)
            if anim is not None:
                return True

        except Exception:
            # If any error occurs, assume no animations
            pass

        return False

    def get_presentation_metadata(self) -> Dict[str, Any]:
        """
        Get metadata about the presentation.

        Returns:
            Dictionary with presentation metadata
        """
        if not self.presentation:
            return {}

        return {
            "title": self.presentation.core_properties.title or "",
            "author": self.presentation.core_properties.author or "",
            "subject": self.presentation.core_properties.subject or "",
            "created": str(self.presentation.core_properties.created) if self.presentation.core_properties.created else "",
            "modified": str(self.presentation.core_properties.modified) if hasattr(self.presentation.core_properties, "modified") and self.presentation.core_properties.modified else "",
            "slide_count": len(self.presentation.slides),
            "slide_width": self.presentation.slide_width,
            "slide_height": self.presentation.slide_height,
        }
