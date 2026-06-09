"""
PPTX Parser - parses PowerPoint files and extracts slide information.
"""

from pathlib import Path
import re
from typing import List, Dict, Any, Optional, Union
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from shuttleslide.pptx_to_html.models import (
    SlideElement, RunElement, BulletProperties, ParagraphElement,
    TextElement, TableElement, ImageElement, ShapeElement, GroupElement,
    SlideBackground, MasterTextStyle, ParsedSlide,
)
from shuttleslide.pptx_to_html.utils.text_sanitizer import sanitize_pptx_text
from shuttleslide.pptx_to_html.utils.namespaces import NAMESPACES, NS_A, NS_R_CLARK, NS_P_CLARK
from shuttleslide.pptx_to_html.utils.units import emu_to_px, emu_to_pt, emu_to_inches, px_to_emu, angle_to_degrees, EMU_PER_INCH
from shuttleslide.pptx_to_html.utils.colors import adjust_color_luminance, resolve_xml_color
from shuttleslide.pptx_to_html.theme_colors import ThemeColorExtractor


class PPTXParser:
    """
    Parser for PPTX files that extracts slide structure and elements.
    """

    def __init__(self, pptx_path: Union[str, Path]):
        """
        Initialize the parser with a PPTX file path.

        Args:
            pptx_path: Path to the PowerPoint file
        """
        self.pptx_path = Path(pptx_path)
        self.presentation: Optional[Presentation] = None
        self.slides: List[ParsedSlide] = []

        # Initialize theme color extractor
        try:
            self.theme_color_extractor = ThemeColorExtractor(str(self.pptx_path))
        except Exception as e:
            print(f"Warning: Could not initialize theme color extractor: {e}")
            self.theme_color_extractor = None

        # XML namespaces used for background and style extraction
        self._ns = NAMESPACES

        # Extract default paragraph spacing from slide master (per-master cache)
        self.default_line_spacing: Optional[float] = None
        self._master_spacing_cache: Dict[int, Optional[float]] = {}  # keyed by id(master)

        # Extract master-level text styles (per-master cache)
        self.master_title_styles: Dict[int, MasterTextStyle] = {}
        self.master_body_styles: Dict[int, MasterTextStyle] = {}
        self._master_text_styles_cache: Dict[int, tuple] = {}  # keyed by id(master)

        # Current slide's master reference (set per-slide during parsing)
        self._current_master = None

        self._extract_master_default_spacing()
        self._extract_master_text_styles()

    def _extract_master_default_spacing(self):
        """Extract default line spacing from the first slide master's body style."""
        try:
            prs = Presentation(str(self.pptx_path))
            if not prs.slide_masters:
                return
            master = prs.slide_masters[0]
            spacing = self._get_master_spacing_for(master)

            # Set as global default
            self.default_line_spacing = spacing
        except Exception:
            pass

    def _get_master_spacing_for(self, master) -> Optional[float]:
        """
        Get default line spacing for a specific master, using cache.

        Args:
            master: python-pptx slide master object

        Returns:
            Default line spacing as float multiplier, or None
        """
        mid = id(master)
        if mid in self._master_spacing_cache:
            return self._master_spacing_cache[mid]

        spacing = None
        try:
            ns = self._ns
            elem = master._element

            body_style = elem.find('.//p:bodyStyle', ns)
            if body_style is not None:
                lvl1 = body_style.find('a:lvl1pPr', ns)
                if lvl1 is not None:
                    ln_spc = lvl1.find('a:lnSpc', ns)
                    if ln_spc is not None:
                        spc_pct = ln_spc.find('a:spcPct', ns)
                        spc_pts = ln_spc.find('a:spcPts', ns)
                        if spc_pct is not None:
                            val = int(spc_pct.get('val', 100000))
                            spacing = val / 100000.0
                        elif spc_pts is not None:
                            val = int(spc_pts.get('val', 0))
                            spacing = val / 100.0
        except Exception:
            pass

        self._master_spacing_cache[mid] = spacing
        return spacing

    def _get_current_master_styles(self) -> tuple:
        """
        Get text styles for the current slide's master.

        Returns:
            Tuple of (title_styles, body_styles) dicts
        """
        if self._current_master is not None:
            return self._get_master_text_styles_for(self._current_master)
        return self.master_title_styles, self.master_body_styles

    def _extract_run_color(self, font) -> Optional[str]:
        """
        Extract color from a font object, handling RGB and theme colors.

        Args:
            font: python-pptx Font object from a run

        Returns:
            Hex color string (e.g., '#FF0000') or None
        """
        if not font.color or font.color.type is None:
            return None

        if font.color.type == 1:  # MSO_COLOR_TYPE.RGB
            return f"#{font.color.rgb}"

        if font.color.type == 2:  # MSO_COLOR_TYPE.SCHEME (theme color)
            if self.theme_color_extractor:
                return self.theme_color_extractor.get_theme_color(
                    int(font.color.theme_color)
                )

        return None

    def _resolve_xml_color(self, color_elem) -> Optional[str]:
        """
        Resolve a DrawingML color element to a hex string.

        Delegates to utils.colors.resolve_xml_color with the theme extractor.
        """
        return resolve_xml_color(color_elem, self.theme_color_extractor)

    def _adjust_color_luminance(self, hex_color: str, factor: float) -> str:
        """Adjust color luminance by multiplying RGB values. Delegates to utils."""
        return adjust_color_luminance(hex_color, factor)

    def _parse_bullet_properties(self, ppr, para_level: int, shape=None) -> BulletProperties:
        """
        Parse bullet properties from a paragraph's <a:pPr> XML element.

        Args:
            ppr: <a:pPr> XML element (can be None)
            para_level: Paragraph indent level (0-8)

        Returns:
            BulletProperties with parsed bullet info
        """
        ns = self._ns

        if ppr is None:
            return BulletProperties(type='inherited')

        # Check for explicit no-bullet
        bu_none = ppr.find('a:buNone', ns)
        if bu_none is not None:
            return BulletProperties(type='none')

        # Check for character bullet
        bu_char = ppr.find('a:buChar', ns)
        if bu_char is not None:
            bullet = BulletProperties(type='char')
            bullet.char = bu_char.get('char', '\u2022')

            # Bullet styling from pPr siblings
            bu_font = ppr.find('a:buFont', ns)
            if bu_font is not None:
                bullet.font_typeface = bu_font.get('typeface')

            bu_sz = ppr.find('a:buSzPct', ns)
            if bu_sz is not None:
                try:
                    bullet.font_size_pct = int(bu_sz.get('val', 100000))
                except (ValueError, TypeError):
                    pass

            bu_clr = ppr.find('a:buClr', ns)
            if bu_clr is not None:
                bullet.color = self._resolve_xml_color(bu_clr)

            return bullet

        # Check for auto-numbered bullet
        bu_auto = ppr.find('a:buAutoNum', ns)
        if bu_auto is not None:
            bullet = BulletProperties(type='autonum')
            bullet.autonum_type = bu_auto.get('type', 'arabicPeriod')
            start_at = bu_auto.get('startAt')
            if start_at:
                try:
                    bullet.autonum_start = int(start_at)
                except (ValueError, TypeError):
                    bullet.autonum_start = 1

            bu_font = ppr.find('a:buFont', ns)
            if bu_font is not None:
                bullet.font_typeface = bu_font.get('typeface')

            bu_sz = ppr.find('a:buSzPct', ns)
            if bu_sz is not None:
                try:
                    bullet.font_size_pct = int(bu_sz.get('val', 100000))
                except (ValueError, TypeError):
                    pass

            bu_clr = ppr.find('a:buClr', ns)
            if bu_clr is not None:
                bullet.color = self._resolve_xml_color(bu_clr)

            return bullet

        # Check for image bullet (buBlip)
        bu_blip = ppr.find('a:buBlip', ns)
        if bu_blip is not None:
            bullet = BulletProperties(type='blip')

            # Resolve the embedded image via relationship
            blip = bu_blip.find('a:blip', ns)
            if blip is not None and shape is not None:
                embed_id = blip.get(f'{NS_R_CLARK}embed')
                if embed_id:
                    try:
                        rel = shape.part.rels[embed_id]
                        image_part = rel.target_part
                        bullet.blip_image_bytes = image_part.blob
                        content_type = image_part.content_type
                        image_type = content_type.split('/')[-1] if '/' in content_type else 'png'
                        if image_type == 'jpeg':
                            image_type = 'jpg'
                        elif image_type == 'x-emf':
                            image_type = 'emf'
                        elif image_type == 'x-wmf':
                            image_type = 'wmf'
                        bullet.blip_image_type = image_type
                    except (KeyError, AttributeError):
                        pass

            # Bullet styling (same pattern as buChar/buAutoNum)
            bu_font = ppr.find('a:buFont', ns)
            if bu_font is not None:
                bullet.font_typeface = bu_font.get('typeface')

            bu_sz = ppr.find('a:buSzPct', ns)
            if bu_sz is not None:
                try:
                    bullet.font_size_pct = int(bu_sz.get('val', 100000))
                except (ValueError, TypeError):
                    pass

            bu_clr = ppr.find('a:buClr', ns)
            if bu_clr is not None:
                bullet.color = self._resolve_xml_color(bu_clr)

            return bullet

        # No bullet element -- inherited from master/style
        return BulletProperties(type='inherited')

    def _resolve_bullet_inheritance(self, bullet: BulletProperties,
                                     para_level: int,
                                     is_title: bool = False,
                                     is_placeholder: bool = False) -> BulletProperties:
        """
        Resolve inherited bullet properties from master styles.

        Only placeholder text boxes inherit bullets from master.
        Non-placeholder text boxes default to no bullet.

        Args:
            bullet: Parsed BulletProperties (type='inherited')
            para_level: Paragraph indent level (0-8)
            is_title: Whether this is a title placeholder
            is_placeholder: Whether this shape is a placeholder

        Returns:
            Resolved BulletProperties
        """
        if bullet.type != 'inherited':
            return bullet

        # Non-placeholder text boxes don't inherit bullets from master
        if not is_placeholder:
            return BulletProperties(type='none')

        # Title placeholders typically don't have bullets
        if is_title:
            return BulletProperties(type='none')

        # Look up master body styles
        _, body_styles = self._get_current_master_styles()
        # Master levels are 1-9, paragraph levels are 0-8
        master_key = para_level + 1
        master_style = body_styles.get(master_key) if body_styles else None

        if master_style is None or master_style.bullet_type is None:
            return BulletProperties(type='none')

        resolved = BulletProperties(type=master_style.bullet_type)
        if master_style.bullet_type == 'char':
            resolved.char = master_style.bullet_char or '\u2022'
        elif master_style.bullet_type == 'autonum':
            resolved.autonum_type = master_style.bullet_autonum_type or 'arabicPeriod'
        resolved.font_typeface = master_style.bullet_font

        return resolved

    def _parse_bg_element(self, bg_elem) -> Optional[SlideBackground]:
        """
        Parse a <p:bg> XML element and return a SlideBackground.

        Handles bgPr (background properties) and bgRef (background reference).

        Args:
            bg_elem: <p:bg> XML element

        Returns:
            SlideBackground or None if background is empty/transparent
        """
        ns = self._ns

        # Check for <p:bgPr> (direct background properties)
        bg_pr = bg_elem.find('p:bgPr', ns)
        if bg_pr is not None:
            return self._parse_fill_from_element(bg_pr)

        # Check for <p:bgRef> (reference to theme background)
        bg_ref = bg_elem.find('p:bgRef', ns)
        if bg_ref is not None:
            idx = bg_ref.get('idx')
            if idx == '0':
                # idx=0 means empty/no fill
                return None

            # bgRef contains a color directly (common pattern)
            # e.g., <p:bgRef idx="1001"><a:schemeClr val="bg1"/></p:bgRef>
            color = self._resolve_xml_color(bg_ref)
            if color:
                return SlideBackground(bg_type='solid', color=color)

        return None

    def _parse_fill_from_element(self, elem) -> Optional[SlideBackground]:
        """
        Parse fill properties from an XML element (bgPr, etc.).

        Handles solidFill, gradFill, and blipFill.

        Args:
            elem: XML element containing fill properties

        Returns:
            SlideBackground or None
        """
        ns = self._ns

        # Solid fill
        solid_fill = elem.find('a:solidFill', ns)
        if solid_fill is not None:
            color = self._resolve_xml_color(solid_fill)
            if color:
                return SlideBackground(bg_type='solid', color=color)

        # Gradient fill
        grad_fill = elem.find('a:gradFill', ns)
        if grad_fill is not None:
            return self._parse_gradient_fill(grad_fill)

        # Image/blip fill
        blip_fill = elem.find('a:blipFill', ns)
        if blip_fill is not None:
            return self._parse_blip_fill_background(blip_fill)

        return None

    def _parse_gradient_fill(self, grad_fill) -> Optional[SlideBackground]:
        """
        Parse <a:gradFill> and convert to CSS linear-gradient.

        Args:
            grad_fill: <a:gradFill> XML element

        Returns:
            SlideBackground with gradient CSS
        """
        ns = self._ns

        # Parse gradient stops
        gs_lst = grad_fill.find('a:gsLst', ns)
        if gs_lst is None:
            return None

        stops = []
        for gs in gs_lst.findall('a:gs', ns):
            pos = gs.get('pos', '0')
            pos_pct = int(pos) / 1000  # XML uses 0-100000, CSS uses 0-100
            color = self._resolve_xml_color(gs)
            if color:
                stops.append(f"{color} {pos_pct}%")

        if not stops:
            return None

        # Parse angle from <a:lin> element
        # OpenXML angle: 60000ths of a degree, counter-clockwise from positive X-axis
        # CSS angle: degrees, clockwise from positive Y-axis (north)
        css_angle = 180  # default: bottom-to-top
        lin = grad_fill.find('a:lin', ns)
        if lin is not None:
            ang = lin.get('ang')
            if ang:
                xml_angle = angle_to_degrees(int(ang))  # Convert to degrees
                # OpenXML: 0 = left-to-right, CSS: 0 = bottom-to-top
                # Conversion: css_angle = (xml_angle - 90) % 360 ... actually:
                # OpenXML measures from X-axis (East) counter-clockwise
                # CSS measures from Y-axis (North) clockwise
                # css_angle = 90 - xml_angle, normalized to 0-360
                css_angle = (90 - xml_angle) % 360

        gradient_css = f"linear-gradient({css_angle:.1f}deg, {', '.join(stops)})"
        return SlideBackground(bg_type='gradient', gradient_css=gradient_css)

    def _parse_blip_fill_background(self, blip_fill) -> Optional[SlideBackground]:
        """
        Parse <a:blipFill> for background image.

        Args:
            blip_fill: <a:blipFill> XML element

        Returns:
            SlideBackground with image data or None
        """
        # This is handled at the slide level where we have access to relationships
        # For now, store a marker; the actual extraction happens in _extract_slide_background
        return None

    def _extract_slide_background(self, slide) -> Optional[SlideBackground]:
        """
        Extract background for a slide following the OpenXML inheritance chain:
        slide -> layout -> master -> default white.

        Args:
            slide: python-pptx slide object

        Returns:
            SlideBackground or None (defaults to white in layout)
        """
        ns = self._ns

        # Check if slide hides master background/shapes
        slide_elem = slide._element
        show_master_sp = slide_elem.get('showMasterSp')

        # Priority 1: Full-slide picture shapes on layout/master
        # Many presentations use a <p:pic> element covering the entire slide
        # as the visual background. These are drawn ON TOP of <p:bg>, so they
        # are what the user actually sees.
        if show_master_sp != '0':
            image_bg = self._find_background_image_shape(slide)
            if image_bg:
                return image_bg

        # Priority 2: <p:bg> element (background fill definition)
        # Check slide's own background
        bg = slide_elem.find('.//p:bg', ns)
        if bg is not None:
            # Check for blipFill (image background) - needs relationship access
            bg_pr = bg.find('p:bgPr', ns)
            if bg_pr is not None:
                blip_fill = bg_pr.find('a:blipFill', ns)
                if blip_fill is not None:
                    image_bg = self._extract_bg_blip_fill(slide, blip_fill)
                    if image_bg:
                        return image_bg

            result = self._parse_bg_element(bg)
            if result:
                return result

        # Check layout/master backgrounds
        if show_master_sp != '0':
            try:
                layout = slide.slide_layout
                if layout is not None:
                    layout_elem = layout._element
                    layout_bg = layout_elem.find('.//p:bg', ns)
                    if layout_bg is not None:
                        result = self._parse_bg_element(layout_bg)
                        if result:
                            return result

                    try:
                        master = layout.slide_master
                        master_elem = master._element
                        master_bg = master_elem.find('.//p:bg', ns)
                        if master_bg is not None:
                            result = self._parse_bg_element(master_bg)
                            if result:
                                return result
                    except Exception:
                        pass
            except Exception:
                pass

        return None

    def _extract_bg_blip_fill(self, slide, blip_fill) -> Optional[SlideBackground]:
        """
        Extract image data from a background blipFill element.

        Args:
            slide: python-pptx slide object
            blip_fill: <a:blipFill> XML element

        Returns:
            SlideBackground with image data or None
        """
        ns = self._ns
        try:
            blip = blip_fill.find('.//a:blip', ns)
            if blip is None:
                return None

            embed_id = blip.get(f'{NS_R_CLARK}embed')
            if not embed_id:
                return None

            rel = slide.part.rels[embed_id]
            image_part = rel.target_part
            image_bytes = image_part.blob
            content_type = image_part.content_type
            image_type = content_type.split('/')[-1] if '/' in content_type else 'png'
            if image_type == 'jpeg':
                image_type = 'jpg'

            return SlideBackground(
                bg_type='image',
                image_data={
                    'image_bytes': image_bytes,
                    'image_type': image_type,
                }
            )
        except Exception:
            return None

    def _find_background_image_shape(self, slide) -> Optional[SlideBackground]:
        """
        Find full-slide picture shapes and overlay shapes on the layout/master
        that serve as the visual background.

        Many presentations place background images as <p:pic> elements in the
        master's spTree rather than using <p:bg>. A semi-transparent rectangle
        overlay may be placed on top to tint/dim the image.

        Args:
            slide: python-pptx slide object

        Returns:
            SlideBackground with image data and optional overlay, or None
        """
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        # Get slide dimensions for size comparison
        slide_w = self.presentation.slide_width if self.presentation else 12192000
        slide_h = self.presentation.slide_height if self.presentation else 6858000
        min_area_ratio = 0.9  # Shape must cover at least 90% of slide
        slide_area = slide_w * slide_h
        threshold = slide_area * min_area_ratio

        # Check layout shapes, then master shapes
        sources = []
        try:
            layout = slide.slide_layout
            if layout is not None:
                sources.append(layout)
                try:
                    master = layout.slide_master
                    if master is not None:
                        sources.append(master)
                except Exception:
                    pass
        except Exception:
            pass

        for source in sources:
            bg_image = None
            overlay_color = None
            overlay_opacity = None

            try:
                for shape in source.shapes:
                    # Check if shape covers most of the slide
                    shape_area = shape.width * shape.height
                    if shape_area < threshold:
                        continue

                    # Position should be near (0,0)
                    if shape.left > slide_w * 0.1 or shape.top > slide_h * 0.1:
                        continue

                    # Picture shape -> background image
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            image = shape.image
                            bg_image = {
                                'image_bytes': image.blob,
                                'image_type': image.ext,
                            }
                        except (ValueError, AttributeError):
                            continue

                    # Non-picture shape with blipFill (image fill) -> also background image
                    elif bg_image is None and hasattr(shape, '_element'):
                        blip = self._extract_blip_fill(shape)
                        if blip is not None:
                            bg_image = blip

                    # Auto-shape with solidFill -> potential overlay
                    # (semi-transparent rectangle on top of background image)
                    elif hasattr(shape, 'fill') and shape.fill:
                        try:
                            if hasattr(shape.fill, 'fore_color'):
                                fc = shape.fill.fore_color
                                color = None
                                if hasattr(fc, 'type') and fc.type is not None:
                                    if fc.type == 1 and hasattr(fc, 'rgb') and fc.rgb:
                                        color = f"#{fc.rgb}"
                                    elif fc.type == 2 and self.theme_color_extractor:
                                        color = self.theme_color_extractor.get_theme_color(int(fc.theme_color))

                                if color:
                                    # Check for alpha in XML
                                    alpha = self._extract_shape_fill_alpha(shape)
                                    if alpha is not None and alpha < 1.0:
                                        overlay_color = color
                                        overlay_opacity = alpha
                        except (AttributeError, TypeError):
                            pass
            except Exception:
                continue

            if bg_image:
                return SlideBackground(
                    bg_type='image',
                    image_data=bg_image,
                    overlay_color=overlay_color,
                    overlay_opacity=overlay_opacity,
                )

        return None

    def _extract_shape_fill_alpha(self, shape) -> Optional[float]:
        """
        Extract alpha value from a shape's solidFill.

        Args:
            shape: python-pptx shape object

        Returns:
            Alpha as 0.0-1.0, or None if no alpha specified
        """
        if not hasattr(shape, '_element'):
            return None

        ns = self._ns
        try:
            # Look for <a:solidFill><a:srgbClr><a:alpha val="..."/>
            # or <a:solidFill><a:schemeClr><a:alpha val="..."/>
            alpha_elem = shape._element.find('.//a:solidFill/*/a:alpha', ns)
            if alpha_elem is not None:
                val = alpha_elem.get('val')
                if val:
                    # OpenXML alpha: 0-100000 where 100000 = 100%
                    return int(val) / 100000.0
        except Exception:
            pass
        return None

    def _extract_master_text_styles(self):
        """Extract default text styles from the first slide master's txStyles element."""
        try:
            prs = Presentation(str(self.pptx_path))
            if not prs.slide_masters:
                return
            master = prs.slide_masters[0]
            title_styles, body_styles = self._get_master_text_styles_for(master)

            # Set as global defaults (for backward compat) and cache
            self.master_title_styles = title_styles
            self.master_body_styles = body_styles
            self._master_text_styles_cache[id(master)] = (title_styles, body_styles)

        except Exception:
            pass

    def _get_master_text_styles_for(self, master) -> tuple:
        """
        Get title and body text styles for a specific master, using cache.

        Args:
            master: python-pptx slide master object

        Returns:
            Tuple of (title_styles_dict, body_styles_dict)
        """
        mid = id(master)
        if mid in self._master_text_styles_cache:
            return self._master_text_styles_cache[mid]

        title_styles = {}
        body_styles = {}
        try:
            ns = self._ns
            elem = master._element

            title_style = elem.find('.//p:titleStyle', ns)
            if title_style is not None:
                title_styles = self._parse_level_styles(title_style, ns)

            body_style = elem.find('.//p:bodyStyle', ns)
            if body_style is not None:
                body_styles = self._parse_level_styles(body_style, ns)

        except Exception:
            pass

        self._master_text_styles_cache[mid] = (title_styles, body_styles)
        return title_styles, body_styles

    def _parse_level_styles(self, style_elem, ns) -> Dict[int, 'MasterTextStyle']:
        """
        Parse paragraph-level styles from a titleStyle or bodyStyle element.

        Args:
            style_elem: <p:titleStyle> or <p:bodyStyle> XML element
            ns: XML namespaces dict

        Returns:
            Dict mapping level (int) to MasterTextStyle
        """
        styles = {}
        for level in range(1, 10):
            lvl_pr = style_elem.find(f'a:lvl{level}pPr', ns)
            if lvl_pr is None:
                continue

            def_rpr = lvl_pr.find('a:defRPr', ns)
            if def_rpr is None:
                continue

            text_style = MasterTextStyle()

            # Font size: sz attribute is in hundredths of a point (4400 = 44pt)
            sz = def_rpr.get('sz')
            if sz:
                text_style.font_size = int(sz) / 100.0

            # Bold
            b = def_rpr.get('b')
            if b == '1':
                text_style.bold = True
            elif b == '0':
                text_style.bold = False

            # Italic
            i = def_rpr.get('i')
            if i == '1':
                text_style.italic = True
            elif i == '0':
                text_style.italic = False

            # Font name from <a:latin typeface="...">
            latin = def_rpr.find('a:latin', ns)
            if latin is not None:
                typeface = latin.get('typeface')
                if typeface and not typeface.startswith('+'):
                    # Skip theme references like "+mj-lt", "+mn-lt"
                    text_style.font_name = typeface

            # Color from <a:solidFill>
            solid_fill = def_rpr.find('a:solidFill', ns)
            if solid_fill is not None:
                text_style.color = self._resolve_xml_color(solid_fill)

            # Bullet properties from the level style
            bu_none = lvl_pr.find('a:buNone', ns)
            bu_char = lvl_pr.find('a:buChar', ns)
            bu_auto = lvl_pr.find('a:buAutoNum', ns)

            if bu_none is not None:
                text_style.bullet_type = None  # Explicitly no bullet
            elif bu_char is not None:
                text_style.bullet_type = 'char'
                text_style.bullet_char = bu_char.get('char', '\u2022')
            elif bu_auto is not None:
                text_style.bullet_type = 'autonum'
                text_style.bullet_autonum_type = bu_auto.get('type', 'arabicPeriod')

            bu_font = lvl_pr.find('a:buFont', ns)
            if bu_font is not None:
                text_style.bullet_font = bu_font.get('typeface')

            styles[level] = text_style

        return styles

    def parse(self) -> List[ParsedSlide]:
        """
        Parse the PPTX file and return list of parsed slides.

        Returns:
            List of parsed slides with their elements
        """
        self.presentation = Presentation(str(self.pptx_path))
        self.slides = []

        for slide_idx, slide in enumerate(self.presentation.slides, start=1):
            parsed_slide = self._parse_slide(slide, slide_idx)
            # Check if slide is hidden
            parsed_slide.hidden = self._is_slide_hidden(slide)
            # Detect animations
            parsed_slide.has_animations = self._detect_slide_animations(slide)
            self.slides.append(parsed_slide)

        return self.slides

    def _parse_slide(self, slide, slide_number: int) -> ParsedSlide:
        """
        Parse a single slide and extract all elements.

        Args:
            slide: python-pptx slide object
            slide_number: Slide number (1-indexed)

        Returns:
            ParsedSlide object with all elements
        """
        # Get slide dimensions from presentation (convert EMU to pixels)
        # PPTX uses EMU (English Metric Units): 1 inch = 914400 EMU
        # Standard screen: 1 inch = 96 pixels
        # Conversion: pixels = EMU / 914400 * 96 = EMU / 9525
        slide_width = emu_to_px(self.presentation.slide_width)
        slide_height = emu_to_px(self.presentation.slide_height)

        # Get layout name
        layout_name = slide.slide_layout.name if slide.slide_layout else "Blank"

        parsed_slide = ParsedSlide(
            slide_number=slide_number,
            layout_name=layout_name,
            width=slide_width,
            height=slide_height,
        )

        # Extract background following inheritance chain (slide -> layout -> master)
        parsed_slide.background = self._extract_slide_background(slide)

        # Set current master for per-master style lookups during shape parsing
        self._current_master = None
        try:
            layout = slide.slide_layout
            if layout is not None:
                self._current_master = layout.slide_master
        except Exception:
            pass

        # Extract all shapes from the slide
        for z_order, shape in enumerate(slide.shapes):
            result = self._parse_shape(shape, z_order)
            if result is not None:
                if isinstance(result, list):
                    parsed_slide.elements.extend(result)
                else:
                    parsed_slide.elements.append(result)

        # Extract non-background shapes from layout and master
        # (logos, decorative elements, etc. that appear on every slide)
        self._add_inherited_shapes(slide, parsed_slide)

        return parsed_slide

    def _add_inherited_shapes(self, slide, parsed_slide: ParsedSlide):
        """
        Add non-background shapes from the slide's layout and master.

        PowerPoint renders layout/master shapes (logos, decorative elements)
        behind slide content. We skip shapes that are part of the background
        (full-slide images/overlays) since those are handled by background
        extraction.

        Args:
            slide: python-pptx slide object
            parsed_slide: ParsedSlide to add inherited shapes to
        """
        slide_w = self.presentation.slide_width if self.presentation else 12192000
        slide_h = self.presentation.slide_height if self.presentation else 6858000
        slide_area = slide_w * slide_h

        # Check if master shapes are hidden
        show_master_sp = slide._element.get('showMasterSp')
        if show_master_sp == '0':
            return

        sources = []
        try:
            layout = slide.slide_layout
            if layout is not None:
                sources.append(layout)
                try:
                    master = layout.slide_master
                    if master is not None:
                        sources.append(master)
                except Exception:
                    pass
        except Exception:
            return

        # Get the starting z-order for inherited shapes (behind slide shapes)
        base_z = 0  # Inherited shapes go behind (lower z-order)

        for source in sources:
            try:
                for shape in source.shapes:
                    shape_area = shape.width * shape.height
                    area_ratio = shape_area / slide_area if slide_area > 0 else 0

                    # Skip full-slide background shapes (handled by background extraction)
                    if area_ratio > 0.9 and shape.left < slide_w * 0.1 and shape.top < slide_h * 0.1:
                        continue

                    # Skip placeholder shapes (title, body, etc.)
                    if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                        continue

                    # Parse the shape
                    result = self._parse_shape(shape, base_z)
                    if result is not None:
                        if isinstance(result, list):
                            parsed_slide.elements.extend(result)
                            base_z -= len(result)
                        else:
                            parsed_slide.elements.append(result)
                            base_z -= 1  # Each inherited shape gets lower z-order
            except Exception:
                continue

    def _shape_has_text(self, shape) -> bool:
        """Check if a shape has actual text content in its text_frame."""
        if not (hasattr(shape, "text_frame") and shape.text_frame):
            return False
        if hasattr(shape.text_frame, "text") and shape.text_frame.text.strip():
            return True
        for para in shape.text_frame.paragraphs:
            if hasattr(para, "text") and para.text.strip():
                return True
        return False

    def _parse_shape(self, shape, z_order: int) -> Optional[Union[SlideElement, List[SlideElement]]]:
        """
        Parse a single shape and return the appropriate element.

        For geometric shapes (FREEFORM, AUTO_SHAPE, etc.) that also contain text,
        returns a list of [ShapeElement, TextElement] so both the visual geometry
        and the rich text content are rendered.

        Args:
            shape: python-pptx shape object
            z_order: Z-order of the shape

        Returns:
            SlideElement, list of SlideElements, or None if shape type is not supported
        """
        # Convert EMU to pixels for element positions
        left = emu_to_px(shape.left)
        top = emu_to_px(shape.top)
        width = emu_to_px(shape.width)
        height = emu_to_px(shape.height)

        # Table (check BEFORE placeholder - tables inside placeholders must be handled as tables)
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return self._parse_table(shape, left, top, width, height, z_order)

        # Placeholder shape (non-table placeholders)
        if shape.is_placeholder:
            return self._parse_placeholder(shape, left, top, width, height, z_order)

        # Picture/Image
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_element = self._parse_image(shape, left, top, width, height, z_order)
            if image_element is None:
                # If image parsing fails, create a placeholder shape
                return ShapeElement(
                    element_type="shape",
                    left=left,
                    top=top,
                    width=width,
                    height=height,
                    z_order=z_order,
                    shape_type="PICTURE_PLACEHOLDER",
                    fill_color="#CCCCCC",
                    line_color="#666666",
                    text="[Image]",
                )
            return image_element

        # Group shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            group_element = self._parse_group(shape, z_order)
            if group_element is None:
                # If group parsing fails, create a placeholder shape
                return ShapeElement(
                    element_type="shape",
                    left=left,
                    top=top,
                    width=width,
                    height=height,
                    z_order=z_order,
                    shape_type="GROUP_PLACEHOLDER",
                    fill_color="#DDDDDD",
                    line_color="#999999",
                    text=f"[Group with {len(shape.shapes)} items]" if hasattr(shape, 'shapes') else "[Group]",
                )
            return group_element

        # Pure text box shapes - no visual geometry to render
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            if self._shape_has_text(shape):
                return self._parse_text_box(shape, left, top, width, height, z_order)
            return None

        # All other shapes (FREEFORM, AUTO_SHAPE, LINE, CALLOUT, etc.)
        # These have visual geometry that must be rendered
        shape_element = self._parse_generic_shape(shape, left, top, width, height, z_order)

        # If the shape also has text, render both geometry AND rich text
        if self._shape_has_text(shape) and shape_element:
            text_element = self._parse_text_box(shape, left, top, width, height, z_order + 1)
            # Clear text on shape element to avoid duplicate text rendering
            shape_element.text = None

            # Only include shape element if it has visible geometry
            # (fill, stroke, image fill, or custom geometry like FREEFORM paths)
            has_visible_geometry = (
                shape_element.fill_color is not None
                or shape_element.line_color is not None
                or (hasattr(shape_element, 'blip_fill') and shape_element.blip_fill)
                or (shape_element.metadata and shape_element.metadata.get('geometry'))
            )
            if has_visible_geometry:
                return [shape_element, text_element]
            else:
                # No visible geometry — just return the text element
                return text_element

        return shape_element

    def _parse_placeholder(
        self, shape, left: float, top: float, width: float, height: float, z_order: int
    ) -> TextElement:
        """Parse a placeholder shape with paragraph support."""
        text = ""
        paragraphs = []

        if hasattr(shape, "text_frame") and shape.text_frame:
            text_frame = shape.text_frame
            text = sanitize_pptx_text(text_frame.text)  # Sanitize for backward compatibility

            # Extract all paragraphs with their formatting
            for para in text_frame.paragraphs:
                # Extract paragraph text and sanitize it
                para_text = sanitize_pptx_text(para.text)

                # Get paragraph level
                para_level = para.level if hasattr(para, 'level') else 0

                # Get paragraph alignment
                alignment = None
                if hasattr(para, 'alignment') and para.alignment is not None:
                    alignment_map = {1: 'left', 2: 'center', 3: 'right', 4: 'justify', 5: 'distribute'}
                    alignment = alignment_map.get(int(para.alignment), None)

                # Extract paragraph spacing
                line_spacing = None
                line_spacing_pts = None
                spacing_before = None
                spacing_after = None

                try:
                    ls = para.line_spacing
                    if ls is not None:
                        if isinstance(ls, float):
                            line_spacing = ls
                        else:
                            line_spacing_pts = ls.pt
                except Exception:
                    pass

                try:
                    sb = para.space_before
                    if sb is not None:
                        spacing_before = sb.pt
                except Exception:
                    pass

                try:
                    sa = para.space_after
                    if sa is not None:
                        spacing_after = sa.pt
                except Exception:
                    pass

                # Extract formatting from all runs
                font_name = None
                font_size = None
                bold = False
                italic = False
                color = None
                runs = []

                if para.runs:
                    for run in para.runs:
                        font = run.font
                        run_bold = font.bold
                        run_italic = font.italic
                        run_font_name = font.name
                        run_font_size = font.size.pt if font.size else None
                        run_color = self._extract_run_color(font)
                        run_text = sanitize_pptx_text(run.text)

                        runs.append(RunElement(
                            text=run_text,
                            bold=run_bold,
                            italic=run_italic,
                            font_name=run_font_name,
                            font_size=run_font_size,
                            color=run_color,
                        ))

                    # Paragraph-level defaults from first run
                    first = runs[0]
                    font_name = first.font_name
                    font_size = first.font_size
                    bold = first.bold if first.bold is not None else False
                    italic = first.italic if first.italic is not None else False
                    color = first.color

                # Apply master defaults when paragraph has no explicit spacing
                if line_spacing is None and line_spacing_pts is None:
                    if self._current_master is not None:
                        line_spacing = self._get_master_spacing_for(self._current_master)
                    if line_spacing is None:
                        line_spacing = self.default_line_spacing

                # Apply master text defaults for missing styles (placeholder uses
                # master title/body styles depending on placeholder type)
                # We check is_title later, so apply body styles as default here
                # and override with title styles after is_title is determined
                _, current_body_styles = self._get_current_master_styles()
                if current_body_styles and para_level in current_body_styles:
                    master_style = current_body_styles[para_level]
                    if font_name is None and master_style.font_name:
                        font_name = master_style.font_name
                    if font_size is None and master_style.font_size:
                        font_size = master_style.font_size
                    if color is None and master_style.color:
                        color = master_style.color

                # Parse bullet properties from OpenXML
                bullet = None
                margin_left_pt = None
                indent_pt = None
                try:
                    p_elem = para._p  # The <a:p> XML element
                    ppr = p_elem.find('a:pPr', self._ns)
                    bullet = self._parse_bullet_properties(ppr, para_level, shape)
                    if ppr is not None:
                        marL = ppr.get('marL')
                        if marL:
                            margin_left_pt = emu_to_pt(int(marL))
                        indent_val = ppr.get('indent')
                        if indent_val:
                            indent_pt = emu_to_pt(int(indent_val))
                except Exception:
                    pass

                paragraphs.append(ParagraphElement(
                    text=para_text,
                    level=para_level,
                    alignment=alignment,
                    font_name=font_name,
                    font_size=font_size,
                    bold=bold,
                    italic=italic,
                    color=color,
                    line_spacing=line_spacing,
                    line_spacing_pts=line_spacing_pts,
                    spacing_before=spacing_before,
                    spacing_after=spacing_after,
                    margin_left=margin_left_pt,
                    indent=indent_pt,
                    bullet=bullet,
                    runs=runs,
                ))

        # Determine if this is a title
        is_title = shape.placeholder_format.type in (
            0,  # Title
            14,  # Centered Title
        ) if hasattr(shape, "placeholder_format") else False

        # Resolve bullet inheritance for all paragraphs
        for p in paragraphs:
            if p.bullet and p.bullet.type == 'inherited':
                p.bullet = self._resolve_bullet_inheritance(
                    p.bullet, p.level, is_title=is_title, is_placeholder=True
                )

        # Extract rotation and transform information
        rotation = None
        vert = None
        flip_h = False
        flip_v = False

        if hasattr(shape, "_element"):
            elem = shape._element
            ns = NS_A

            # Check body properties for vertical text and alignment
            body_pr = elem.find('.//a:bodyPr', ns)
            vertical_align = None
            if body_pr is not None:
                vert = body_pr.get('vert')  # eaVert, mongolianVert, etc.

                # Extract vertical alignment from anchor attribute
                anchor = body_pr.get('anchor')
                if anchor:
                    anchor_map = {'t': 'top', 'ctr': 'middle', 'b': 'bottom'}
                    vertical_align = anchor_map.get(anchor)

            # Check transformation for flip and rotation
            xfrm = elem.find('.//a:xfrm', ns)
            if xfrm is not None:
                flip_h = xfrm.get('flipH') == '1'
                flip_v = xfrm.get('flipV') == '1'
                # Rotation is in EMU (1/60,000 of a degree)
                rot_emu = xfrm.get('rot')
                if rot_emu:
                    try:
                        rotation = angle_to_degrees(float(rot_emu))
                    except (ValueError, TypeError):
                        pass

        return TextElement(
            element_type="text",
            left=left,
            top=top,
            width=width,
            height=height,
            z_order=z_order,
            text=text,
            paragraphs=paragraphs,
            is_title=is_title,
            metadata={"placeholder_type": shape.placeholder_format.type if hasattr(shape, "placeholder_format") else None},
            rotation=rotation,
            vert=vert,
            flip_h=flip_h,
            flip_v=flip_v,
            vertical_align=vertical_align,
        )

    def _parse_text_box(
        self, shape, left: float, top: float, width: float, height: float, z_order: int
    ) -> TextElement:
        """Parse a text box shape with paragraph-level structure."""
        text_frame = shape.text_frame
        text = sanitize_pptx_text(text_frame.text)  # Sanitize for backward compatibility

        # Extract all paragraphs with their formatting
        paragraphs = []
        for para in text_frame.paragraphs:
            # Extract paragraph text and sanitize it
            para_text = sanitize_pptx_text(para.text)

            # Get paragraph level
            para_level = para.level if hasattr(para, 'level') else 0

            # Get paragraph alignment
            alignment = None
            if hasattr(para, 'alignment') and para.alignment is not None:
                alignment_map = {1: 'left', 2: 'center', 3: 'right', 4: 'justify', 5: 'distribute'}
                alignment = alignment_map.get(int(para.alignment), None)

            # Extract paragraph spacing
            line_spacing = None
            line_spacing_pts = None
            spacing_before = None
            spacing_after = None

            try:
                ls = para.line_spacing
                if ls is not None:
                    if isinstance(ls, float):
                        line_spacing = ls
                    else:
                        line_spacing_pts = ls.pt
            except Exception:
                pass

            try:
                sb = para.space_before
                if sb is not None:
                    spacing_before = sb.pt
            except Exception:
                pass

            try:
                sa = para.space_after
                if sa is not None:
                    spacing_after = sa.pt
            except Exception:
                pass

            # Extract formatting from all runs
            font_name = None
            font_size = None
            bold = False
            italic = False
            color = None
            runs = []

            if para.runs:
                for run in para.runs:
                    font = run.font
                    run_bold = font.bold
                    run_italic = font.italic
                    run_font_name = font.name
                    run_font_size = font.size.pt if font.size else None
                    run_color = self._extract_run_color(font)
                    run_text = sanitize_pptx_text(run.text)

                    runs.append(RunElement(
                        text=run_text,
                        bold=run_bold,
                        italic=run_italic,
                        font_name=run_font_name,
                        font_size=run_font_size,
                        color=run_color,
                    ))

                # Paragraph-level defaults from first run
                first = runs[0]
                font_name = first.font_name
                font_size = first.font_size
                bold = first.bold if first.bold is not None else False
                italic = first.italic if first.italic is not None else False
                color = first.color

            # Apply master defaults when paragraph has no explicit spacing
            if line_spacing is None and line_spacing_pts is None:
                line_spacing = self.default_line_spacing

            # Apply master text defaults for missing styles (text boxes use body styles)
            _, current_body_styles = self._get_current_master_styles()
            if current_body_styles and para_level in current_body_styles:
                master_style = current_body_styles[para_level]
                if font_name is None and master_style.font_name:
                    font_name = master_style.font_name
                if font_size is None and master_style.font_size:
                    font_size = master_style.font_size
                if color is None and master_style.color:
                    color = master_style.color

            # Parse bullet properties from OpenXML
            bullet = None
            try:
                p_elem = para._p  # The <a:p> XML element
                ppr = p_elem.find('a:pPr', self._ns)
                bullet = self._parse_bullet_properties(ppr, para_level, shape)
                margin_left_pt = None
                indent_pt = None
                if ppr is not None:
                    marL = ppr.get('marL')
                    if marL:
                        margin_left_pt = emu_to_pt(int(marL))
                    indent_val = ppr.get('indent')
                    if indent_val:
                        indent_pt = emu_to_pt(int(indent_val))
            except Exception:
                pass

            paragraphs.append(ParagraphElement(
                text=para_text,
                level=para_level,
                alignment=alignment,
                font_name=font_name,
                font_size=font_size,
                bold=bold,
                italic=italic,
                color=color,
                line_spacing=line_spacing,
                line_spacing_pts=line_spacing_pts,
                spacing_before=spacing_before,
                spacing_after=spacing_after,
                margin_left=margin_left_pt,
                indent=indent_pt,
                bullet=bullet,
                runs=runs,
            ))

        # Resolve bullet inheritance for text box paragraphs
        # Non-placeholder text boxes don't inherit bullets from master
        for p in paragraphs:
            if p.bullet and p.bullet.type == 'inherited':
                p.bullet = self._resolve_bullet_inheritance(
                    p.bullet, p.level, is_title=False, is_placeholder=False
                )

        # For backward compatibility, use first paragraph's formatting
        font_name = None
        font_size = None
        bold = False
        italic = False
        color = None

        if paragraphs and paragraphs[0].runs:
            first_para = paragraphs[0]
            font_name = first_para.font_name
            font_size = first_para.font_size
            bold = first_para.bold
            italic = first_para.italic
            color = first_para.color

        # Extract rotation and transform information
        rotation = None
        vert = None
        flip_h = False
        flip_v = False

        if hasattr(shape, "_element"):
            elem = shape._element
            ns = NS_A

            # Check body properties for vertical text and alignment
            body_pr = elem.find('.//a:bodyPr', ns)
            vertical_align = None
            if body_pr is not None:
                vert = body_pr.get('vert')  # eaVert, mongolianVert, etc.

                # Extract vertical alignment from anchor attribute
                anchor = body_pr.get('anchor')
                if anchor:
                    anchor_map = {'t': 'top', 'ctr': 'middle', 'b': 'bottom'}
                    vertical_align = anchor_map.get(anchor)

            # Check transformation for flip and rotation
            xfrm = elem.find('.//a:xfrm', ns)
            if xfrm is not None:
                flip_h = xfrm.get('flipH') == '1'
                flip_v = xfrm.get('flipV') == '1'
                # Rotation is in EMU (1/60,000 of a degree)
                rot_emu = xfrm.get('rot')
                if rot_emu:
                    try:
                        rotation = angle_to_degrees(float(rot_emu))
                    except (ValueError, TypeError):
                        pass

        # Extract outline/border properties (line_color and line_width)
        line_color = None
        line_width = None

        # Check for noFill in XML BEFORE accessing python-pptx line properties
        # Accessing shape.line.color triggers python-pptx to resolve style references,
        # which can modify the XML and replace <noFill/> with <solidFill/>
        _line_has_noFill = False
        if hasattr(shape, '_element'):
            try:
                _ns = {'a': NAMESPACES['a'], 'p': NAMESPACES['p']}
                _spPr = shape._element.find('./p:spPr', _ns)
                if _spPr is not None:
                    _ln = _spPr.find('./a:ln', _ns)
                    if _ln is not None:
                        _noFill = _ln.find('./a:noFill', _ns)
                        if _noFill is not None:
                            _line_has_noFill = True
            except Exception:
                pass

        # Extract line color from shape.line (same logic as _parse_generic_shape)
        if not _line_has_noFill and hasattr(shape, "line") and shape.line:
            try:
                # Try color attribute first
                if hasattr(shape.line, "color") and shape.line.color:
                    line_color_obj = shape.line.color

                    # Try to get theme color first
                    if hasattr(line_color_obj, "theme_color") and line_color_obj.theme_color is not None:
                        theme_color = line_color_obj.theme_color
                        if self.theme_color_extractor:
                            theme_rgb = self.theme_color_extractor.get_theme_color(theme_color)
                            if theme_rgb:
                                line_color = theme_rgb

                    # If no theme color, try direct RGB
                    if not line_color and hasattr(line_color_obj, "rgb") and line_color_obj.rgb:
                        rgb_obj = line_color_obj.rgb
                        # Handle RGBColor objects
                        if hasattr(rgb_obj, '__class__') and 'RGBColor' in str(rgb_obj.__class__):
                            rgb_str = str(rgb_obj).strip()
                            if len(rgb_str) == 6 and rgb_str.isalnum():
                                line_color = f"#{rgb_str}"
                        # Handle string RGB values
                        elif isinstance(rgb_obj, str):
                            rgb_str = rgb_obj.strip()
                            if len(rgb_str) >= 6:
                                line_color = f"#{rgb_str[:6]}"
                        # Handle integer RGB values
                        elif isinstance(rgb_obj, int) and rgb_obj > 0:
                            line_color = f"#{rgb_obj:06x}"
            except (AttributeError, TypeError, ValueError):
                pass

        # If no direct line color, try to extract from style/lnRef (theme style reference)
        if not _line_has_noFill and not line_color and hasattr(shape, '_element'):
            try:
                ns_xml = {'a': NAMESPACES['a'], 'p': NAMESPACES['p']}

                # Find style element first, then lnRef child
                style = shape._element.find('.//p:style', ns_xml)
                if style is not None:
                    ln_ref = style.find('.//a:lnRef', ns_xml)
                    if ln_ref is not None:
                        # Check for schemeClr (theme color reference)
                        scheme_clr = ln_ref.find('.//a:schemeClr', ns_xml)
                        if scheme_clr is not None:
                            val = scheme_clr.get('val')  # e.g., "accent1", "accent2", etc.
                            if val and self.theme_color_extractor:
                                # Convert scheme color name to theme color enum
                                theme_color_name = val.upper()
                                # Add underscore between letter and number (e.g., "ACCENT3" -> "ACCENT_3")
                                theme_color_name = re.sub(r'([A-Z]+)(\d+)', r'\1_\2', theme_color_name)
                                try:
                                    from pptx.enum.dml import MSO_THEME_COLOR
                                    theme_color_enum = getattr(MSO_THEME_COLOR, theme_color_name)
                                    theme_rgb = self.theme_color_extractor.get_theme_color(theme_color_enum)
                                    if theme_rgb:
                                        line_color = theme_rgb
                                except (AttributeError, ValueError):
                                    pass
            except Exception:
                pass

        # Extract line width from XML (<a:ln w="...">)
        if hasattr(shape, '_element'):
            try:
                ns_xml = NS_A
                ln_elem = shape._element.find('.//a:ln', ns_xml)
                if ln_elem is not None:
                    w_attr = ln_elem.get('w')
                    if w_attr:
                        line_width = emu_to_px(int(w_attr))
            except Exception:
                pass

        return TextElement(
            element_type="text",
            left=left,
            top=top,
            width=width,
            height=height,
            z_order=z_order,
            text=text,
            paragraphs=paragraphs,
            font_name=font_name,
            font_size=font_size,
            bold=bold,
            italic=italic,
            color=color,
            rotation=rotation,
            vert=vert,
            flip_h=flip_h,
            flip_v=flip_v,
            vertical_align=vertical_align,
            line_color=line_color,
            line_width=line_width,
        )

    def _parse_table(
        self, shape, left: float, top: float, width: float, height: float, z_order: int
    ) -> TableElement:
        """Parse a table shape."""
        table = shape.table
        rows = len(table.rows)
        cols = len(table.columns)

        # Extract table data
        data = []
        cell_styles = []

        ns = NS_A

        for row in table.rows:
            row_data = []
            row_styles = []
            for cell in row.cells:
                row_data.append(cell.text)
                # Extract cell styling
                cell_style = {
                    "background_color": None,
                    "border_color": None,
                }

                # Try python-pptx fill API first
                bg_color = None
                try:
                    if cell.fill and hasattr(cell.fill, "type") and cell.fill.type is not None:
                        if hasattr(cell.fill, "fore_color") and cell.fill.fore_color is not None:
                            if hasattr(cell.fill.fore_color, "rgb") and cell.fill.fore_color.rgb:
                                bg_color = f"#{cell.fill.fore_color.rgb}"
                except Exception:
                    pass

                # Fallback: parse XML directly for cell background
                if bg_color is None and hasattr(cell, '_tc'):
                    tc_pr = cell._tc.find('a:tcPr', ns)
                    if tc_pr is not None:
                        # Check for solidFill
                        solid_fill = tc_pr.find('a:solidFill', ns)
                        if solid_fill is not None:
                            srgb = solid_fill.find('a:srgbClr', ns)
                            if srgb is not None:
                                bg_color = f"#{srgb.get('val')}"
                            else:
                                scheme_clr = solid_fill.find('a:schemeClr', ns)
                                if scheme_clr is not None:
                                    val = scheme_clr.get('val')
                                    if val and self.theme_color_extractor:
                                        bg_color = self.theme_color_extractor.get_theme_color(val)

                cell_style["background_color"] = bg_color
                row_styles.append(cell_style)
            data.append(row_data)
            cell_styles.append(row_styles)

        return TableElement(
            element_type="table",
            left=left,
            top=top,
            width=width,
            height=height,
            z_order=z_order,
            rows=rows,
            cols=cols,
            data=data,
            cell_styles=cell_styles,
        )

    def _parse_image(
        self, shape, left: float, top: float, width: float, height: float, z_order: int
    ) -> Optional[ImageElement]:
        """Parse an image shape, including srcRect cropping and clrChange effects."""
        try:
            # Get image bytes
            image = shape.image
            image_bytes = image.blob
            image_type = image.ext

            # Get alt text
            alt_text = "" if not hasattr(shape, "alt_text") else shape.alt_text

            src_rect = None
            clr_change = None
            scene3d_camera = None
            fill_mode = "stretch"

            # Extract special properties from XML
            ns = self._ns
            if hasattr(shape, '_element'):
                elem = shape._element

                # Extract srcRect (image cropping)
                src_rect_elem = elem.find('.//a:srcRect', ns)
                if src_rect_elem is not None:
                    src_rect = {}
                    for attr in ['l', 't', 'r', 'b']:
                        val = src_rect_elem.get(attr)
                        if val:
                            src_rect[attr] = int(val)
                    if not src_rect:
                        src_rect = None

                # Extract clrChange (color replacement, typically white→transparent)
                clr_change_elem = elem.find('.//a:blip/a:clrChange', ns)
                if clr_change_elem is not None:
                    clr_from = clr_change_elem.find('a:clrFrom/a:srgbClr', ns)
                    clr_to = clr_change_elem.find('a:clrTo/a:srgbClr', ns)
                    if clr_from is not None and clr_to is not None:
                        from_color = '#' + clr_from.get('val', '')
                        # Check for alpha=0 in clrTo (means transparent)
                        alpha_elems = clr_to.findall('a:alpha', ns)
                        is_transparent = any(
                            int(a.get('val', '100000')) == 0 for a in alpha_elems
                        )
                        if is_transparent:
                            clr_change = {'from': from_color, 'to': 'transparent'}

                # Extract scene3d camera preset
                scene3d = elem.find('.//a:scene3d/a:camera', ns)
                if scene3d is not None:
                    prst = scene3d.get('prst')
                    if prst:
                        scene3d_camera = prst

                # Extract fill mode from <a:blipFill>
                blipFill_elem = elem.find('p:blipFill', ns)
                if blipFill_elem is not None:
                    if blipFill_elem.find('a:stretch', ns) is not None:
                        fill_mode = "stretch"
                    elif blipFill_elem.find('a:tile', ns) is not None:
                        fill_mode = "tile"
                    else:
                        fill_mode = "none"

            # Calculate PPT image scale for scene3d images.
            # PPT "Scale" = shape_EMU / (cropped_img_px * 914400 / img_dpi)
            # This captures user stretching applied before the 3D transform.
            scale_w = None
            scale_h = None
            if scene3d_camera:
                try:
                    from PIL import Image as PILImage
                    import io as _io
                    pil_img = PILImage.open(_io.BytesIO(image.blob))
                    img_w, img_h = pil_img.size
                    dpi = pil_img.info.get('dpi', (96, 96))
                    dpi_x = dpi[0] if dpi and dpi[0] > 0 else 96
                    dpi_y = dpi[1] if dpi and dpi[1] > 0 else 96
                    pil_img.close()

                    # srcRect crop factors (1/100000ths)
                    sl = (src_rect or {}).get('l', 0)
                    sr = (src_rect or {}).get('r', 0)
                    st = (src_rect or {}).get('t', 0)
                    sb = (src_rect or {}).get('b', 0)
                    cropped_w_px = img_w * (100000 - sl - sr) / 100000
                    cropped_h_px = img_h * (100000 - st - sb) / 100000

                    # Cropped image size in EMU (using image DPI)
                    cropped_w_emu = cropped_w_px * EMU_PER_INCH / dpi_x
                    cropped_h_emu = cropped_h_px * EMU_PER_INCH / dpi_y

                    # Shape dimensions in EMU (width/height are in px = EMU/9525)
                    shape_w_emu = px_to_emu(width)
                    shape_h_emu = px_to_emu(height)

                    if cropped_w_emu > 0 and cropped_h_emu > 0:
                        scale_w = round(shape_w_emu / cropped_w_emu, 4)
                        scale_h = round(shape_h_emu / cropped_h_emu, 4)
                except Exception:
                    pass

            # Apply srcRect cropping with Pillow
            if src_rect:
                image_bytes = self._crop_image_src_rect(image_bytes, src_rect)
                # Update dimensions after cropping - the shape's width/height
                # in PPTX already reflects the cropped area, so no change needed

            # Apply clrChange with Pillow
            if clr_change and clr_change.get('to') == 'transparent':
                image_bytes = self._apply_color_change(
                    image_bytes, clr_change['from'], tolerance=30
                )

            # Ensure image_type is PNG after Pillow processing (transparency requires PNG)
            if clr_change or src_rect:
                image_type = 'png'

            element = ImageElement(
                element_type="image",
                left=left,
                top=top,
                width=width,
                height=height,
                z_order=z_order,
                image_bytes=image_bytes,
                image_type=image_type,
                alt_text=alt_text,
                src_rect=src_rect,
                clr_change=clr_change,
                fill_mode=fill_mode,
                scale_w=scale_w,
                scale_h=scale_h,
            )

            # Store scene3d in metadata for CSS rendering
            if scene3d_camera:
                element.metadata = element.metadata or {}
                element.metadata['scene3d_camera'] = scene3d_camera

            return element
        except (ValueError, AttributeError):
            # No embedded image or other error
            return None

    def _crop_image_src_rect(self, image_bytes: bytes, src_rect: dict) -> bytes:
        """Crop image according to OpenXML srcRect percentages.

        srcRect values are in 1/100000ths of the image dimension.
        l=10000 means crop 10% from left, r=20000 means crop 20% from right.
        """
        try:
            from PIL import Image
            import io

            img = Image.open(io.BytesIO(image_bytes))
            w, h = img.size

            l = int(src_rect.get('l', 0) / 100000 * w)
            t = int(src_rect.get('t', 0) / 100000 * h)
            r = int(src_rect.get('r', 0) / 100000 * w)
            b = int(src_rect.get('b', 0) / 100000 * h)

            # Clamp to image bounds
            l = max(0, min(l, w))
            t = max(0, min(t, h))
            r = max(0, min(r, w - l))
            b = max(0, min(b, h - t))

            cropped = img.crop((l, t, w - r, h - b))
            output = io.BytesIO()
            fmt = 'PNG' if img.mode == 'RGBA' else (img.format or 'PNG')
            cropped.save(output, format=fmt)
            return output.getvalue()
        except Exception:
            return image_bytes

    def _apply_color_change(self, image_bytes: bytes, from_color: str,
                            tolerance: int = 30) -> bytes:
        """Apply color replacement to image: make near-'from_color' pixels transparent."""
        try:
            from PIL import Image
            import io

            img = Image.open(io.BytesIO(image_bytes)).convert('RGBA')
            pixels = img.load()

            # Parse target color
            r_target = int(from_color[1:3], 16)
            g_target = int(from_color[3:5], 16)
            b_target = int(from_color[5:7], 16)

            w, h = img.size
            for y in range(h):
                for x in range(w):
                    r, g, b, a = pixels[x, y]
                    if (abs(r - r_target) <= tolerance and
                            abs(g - g_target) <= tolerance and
                            abs(b - b_target) <= tolerance):
                        pixels[x, y] = (r, g, b, 0)  # Make transparent

            output = io.BytesIO()
            img.save(output, format='PNG')
            return output.getvalue()
        except Exception:
            return image_bytes

    def _parse_group(self, shape, z_order: int, depth: int = 0) -> Optional[GroupElement]:
        """Parse a group shape by recursively parsing children with coordinate transformation.

        The group has two coordinate systems:
        - Group position (off/ext): where the group sits on the slide
        - Child coordinate system (chOff/chExt): the coordinate space children use

        Children's coordinates are transformed: slide_pos = group_off + (child_pos - chOff) * scale
        Then made relative to the group's top-left for CSS rendering inside a container.
        """
        if depth > 5:
            return None  # Safety: prevent infinite recursion

        try:
            ns = self._ns
            elem = shape._element

            # Extract group transform from XML
            xfrm = elem.find('.//p:grpSpPr/a:xfrm', ns)
            if xfrm is None:
                return None

            off = xfrm.find('a:off', ns)
            ext = xfrm.find('a:ext', ns)
            ch_off = xfrm.find('a:chOff', ns)
            ch_ext = xfrm.find('a:chExt', ns)

            if off is None or ext is None:
                return None

            group_off_x = int(off.get('x', '0'))
            group_off_y = int(off.get('y', '0'))
            group_ext_cx = int(ext.get('cx', '0'))
            group_ext_cy = int(ext.get('cy', '0'))

            # Child coordinate system - defaults to group extent if not specified
            co_x = int(ch_off.get('x', '0')) if ch_off is not None else 0
            co_y = int(ch_off.get('y', '0')) if ch_off is not None else 0
            ce_cx = int(ch_ext.get('cx', str(group_ext_cx))) if ch_ext is not None else group_ext_cx
            ce_cy = int(ch_ext.get('cy', str(group_ext_cy))) if ch_ext is not None else group_ext_cy

            # Compute scale factors
            scale_x = group_ext_cx / ce_cx if ce_cx else 1.0
            scale_y = group_ext_cy / ce_cy if ce_cy else 1.0

            # Group position in pixels
            group_left_px = emu_to_px(group_off_x)
            group_top_px = emu_to_px(group_off_y)
            group_width_px = emu_to_px(group_ext_cx)
            group_height_px = emu_to_px(group_ext_cy)

            # Parse children recursively
            children = []
            child_z = z_order

            if not hasattr(shape, 'shapes'):
                return None

            for child_shape in shape.shapes:
                child_elem = self._parse_shape(child_shape, child_z)
                if child_elem is None:
                    continue

                # _parse_shape may return a list (e.g. shape + text)
                if isinstance(child_elem, list):
                    for item in child_elem:
                        self._transform_child_to_group_relative(
                            item, group_off_x, group_off_y,
                            scale_x, scale_y, co_x, co_y,
                            group_left_px, group_top_px
                        )
                        children.append(item)
                        child_z += 1
                else:
                    self._transform_child_to_group_relative(
                        child_elem, group_off_x, group_off_y,
                        scale_x, scale_y, co_x, co_y,
                        group_left_px, group_top_px
                    )
                    children.append(child_elem)
                    child_z += 1

            return GroupElement(
                element_type="group",
                left=group_left_px,
                top=group_top_px,
                width=group_width_px,
                height=group_height_px,
                z_order=z_order,
                children=children,
                metadata={"group_shape": True, "child_count": len(children)}
            )

        except Exception:
            return None

    def _transform_child_to_group_relative(self, child, group_off_x, group_off_y,
                                            scale_x, scale_y, ch_off_x, ch_off_y,
                                            group_left_px, group_top_px):
        """Transform a child element's coordinates from group child-space to group-relative pixels.

        Steps:
        1. Convert child's pixel coords back to EMU
        2. Map from child-space to slide-space using group transform
        3. Convert back to pixels
        4. Make relative to group's top-left
        """
        # Convert child's px back to EMU
        child_x_emu = px_to_emu(child.left)
        child_y_emu = px_to_emu(child.top)
        child_w_emu = px_to_emu(child.width)
        child_h_emu = px_to_emu(child.height)

        # Map from child-space to slide-space
        slide_x = group_off_x + (child_x_emu - ch_off_x) * scale_x
        slide_y = group_off_y + (child_y_emu - ch_off_y) * scale_y
        slide_w = child_w_emu * scale_x
        slide_h = child_h_emu * scale_y

        # Convert to pixels and make relative to group top-left
        child.left = emu_to_px(slide_x) - group_left_px
        child.top = emu_to_px(slide_y) - group_top_px
        child.width = emu_to_px(slide_w)
        child.height = emu_to_px(slide_h)

        # Adjust image scale for group coordinate transform.
        # scale_w/h was calculated in child-space; the group scale factor
        # converts it to visual (slide-space) scale.
        if isinstance(child, ImageElement) and child.scale_w is not None:
            child.scale_w = child.scale_w * scale_x
            child.scale_h = child.scale_h * scale_y

        # For nested groups, propagate the parent group's scale to all
        # descendants.  Without this, nested-group children only carry the
        # inner group's scale and miss the outer group's scale, causing
        # dimensions to be off (e.g. 2x too large when outer scale = 0.5).
        if isinstance(child, GroupElement):
            self._propagate_group_scale(child, scale_x, scale_y)

    def _propagate_group_scale(self, group: GroupElement, scale_x: float, scale_y: float):
        """Recursively scale all descendants of a group by the parent's scale factors."""
        for sub in group.children:
            sub.left *= scale_x
            sub.top *= scale_y
            sub.width *= scale_x
            sub.height *= scale_y
            if isinstance(sub, GroupElement):
                self._propagate_group_scale(sub, scale_x, scale_y)

    def _parse_generic_shape(
        self, shape, left: float, top: float, width: float, height: float, z_order: int
    ) -> Optional[ShapeElement]:
        """Parse a generic shape (rectangles, circles, lines, etc.)."""
        # Get shape type name
        # python-pptx EnumValue doesn't have .name; str() gives e.g. "AUTO_SHAPE (1)"
        shape_type_str = str(shape.shape_type)
        # Strip the enum value suffix like " (1)" for cleaner type names
        shape_type = re.sub(r'\s*\(\d+\)$', '', shape_type_str)

        # For AUTO_SHAPE, extract the actual preset geometry name from XML
        # (e.g., "cloudCallout", "roundedRectCallout") which is more specific
        preset_name = None
        if shape_type == "AUTO_SHAPE" and hasattr(shape, "_element"):
            try:
                ns = NS_A
                prst_geom = shape._element.find('.//a:prstGeom', ns)
                if prst_geom is not None:
                    preset_name = prst_geom.get('prst')
                    if preset_name:
                        shape_type = preset_name  # Use specific preset name instead of generic "AUTO_SHAPE"
            except Exception:
                pass

        # Similarly for FREEFORM shapes, try to detect if they have custom geometry
        if shape_type == "FREEFORM" and hasattr(shape, "_element"):
            # FREEFORM shapes already have custGeom which will be handled by geometry extraction
            pass

        # Extract fill and line colors
        fill_color = None
        line_color = None

        # Detect explicit noFill (e.g., FREEFORM shapes with stroke-only outlines)
        # When fill.type is BACKGROUND, it means "no fill, show background" (<a:noFill/>)
        if hasattr(shape, "fill") and shape.fill:
            try:
                if shape.fill.type is not None and str(shape.fill.type) == 'BACKGROUND (5)':
                    fill_color = "none"
            except (AttributeError, TypeError):
                pass

        # Extract fill color (handle different fill types, theme colors, and RGBColor objects)
        if fill_color is None and hasattr(shape, "fill") and shape.fill:
            try:
                # Check if fill has fore_color attribute
                if hasattr(shape.fill, "fore_color"):
                    fore_color = shape.fill.fore_color

                    # Try to get theme color first (if it's a theme color reference)
                    if hasattr(fore_color, "theme_color") and fore_color.theme_color is not None:
                        theme_color = fore_color.theme_color
                        # Try to get theme color value
                        if self.theme_color_extractor:
                            # Convert theme color enum to RGB
                            theme_rgb = self.theme_color_extractor.get_theme_color(theme_color)
                            if theme_rgb:
                                fill_color = theme_rgb

                    # If no theme color, try direct RGB
                    if not fill_color and hasattr(fore_color, "rgb") and fore_color.rgb:
                        rgb_obj = fore_color.rgb
                        # Handle RGBColor objects (pptx.dml.color.RGBColor)
                        if hasattr(rgb_obj, '__class__') and 'RGBColor' in str(rgb_obj.__class__):
                            # RGBColor object - convert to string representation
                            rgb_str = str(rgb_obj).strip()
                            if len(rgb_str) == 6 and rgb_str.isalnum():
                                fill_color = f"#{rgb_str}"
                        # Handle string RGB values
                        elif isinstance(rgb_obj, str):
                            rgb_str = rgb_obj.strip()
                            if len(rgb_str) >= 6:
                                fill_color = f"#{rgb_str[:6]}"
                        # Handle integer RGB values
                        elif isinstance(rgb_obj, int) and rgb_obj > 0:
                            fill_color = f"#{rgb_obj:06x}"
            except (AttributeError, TypeError, ValueError) as e:
                fill_color = None

        # Check for noFill in XML BEFORE accessing python-pptx line properties
        # Accessing shape.line.color triggers python-pptx to resolve style references,
        # which can modify the XML and replace <noFill/> with <solidFill/>
        _line_has_noFill = False
        if hasattr(shape, '_element'):
            try:
                _ns = {'a': NAMESPACES['a'], 'p': NAMESPACES['p']}
                _spPr = shape._element.find('./p:spPr', _ns)
                if _spPr is not None:
                    _ln = _spPr.find('./a:ln', _ns)
                    if _ln is not None:
                        _noFill = _ln.find('./a:noFill', _ns)
                        if _noFill is not None:
                            _line_has_noFill = True
            except Exception:
                pass

        # Extract line color (handle theme colors, different color types, and RGBColor objects)
        if not _line_has_noFill and hasattr(shape, "line") and shape.line:
            try:
                # Try color attribute first
                if hasattr(shape.line, "color") and shape.line.color:
                    line_color_obj = shape.line.color

                    # Try to get theme color first (if it's a theme color reference)
                    if hasattr(line_color_obj, "theme_color") and line_color_obj.theme_color is not None:
                        theme_color = line_color_obj.theme_color
                        # Try to get theme color value
                        if self.theme_color_extractor:
                            # Convert theme color enum to RGB
                            theme_rgb = self.theme_color_extractor.get_theme_color(theme_color)
                            if theme_rgb:
                                line_color = theme_rgb

                    # If no theme color, try direct RGB
                    if not line_color and hasattr(line_color_obj, "rgb") and line_color_obj.rgb:
                        rgb_obj = line_color_obj.rgb
                        # Handle RGBColor objects
                        if hasattr(rgb_obj, '__class__') and 'RGBColor' in str(rgb_obj.__class__):
                            rgb_str = str(rgb_obj).strip()
                            if len(rgb_str) == 6 and rgb_str.isalnum():
                                line_color = f"#{rgb_str}"
                        # Handle string RGB values
                        elif isinstance(rgb_obj, str):
                            rgb_str = rgb_obj.strip()
                            if len(rgb_str) >= 6:
                                line_color = f"#{rgb_str[:6]}"
                        # Handle integer RGB values
                        elif isinstance(rgb_obj, int) and rgb_obj > 0:
                            line_color = f"#{rgb_obj:06x}"
            except (AttributeError, TypeError, ValueError) as e:
                line_color = None

        # If no direct line color, try to extract from style/lnRef (theme style reference)
        if not _line_has_noFill and not line_color and hasattr(shape, '_element'):
            try:
                ns = {'a': NAMESPACES['a'], 'p': NAMESPACES['p']}

                # Find style element first, then lnRef child
                style = shape._element.find('.//p:style', ns)
                if style is not None:
                    # Find style element first, then lnRef child
                    style = shape._element.find('.//p:style', ns)
                    if style is not None:
                        ln_ref = style.find('.//a:lnRef', ns)
                        if ln_ref is not None:
                            # Check for schemeClr (theme color reference)
                            scheme_clr = ln_ref.find('.//a:schemeClr', ns)
                            if scheme_clr is not None:
                                val = scheme_clr.get('val')  # e.g., "accent1", "accent2", etc.
                                if val and self.theme_color_extractor:
                                    # Convert scheme color name to theme color enum
                                    # val is like "accent3", we need to convert to MSO_THEME_COLOR.ACCENT_3
                                    # The scheme color name uses lowercase without underscore, but the enum uses uppercase WITH underscore
                                    theme_color_name = val.upper()
                                    # Add underscore between letter and number (e.g., "ACCENT3" -> "ACCENT_3")
                                    theme_color_name = re.sub(r'([A-Z]+)(\d+)', r'\1_\2', theme_color_name)
                                    try:
                                        from pptx.enum.dml import MSO_THEME_COLOR
                                        theme_color_enum = getattr(MSO_THEME_COLOR, theme_color_name)
                                        theme_rgb = self.theme_color_extractor.get_theme_color(theme_color_enum)
                                        if theme_rgb:
                                            line_color = theme_rgb
                                    except (AttributeError, ValueError):
                                        pass
            except Exception:
                pass

        # Extract line dash style from XML (<a:ln><a:prstDash val="..."/>)
        dash_style = None
        if hasattr(shape, '_element'):
            try:
                ns = NS_A
                prst_dash = shape._element.find('.//a:ln/a:prstDash', ns)
                if prst_dash is not None:
                    dash_style = prst_dash.get('val')
            except Exception:
                pass

        # Check if shape has text
        text = None
        if hasattr(shape, "text_frame") and shape.text_frame:
            text = shape.text_frame.text

        # Extract blipFill (image fill) - some shapes have images as fill
        blip_fill = self._extract_blip_fill(shape)

        # Extract geometry data for SVG generation
        geometry_data = self._extract_shape_geometry(shape)

        # Build metadata
        metadata = {}
        if geometry_data:
            metadata["geometry"] = geometry_data
        if preset_name:
            metadata["preset_name"] = preset_name

        # Extract flipH/flipV, rotation, scene3d, and line width from XML
        flip_h = False
        flip_v = False
        rotation = None
        if hasattr(shape, '_element'):
            try:
                ns_xml = self._ns
                xfrm = shape._element.find('.//a:xfrm', ns_xml)
                if xfrm is not None:
                    flip_h = xfrm.get('flipH') == '1'
                    flip_v = xfrm.get('flipV') == '1'
                    rot_val = xfrm.get('rot')
                    if rot_val:
                        rotation = angle_to_degrees(float(rot_val))

                # Extract scene3d camera preset
                camera = shape._element.find('.//a:scene3d/a:camera', ns_xml)
                if camera is not None:
                    prst = camera.get('prst')
                    if prst:
                        metadata['scene3d_camera'] = prst

                # Extract explicit line width (<a:ln w="...">)
                ln_elem = shape._element.find('.//a:ln', ns_xml)
                if ln_elem is not None:
                    w_attr = ln_elem.get('w')
                    if w_attr:
                        metadata['line_width'] = emu_to_px(int(w_attr))
            except Exception:
                pass

        return ShapeElement(
            element_type="shape",
            left=left,
            top=top,
            width=width,
            height=height,
            z_order=z_order,
            shape_type=shape_type,
            fill_color=fill_color,
            line_color=line_color,
            dash_style=dash_style,
            text=text,
            blip_fill=blip_fill,
            flip_h=flip_h,
            flip_v=flip_v,
            rotation=rotation,
            metadata=metadata,
        )

    def _extract_blip_fill(self, shape) -> Optional[Dict[str, Any]]:
        """
        Extract blipFill (image fill) data from a shape.

        Some shapes use an image as their fill (e.g., callout shapes with
        handwritten-style content). The image is referenced via relationship ID.

        Args:
            shape: python-pptx shape object

        Returns:
            Dictionary with {image_bytes, image_type} or None
        """
        if not hasattr(shape, '_element'):
            return None

        try:
            ns = {'a': NAMESPACES['a'], 'r': NAMESPACES['r']}
            elem = shape._element

            # Look for blipFill with embedded blip
            blip_fill = elem.find('.//a:blipFill', ns)
            if blip_fill is None:
                return None

            blip = blip_fill.find('.//a:blip', ns)
            if blip is None:
                return None

            # Check for embedded image (r:embed) or linked image (r:link)
            embed_id = blip.get(f'{NS_R_CLARK}embed')
            if not embed_id:
                return None

            # Get the image from the relationship
            try:
                rel = shape.part.rels[embed_id]
                image_part = rel.target_part
                image_bytes = image_part.blob
                # Determine image type from content type
                content_type = image_part.content_type
                image_type = content_type.split('/')[-1] if '/' in content_type else 'png'
                # Normalize common types
                if image_type == 'jpeg':
                    image_type = 'jpg'
                elif image_type == 'x-emf':
                    image_type = 'emf'
                elif image_type == 'x-wmf':
                    image_type = 'wmf'

                return {
                    'image_bytes': image_bytes,
                    'image_type': image_type,
                }
            except (KeyError, AttributeError):
                return None

        except Exception:
            return None

    def _extract_shape_geometry(self, shape) -> Optional[Dict[str, Any]]:
        """
        Extract geometry data from PowerPoint shape for SVG generation.

        Args:
            shape: python-pptx shape object

        Returns:
            Dictionary with geometry data or None
        """
        if not hasattr(shape, '_element'):
            return None

        try:
            elem = shape._element
            namespaces = NS_A

            geometry_data = {}

            # Check for preset geometry
            prst_geom = elem.find('.//a:prstGeom', namespaces)
            if prst_geom is not None:
                geometry_data['type'] = 'preset'
                geometry_data['prst'] = prst_geom.get('prst')

                # Extract adjustment values from avLst
                av_lst = prst_geom.find('a:avLst', namespaces)
                if av_lst is not None:
                    adjustments = []
                    for gd in av_lst.findall('a:gd', namespaces):
                        name = gd.get('name')
                        fmla = gd.get('fmla', '')
                        if fmla.startswith('val '):
                            try:
                                value = int(fmla[4:])
                                adjustments.append({'name': name, 'value': value})
                            except ValueError:
                                pass
                    if adjustments:
                        geometry_data['adjustments'] = adjustments

            # Check for custom geometry (freeform shapes)
            cust_geom = elem.find('.//a:custGeom', namespaces)
            if cust_geom is not None:
                geometry_data['type'] = 'custom'

                # Extract path data
                path_lst = cust_geom.find('a:pathLst', namespaces)
                if path_lst is not None:
                    paths = []
                    for path in path_lst.findall('a:path', namespaces):
                        path_data = {
                            'w': path.get('w'),
                            'h': path.get('h'),
                            'commands': []
                        }

                        # Parse path commands
                        for child in path:
                            tag_name = child.tag.split('}')[-1]

                            if tag_name == 'moveTo':
                                pt = child.find('a:pt', namespaces)
                                if pt is not None:
                                    path_data['commands'].append({
                                        'type': 'M',
                                        'x': pt.get('x'),
                                        'y': pt.get('y')
                                    })

                            elif tag_name == 'lnTo':
                                pt = child.find('a:pt', namespaces)
                                if pt is not None:
                                    path_data['commands'].append({
                                        'type': 'L',
                                        'x': pt.get('x'),
                                        'y': pt.get('y')
                                    })

                            elif tag_name == 'cubicBezTo':
                                points = child.findall('a:pt', namespaces)
                                if len(points) >= 3:
                                    path_data['commands'].append({
                                        'type': 'C',
                                        'x1': points[0].get('x'),
                                        'y1': points[0].get('y'),
                                        'x2': points[1].get('x'),
                                        'y2': points[1].get('y'),
                                        'x': points[2].get('x'),
                                        'y': points[2].get('y')
                                    })

                            elif tag_name == 'quadBezTo':
                                points = child.findall('a:pt', namespaces)
                                if len(points) >= 2:
                                    path_data['commands'].append({
                                        'type': 'Q',
                                        'x1': points[0].get('x'),
                                        'y1': points[0].get('y'),
                                        'x': points[1].get('x'),
                                        'y': points[1].get('y')
                                    })

                            elif tag_name == 'arcTo':
                                path_data['commands'].append({
                                    'type': 'A',
                                    'wR': child.get('wR'),
                                    'hR': child.get('hR'),
                                    'stAng': child.get('stAng'),
                                    'swAng': child.get('swAng'),
                                    'x': child.get('x'),
                                    'y': child.get('y')
                                })

                            elif tag_name == 'close':
                                path_data['commands'].append({'type': 'Z'})

                        if path_data['commands']:
                            paths.append(path_data)

                    geometry_data['paths'] = paths

            return geometry_data if geometry_data else None

        except Exception:
            # If geometry extraction fails, return None
            return None

    def _is_slide_hidden(self, slide) -> bool:
        """
        Detect if a slide is marked as hidden in PowerPoint.

        Args:
            slide: python-pptx slide object

        Returns:
            True if slide is hidden, False otherwise
        """
        try:
            # Try to access the slide's XML element directly
            elem = slide._element

            # Check for common hidden attributes
            # PowerPoint uses different attributes to mark slides as hidden
            for attr in ("show", "hidden"):
                val = elem.get(attr)
                if val is not None:
                    # "show" attribute: 0 means hidden, 1 means visible
                    # "hidden" attribute: 1 means hidden, 0/none means visible
                    if attr == "show":
                        return str(val) == "0"
                    else:  # "hidden" attribute
                        return str(val) == "1"

            # Try namespace-specific attributes
            # PowerPoint uses the PresentationML namespace
            ns = NS_P_CLARK
            for attr in ("show", "hidden"):
                val = elem.get(f"{ns}{attr}")
                if val is not None:
                    if attr == "show":
                        return str(val) == "0"
                    else:  # "hidden" attribute
                        return str(val) == "1"

        except Exception:
            # If any error occurs, assume slide is not hidden
            pass

        return False

    def _detect_slide_animations(self, slide) -> bool:
        """
        Detect if a slide has animations.

        Args:
            slide: python-pptx slide object

        Returns:
            True if slide has animations, False otherwise
        """
        try:
            # Try to access the slide's XML element directly
            elem = slide._element

            # Define the PresentationML namespace
            ns = {"p": NAMESPACES['p']}

            # Look for timing elements which indicate animations
            # In PowerPoint XML, animations are stored in timing elements
            timing = elem.find('.//p:timing', ns)
            if timing is not None:
                return True

            # Also check for animation elements directly
            anim = elem.find('.//p:anim', ns)
            if anim is not None:
                return True

        except Exception:
            # If any error occurs, assume no animations
            pass

        return False

    def get_presentation_metadata(self) -> Dict[str, Any]:
        """
        Get metadata about the presentation.

        Returns:
            Dictionary with presentation metadata
        """
        if not self.presentation:
            return {}

        return {
            "title": self.presentation.core_properties.title or "",
            "author": self.presentation.core_properties.author or "",
            "subject": self.presentation.core_properties.subject or "",
            "created": str(self.presentation.core_properties.created) if self.presentation.core_properties.created else "",
            "modified": str(self.presentation.core_properties.modified) if hasattr(self.presentation.core_properties, "modified") and self.presentation.core_properties.modified else "",
            "slide_count": len(self.presentation.slides),
            "slide_width": self.presentation.slide_width,
            "slide_height": self.presentation.slide_height,
        }
