"""
测试脚本：精确测量PPT和HTML的段前段后距差异

目标：
1. 生成测试PPT文件，包含不同的段前段后距值
2. 将PPT转换为HTML
3. 使用浏览器自动化测量实际段落间距
4. 计算最佳调整比例系数

测试范围：
- 段前距：0pt, 5pt, 10pt, 15pt, 20pt
- 段后距：0pt, 5pt, 10pt, 15pt, 20pt
- 字体大小：12pt, 16pt, 20pt
"""

import asyncio
import json
from pathlib import Path
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt


class ParagraphSpacingTester:
    """段前段后距精确测量工具"""

    def __init__(self, output_dir: Path = None):
        self.results = []
        # 测试不同的段前段后距值（覆盖常用范围）
        self.test_spacing_before_values = [0, 5, 10, 15, 20]
        self.test_spacing_after_values = [0, 5, 10]
        # 测试不同的字体大小
        self.test_font_sizes = [12, 16, 20]
        self.output_dir = output_dir or Path.cwd() / 'test_output'
        self.output_dir.mkdir(exist_ok=True)

    def create_test_ppt(self, spacing_before: int, spacing_after: int, font_size: int) -> Path:
        """
        创建测试PPT文件

        Args:
            spacing_before: 段前距（点）
            spacing_after: 段后距（点）
            font_size: 字体大小（点）

        Returns:
            PPT文件路径
        """
        prs = Presentation()
        # 使用空白布局
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # 添加文本框
        textbox = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(8), Inches(5)
        )
        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        # 添加多行文本用于测量
        test_lines = [
            "Paragraph 1: First paragraph for spacing measurement",
            "Paragraph 2: ABCDEFGHIJKLMNOPQRSTUVWXYZ",
            "Paragraph 3: The quick brown fox jumps over lazy dog",
            "Paragraph 4: Testing spacing between paragraphs",
            "Paragraph 5: Final paragraph for measurement"
        ]

        for i, line in enumerate(test_lines):
            if i > 0:
                text_frame.add_paragraph()

            # 设置段落文本和格式
            para = text_frame.paragraphs[i]
            para.text = line
            para.font.size = Pt(font_size)
            para.font.name = 'Arial'

            # 设置段前段后距
            # 注意：python-pptx的单位是English Metric Units (EMUs)
            # 1 inch = 914400 EMUs
            # 1 point = 1/72 inch = 914400/72 = 12700 EMUs
            if spacing_before > 0:
                para.space_before = Pt(spacing_before)
            if spacing_after > 0:
                para.space_after = Pt(spacing_after)

        # 保存PPT
        filename = f"test_para_spacing_before{spacing_before}_after{spacing_after}_{font_size}pt.pptx"
        ppt_path = self.output_dir / filename
        prs.save(ppt_path)
        return ppt_path

    async def measure_html_paragraph_spacing(self, html_path: Path) -> dict:
        """
        使用Playwright测量HTML中的实际段落间距

        Args:
            html_path: HTML文件路径

        Returns:
            测量结果字典
        """
        try:
            from playwright.async_api import async_playwright
        except ImportError:
            print("警告: 未安装playwright，跳过浏览器测量")
            print("请运行: pip install playwright && playwright install chromium")
            return {}

        async with async_playwright() as p:
            browser = await p.chromium.launch()
            page = await browser.new_page()

            await page.goto(f"file://{html_path.absolute()}")

            # 测量段落间距
            measurements = await page.evaluate("""
                () => {
                    const paragraphs = document.querySelectorAll('p.text-paragraph, p');
                    const results = [];

                    for (let i = 1; i < paragraphs.length; i++) {
                        const prev = paragraphs[i - 1].getBoundingClientRect();
                        const curr = paragraphs[i].getBoundingClientRect();

                        // 测量段落之间的距离（从上一段底部到当前段落顶部）
                        const paragraph_spacing = curr.top - prev.bottom;

                        // 获取段落样式
                        const prevMarginBottom = window.getComputedStyle(paragraphs[i - 1]).marginBottom;
                        const currMarginTop = window.getComputedStyle(paragraphs[i]).marginTop;

                        results.push({
                            paragraph_number: i + 1,
                            paragraph_spacing: paragraph_spacing,
                            prev_margin_bottom: prevMarginBottom,
                            curr_margin_top: currMarginTop,
                            prev_bottom: prev.bottom,
                            curr_top: curr.top
                        });
                    }

                    return results;
                }
            """)

            await browser.close()
            return measurements

    async def run_test_for_spacing(self, spacing_before: int, spacing_after: int, font_size: int):
        """
        测试特定段前段后距和字体大小

        Args:
            spacing_before: PPT段前距（点）
            spacing_after: PPT段后距（点）
            font_size: 字体大小
        """
        print(f"\n测试段前距 {spacing_before}pt, 段后距 {spacing_after}pt, 字体 {font_size}pt...")

        # 1. 创建测试PPT
        ppt_path = self.create_test_ppt(spacing_before, spacing_after, font_size)

        # 2. 转换为HTML
        from shuttleslide.pptx_to_html.parser import PPTXParser
        from shuttleslide.pptx_to_html.layouts.slideshow import SlideshowLayout

        parser = PPTXParser(ppt_path)
        slides = parser.parse()
        layout = SlideshowLayout()
        html = layout.convert(slides)

        html_path = ppt_path.with_suffix('.html')
        html_path.write_text(html, encoding='utf-8')

        # 3. 测量HTML段落间距
        measurements = await self.measure_html_paragraph_spacing(html_path)

        if not measurements:
            print(f"  跳过（未安装playwright）")
            # 清理测试文件
            ppt_path.unlink(missing_ok=True)
            html_path.unlink(missing_ok=True)
            return

        # 4. 记录结果
        if measurements:
            # 计算平均段落间距
            avg_spacing = sum(m['paragraph_spacing'] for m in measurements) / len(measurements)

            # PPT理论段落间距 = spacing_before + spacing_after
            # 但在0pt时，PPT使用默认紧凑间距
            theoretical_ppt_spacing = spacing_before + spacing_after

            result = {
                'ppt_spacing_before': spacing_before,
                'ppt_spacing_after': spacing_after,
                'font_size': font_size,
                'ppt_total_spacing': theoretical_ppt_spacing,
                'html_avg_paragraph_spacing': round(avg_spacing, 2),
                'measurements': measurements,
                'num_paragraphs': len(measurements) + 1
            }

            self.results.append(result)
            print(f"  PPT段前距: {spacing_before}pt, 段后距: {spacing_after}pt")
            print(f"  PPT总间距: {theoretical_ppt_spacing}pt")
            print(f"  HTML平均段落间距: {avg_spacing:.2f}px")
            print(f"  测量段落数: {len(measurements) + 1}")
        else:
            print(f"  测量失败")

        # 5. 清理测试文件
        ppt_path.unlink(missing_ok=True)
        html_path.unlink(missing_ok=True)

    async def run_all_tests(self):
        """运行所有测试"""
        print("开始段前段后距精确测量测试...")
        print("=" * 70)
        print(f"测试参数:")
        print(f"  段前距: {self.test_spacing_before_values}pt")
        print(f"  段后距: {self.test_spacing_after_values}pt")
        print(f"  字体大小: {self.test_font_sizes}pt")
        total_tests = (len(self.test_spacing_before_values) *
                      len(self.test_spacing_after_values) *
                      len(self.test_font_sizes))
        print(f"  总测试数: {total_tests}")
        print("=" * 70)

        # 先检查playwright是否可用
        try:
            from playwright.async_api import async_playwright
        except ImportError:
            print("\n错误: 未安装playwright")
            print("请先安装依赖:")
            print("  pip install playwright")
            print("  playwright install chromium")
            return

        for spacing_before in self.test_spacing_before_values:
            for spacing_after in self.test_spacing_after_values:
                for font_size in self.test_font_sizes:
                    await self.run_test_for_spacing(spacing_before, spacing_after, font_size)

        if not self.results:
            print("\n没有测试结果，请检查playwright安装")
            return

        print("\n" + "=" * 70)
        print("测试完成！正在分析结果...")

        self.analyze_results()

    def analyze_results(self):
        """分析测试结果，计算最佳调整比例系数"""
        print("\n=== 段前段后距测试结果分析 ===\n")

        # 按PPT总间距分组
        spacing_groups = {}
        for result in self.results:
            total_spacing = result['ppt_total_spacing']
            if total_spacing not in spacing_groups:
                spacing_groups[total_spacing] = []
            spacing_groups[total_spacing].append(result)

        # 计算每个PPT间距值对应的HTML间距
        spacing_analysis = {}

        for total_spacing in sorted(spacing_groups.keys()):
            group_results = spacing_groups[total_spacing]

            # 计算该组的平均HTML段落间距
            html_distances = [r['html_avg_paragraph_spacing'] for r in group_results]
            avg_html_distance = sum(html_distances) / len(html_distances)

            spacing_analysis[total_spacing] = {
                'ppt_spacing': total_spacing,
                'avg_html_spacing_px': avg_html_distance,
                'sample_count': len(group_results),
                'html_distances': html_distances
            }

            print(f"PPT总间距: {total_spacing}pt")
            print(f"  HTML平均段落间距: {avg_html_distance:.2f}px")
            print(f"  测量样本数: {len(group_results)}")

        # 计算比例系数
        print("\n=== 比例系数分析 ===\n")

        ratios = []
        ratio_details = []

        for spacing, data in spacing_analysis.items():
            if spacing == 0:
                # 0pt时使用特殊的调整逻辑
                continue

            # 计算比例：HTML间距 / PPT间距
            # 注意单位转换：PPT使用pt，HTML测量的是px
            # 1pt ≈ 1.333px (96 DPI下)
            # 所以PPT的X pt在屏幕上应该显示为 X * 1.333 px

            ppt_px = spacing * 1.333  # 理论屏幕像素
            html_px = data['avg_html_spacing_px']

            # 实际渲染比例
            if ppt_px > 0:
                ratio = html_px / ppt_px
                ratios.append(ratio)
                ratio_details.append({
                    'ppt_spacing_pt': spacing,
                    'ppt_theoretical_px': ppt_px,
                    'html_actual_px': html_px,
                    'ratio': ratio
                })

                print(f"PPT {spacing}pt (理论 {ppt_px:.2f}px) → HTML {html_px:.2f}px")
                print(f"  比例系数: {ratio:.3f}")

        if ratios:
            avg_ratio = sum(ratios) / len(ratios)
            print(f"\n全局平均比例系数: {avg_ratio:.3f}")
            print(f"解释: HTML的段落间距是PPT的 {avg_ratio:.1%} 倍")

            # 推荐的调整系数（如果要达到视觉匹配）
            # 如果HTML是PPT的1.2倍，那么我们需要将HTML间距除以1.2
            recommended_factor = round(1 / avg_ratio, 3)
            print(f"\n推荐使用: PARAGRAPH_SPACING_RATIO = {recommended_factor}")
            print(f"  (将PPT间距值乘以 {recommended_factor} 后用于CSS margin)")

            # 特别分析0pt的情况
            if 0 in spacing_analysis:
                zero_pt_data = spacing_analysis[0]
                print(f"\n=== 0pt特殊情况分析 ===")
                print(f"PPT 0pt → HTML {zero_pt_data['avg_html_spacing_px']:.2f}px")
                print(f"这说明PPT的0pt相当于HTML的负margin")
                print(f"当前代码使用: {self.results[0].get('adjustment', 'N/A')}")

        # 保存结果
        self.save_results(spacing_analysis, ratio_details, locals().get('recommended_factor', None))

    def save_results(self, spacing_analysis: dict, ratio_details: list, recommended: float = None):
        """保存测试结果到文件"""
        output = {
            'timestamp': datetime.now().isoformat(),
            'test_parameters': {
                'spacing_before_values': self.test_spacing_before_values,
                'spacing_after_values': self.test_spacing_after_values,
                'font_sizes': self.test_font_sizes
            },
            'detailed_results': self.results,
            'spacing_analysis': spacing_analysis,
            'ratio_analysis': ratio_details,
            'recommendations': {
                'paragraph_spacing_ratio': recommended
            }
        }

        output_path = self.output_dir / 'paragraph_spacing_test_results.json'
        output_path.write_text(json.dumps(output, indent=2, ensure_ascii=False), encoding='utf-8')
        print(f"\n详细结果已保存到: {output_path}")

        # 保存摘要报告
        self.save_summary_report(output_path, recommended, spacing_analysis)

    def save_summary_report(self, json_path: Path, recommended: float, analysis: dict):
        """保存人类可读的摘要报告"""
        report_path = json_path.with_suffix('.txt')
        lines = [
            "段前段后距精确测量测试报告",
            "=" * 70,
            f"测试时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
            f"测试范围: 段前距 {self.test_spacing_before_values}pt, "
            f"段后距 {self.test_spacing_after_values}pt, "
            f"字体 {self.test_font_sizes}pt",
            "",
            "测试结果:",
            f"  总测试数: {len(self.results)}",
            "",
            "各PPT间距对应的HTML间距:",
        ]

        for spacing in sorted(analysis.keys()):
            data = analysis[spacing]
            lines.append(f"  PPT {spacing}pt → HTML {data['avg_html_spacing_px']:.2f}px "
                        f"(样本数: {data['sample_count']})")

        if recommended:
            lines.extend([
                "",
                "应用建议:",
                f"  推荐比例系数: PARAGRAPH_SPACING_RATIO = {recommended}",
                "",
                "使用方法:",
                "  在 text.py 中修改段落间距计算逻辑:",
                "  ```python",
                "  if para.spacing_before is not None and para.spacing_before > 0:",
                "      adjusted = para.spacing_before * PARAGRAPH_SPACING_RATIO",
                "      styles.append(f'margin-top: {adjusted}pt !important')",
                "  ```",
                "",
                "详细数据请查看JSON文件",
            ])

        report_path.write_text('\n'.join(lines), encoding='utf-8')
        print(f"摘要报告已保存到: {report_path}")


async def main():
    """主函数"""
    tester = ParagraphSpacingTester()
    await tester.run_all_tests()


if __name__ == '__main__':
    asyncio.run(main())
