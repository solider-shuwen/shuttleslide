"""
测试脚本：精确测量PPT和HTML的行距差异
使用Playwright控制浏览器，测量实际渲染的行距

目标：
1. 生成测试PPT文件，包含不同行距值
2. 将PPT转换为HTML
3. 使用浏览器自动化测量实际行距
4. 计算最佳调整系数
"""

import asyncio
import json
from pathlib import Path
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt


class LineSpacingTester:
    """行距精确测量工具"""

    def __init__(self, output_dir: Path = None):
        self.results = []
        # 测试不同的行距值（覆盖常用范围）
        self.test_spacing_values = [0.8, 0.85, 0.9, 0.95, 1.0, 1.1, 1.2, 1.5]
        # 测试不同的字体大小
        self.test_font_sizes = [12, 16, 20, 24]
        self.output_dir = output_dir or Path.cwd() / 'test_output'
        self.output_dir.mkdir(exist_ok=True)

    def create_test_ppt(self, spacing_value: float, font_size: int) -> Path:
        """
        创建测试PPT文件

        Args:
            spacing_value: 行距值（如0.9表示90%）
            font_size: 字体大小（点）

        Returns:
            PPT文件路径
        """
        prs = Presentation()
        # 使用空白布局
        blank_layout = prs.slide_layouts[6]  # 空白布局通常是第7个（索引6）
        slide = prs.slides.add_slide(blank_layout)

        # 添加文本框
        textbox = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(8), Inches(5)
        )
        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        # 添加多行文本用于测量
        test_lines = [
            "Line 1: Test text for baseline measurement",
            "Line 2: ABCDEFGHIJKLMNOPQRSTUVWXYZ",
            "Line 3: The quick brown fox jumps",
            "Line 4: abcdefghijklmnopqrstuvwxyz",
            "Line 5: End of test content"
        ]

        for i, line in enumerate(test_lines):
            if i > 0:
                text_frame.add_paragraph()

            # 设置段落文本和格式
            para = text_frame.paragraphs[i]
            para.text = line
            para.line_spacing = spacing_value
            para.font.size = Pt(font_size)
            para.font.name = 'Arial'  # 使用通用字体避免度量差异

        # 保存PPT
        filename = f"test_spacing_{spacing_value}_{font_size}pt.pptx"
        ppt_path = self.output_dir / filename
        prs.save(ppt_path)
        return ppt_path

    async def measure_html_line_height(self, html_path: Path) -> list:
        """
        使用Playwright测量HTML中的实际行距

        Args:
            html_path: HTML文件路径

        Returns:
            测量结果列表
        """
        try:
            from playwright.async_api import async_playwright
        except ImportError:
            print("警告: 未安装playwright，跳过浏览器测量")
            print("请运行: pip install playwright && playwright install chromium")
            return []

        async with async_playwright() as p:
            browser = await p.chromium.launch()
            page = await browser.new_page()

            await page.goto(f"file://{html_path.absolute()}")

            # 测量行距
            measurements = await page.evaluate("""
                () => {
                    const paragraphs = document.querySelectorAll('p.text-paragraph, p');
                    const results = [];

                    for (let i = 1; i < paragraphs.length; i++) {
                        const prev = paragraphs[i - 1].getBoundingClientRect();
                        const curr = paragraphs[i].getBoundingClientRect();

                        // 测量基线距离（近似为顶部距离）
                        const baseline_distance = curr.top - prev.top;

                        results.push({
                            line_number: i + 1,
                            baseline_distance: baseline_distance,
                            line_height: window.getComputedStyle(paragraphs[i]).lineHeight
                        });
                    }

                    return results;
                }
            """)

            await browser.close()
            return measurements

    async def run_test_for_spacing(self, spacing_value: float, font_size: int):
        """
        测试特定行距值和字体大小

        Args:
            spacing_value: PPT行距值
            font_size: 字体大小
        """
        print(f"\n测试行距 {spacing_value}, 字体 {font_size}pt...")

        # 1. 创建测试PPT
        ppt_path = self.create_test_ppt(spacing_value, font_size)

        # 2. 转换为HTML
        from shuttleslide.pptx_to_html.parser import PPTXParser
        from shuttleslide.pptx_to_html.layouts.slideshow import SlideshowLayout

        parser = PPTXParser(ppt_path)
        slides = parser.parse()
        layout = SlideshowLayout()
        html = layout.convert(slides)

        html_path = ppt_path.with_suffix('.html')
        html_path.write_text(html, encoding='utf-8')

        # 3. 测量HTML行距
        measurements = await self.measure_html_line_height(html_path)

        if not measurements:
            print(f"  跳过（未安装playwright）")
            # 清理测试文件
            ppt_path.unlink(missing_ok=True)
            html_path.unlink(missing_ok=True)
            return

        # 4. 记录结果
        if measurements:
            avg_baseline_distance = sum(m['baseline_distance'] for m in measurements) / len(measurements)

            result = {
                'ppt_spacing': spacing_value,
                'font_size': font_size,
                'html_avg_baseline_distance': round(avg_baseline_distance, 2),
                'measurements': measurements,
                'num_lines': len(measurements) + 1
            }

            self.results.append(result)
            print(f"  PPT行距: {spacing_value}")
            print(f"  HTML平均基线距离: {avg_baseline_distance:.2f}px")
            print(f"  测量行数: {len(measurements) + 1}")
        else:
            print(f"  测量失败")

        # 5. 清理测试文件
        ppt_path.unlink(missing_ok=True)
        html_path.unlink(missing_ok=True)

    async def run_all_tests(self):
        """运行所有测试"""
        print("开始行距精确测量测试...")
        print("=" * 60)
        print(f"测试参数:")
        print(f"  行距值: {self.test_spacing_values}")
        print(f"  字体大小: {self.test_font_sizes}pt")
        print(f"  总测试数: {len(self.test_spacing_values) * len(self.test_font_sizes)}")
        print("=" * 60)

        # 先检查playwright是否可用
        try:
            from playwright.async_api import async_playwright
        except ImportError:
            print("\n错误: 未安装playwright")
            print("请先安装依赖:")
            print("  pip install playwright")
            print("  playwright install chromium")
            return

        for spacing in self.test_spacing_values:
            for font_size in self.test_font_sizes:
                await self.run_test_for_spacing(spacing, font_size)

        if not self.results:
            print("\n没有测试结果，请检查playwright安装")
            return

        print("\n" + "=" * 60)
        print("测试完成！正在分析结果...")

        self.analyze_results()

    def analyze_results(self):
        """分析测试结果，计算最佳调整系数"""
        print("\n=== 测试结果分析 ===\n")

        # 计算每个PPT行距值对应的最佳调整系数
        adjustment_factors = {}

        for spacing in self.test_spacing_values:
            spacing_results = [r for r in self.results if r['ppt_spacing'] == spacing]

            if not spacing_results:
                continue

            # 计算理论基线距离（PPT）
            # 假设PPT使用 baseline_distance = font_size × spacing × 1.2
            theoretical_distances = {
                font_size: font_size * spacing * 1.2
                for font_size in self.test_font_sizes
            }

            # 计算每个测试的调整系数
            factors = []
            for result in spacing_results:
                font_size = result['font_size']
                html_distance = result['html_avg_baseline_distance']
                theoretical = theoretical_distances[font_size]

                # CSS line-height = baseline_distance / font_size × adjustment
                # adjustment = (html_distance / font_size) / spacing
                factor = (html_distance / font_size) / spacing
                factors.append(factor)

            avg_factor = sum(factors) / len(factors)
            adjustment_factors[spacing] = round(avg_factor, 3)

            print(f"PPT行距 {spacing}:")
            print(f"  HTML实际基线距离: {[r['html_avg_baseline_distance'] for r in spacing_results]}px")
            print(f"  计算调整系数: {[f'{f:.3f}' for f in factors]}")
            print(f"  平均调整系数: {avg_factor:.3f}")

        # 全局分析
        global_avg = sum(adjustment_factors.values()) / len(adjustment_factors)
        print(f"\n全局平均调整系数: {global_avg:.3f}")

        # 推荐值
        recommended = round(global_avg, 3)
        print(f"\n推荐使用: LINE_HEIGHT_ADJUSTMENT = {recommended}")

        # 分析不同范围的建议
        tight_avg = sum(v for k, v in adjustment_factors.items() if k < 1.0) / len([k for k in adjustment_factors if k < 1.0])
        normal_avg = sum(v for k, v in adjustment_factors.items() if 1.0 <= k <= 1.2) / len([k for k in adjustment_factors if 1.0 <= k <= 1.2])
        print(f"\n分段分析:")
        print(f"  紧凑行距 (<1.0) 平均系数: {tight_avg:.3f}")
        if 1.0 <= k <= 1.2:
            print(f"  正常行距 (1.0-1.2) 平均系数: {normal_avg:.3f}")

        # 保存结果
        self.save_results(adjustment_factors, recommended)

    def save_results(self, adjustment_factors: dict, recommended: float):
        """保存测试结果到文件"""
        output = {
            'timestamp': datetime.now().isoformat(),
            'test_spacing_values': self.test_spacing_values,
            'test_font_sizes': self.test_font_sizes,
            'detailed_results': self.results,
            'adjustment_factors_by_spacing': adjustment_factors,
            'recommended_global_factor': recommended
        }

        output_path = self.output_dir / 'line_spacing_test_results.json'
        output_path.write_text(json.dumps(output, indent=2), encoding='utf-8')
        print(f"\n详细结果已保存到: {output_path}")

        # 同时保存一份摘要报告
        self.save_summary_report(output_path, recommended, adjustment_factors)

    def save_summary_report(self, json_path: Path, recommended: float, factors: dict):
        """保存人类可读的摘要报告"""
        report_path = json_path.with_suffix('.txt')
        lines = [
            "行距精确测量测试报告",
            "=" * 60,
            f"测试时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
            f"测试范围: 行距 {self.test_spacing_values}, 字体 {self.test_font_sizes}pt",
            "",
            "测试结果:",
            f"  总测试数: {len(self.results)}",
            f"  推荐调整系数: {recommended}",
            "",
            "各行距值对应的调整系数:",
        ]

        for spacing in sorted(factors.keys()):
            lines.append(f"  PPT行距 {spacing} -> 调整系数 {factors[spacing]}")

        lines.extend([
            "",
            "应用建议:",
            f"  在 text.py 中设置: LINE_HEIGHT_ADJUSTMENT = {recommended}",
            "",
            "详细数据请查看JSON文件",
        ])

        report_path.write_text('\n'.join(lines), encoding='utf-8')
        print(f"摘要报告已保存到: {report_path}")


async def main():
    """主函数"""
    tester = LineSpacingTester()
    await tester.run_all_tests()


if __name__ == '__main__':
    asyncio.run(main())
